#!/usr/bin/env python3
"""
Build a native, editable PowerPoint file for the Week 4 Advanced Workshop deck.

Source of truth: docs/decks/week-4-advanced.html
Output:          docs/pptx/week-4-advanced.pptx

The PPTX:
  - 16:9 widescreen (13.333 x 7.5 in)
  - Real PowerPoint shapes / text frames / tables (no images-of-slides)
  - Native speaker notes pulled from the <aside class="notes"> block
  - PowerPoint-safe fonts (Calibri / Calibri Light / Consolas)
  - A careful approximation of Week 4's dark palette + yellow BREAK accent

Self-contained: no imports from sibling build_pptx_week_*.py scripts.
"""

from __future__ import annotations

import re
import sys
from pathlib import Path
from typing import Iterable

from bs4 import BeautifulSoup, NavigableString, Tag
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree


# ─────────────────────────────────────────────────────────────────────────────
# Paths
# ─────────────────────────────────────────────────────────────────────────────
ROOT = Path(__file__).resolve().parent.parent
SRC_HTML = ROOT / "docs" / "decks" / "week-4-advanced.html"
OUT_PPTX = ROOT / "docs" / "pptx" / "week-4-advanced.pptx"


# ─────────────────────────────────────────────────────────────────────────────
# Palette (Week 4 — dark base, scarlet, yellow BREAK accent)
# ─────────────────────────────────────────────────────────────────────────────
SCARLET       = RGBColor(0xCC, 0x00, 0x00)
SCARLET_DARK  = RGBColor(0xA3, 0x00, 0x00)
GOLD          = RGBColor(0xF5, 0xD1, 0x30)
GOLD_SOFT     = RGBColor(0xFF, 0xF3, 0xB0)
GOLD_DARK     = RGBColor(0xD4, 0xB1, 0x1A)
INK           = RGBColor(0x1A, 0x1A, 0x1A)
INK_SOFT      = RGBColor(0x4A, 0x4A, 0x4A)
INK_MUTED     = RGBColor(0x6E, 0x6E, 0x6E)
PAPER         = RGBColor(0xFF, 0xFF, 0xFF)
PAPER_WARM    = RGBColor(0xF8, 0xF7, 0xF5)
RULE          = RGBColor(0xD9, 0xD8, 0xD4)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_WHITE    = RGBColor(0xF0, 0xF0, 0xF0)


# Type stack — PowerPoint-safe
FONT_DISPLAY = "Calibri"
FONT_BODY    = "Calibri"
FONT_LIGHT   = "Calibri Light"


# Slide geometry (16:9 widescreen, EMU)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ─────────────────────────────────────────────────────────────────────────────
# Low-level helpers
# ─────────────────────────────────────────────────────────────────────────────
def _text_of(el: Tag | None) -> str:
    if el is None:
        return ""
    # Preserve inline emphasis as plain text — formatting handled by callers.
    txt = el.get_text(separator=" ", strip=True)
    # Collapse repeated whitespace; normalize unicode dashes.
    txt = re.sub(r"\s+", " ", txt)
    return txt


def _add_rect(slide, x, y, w, h, fill: RGBColor | None, line: RGBColor | None = None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    return shp


def _estimate_visible_lines(text: str, box_w_in: float, font_size_pt: float) -> int:
    """Conservative visible-line count for bold sans display text.
    Mirrors the helper in scripts/build_pptx_week_6.py so downstream layout
    (subheads, columns, timeboxes) can be sized from the actual line count."""
    if not text:
        return 1
    char_w = 0.62 * float(font_size_pt) / 72.0
    chars_per_line = max(6, int(float(box_w_in) / char_w))
    total = 0
    for chunk in str(text).split("\n"):
        n = max(1, len(chunk))
        total += (n + chars_per_line - 1) // chars_per_line
    return max(1, total)


def _add_text(slide, x, y, w, h, *,
              text: str = "",
              font: str = FONT_BODY,
              size: int = 18,
              bold: bool = False,
              italic: bool = False,
              color: RGBColor = INK,
              anchor: str = "top",
              align: str = "left",
              auto_size: bool = False,
              line_spacing: float | None = None,
              fill: RGBColor | None = None):
    box = slide.shapes.add_textbox(x, y, w, h)
    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    else:
        box.fill.background()
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    if anchor == "top":
        tf.vertical_anchor = MSO_ANCHOR.TOP
    elif anchor == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    p = tf.paragraphs[0]
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    elif align == "justify":
        p.alignment = PP_ALIGN.JUSTIFY
    else:
        p.alignment = PP_ALIGN.LEFT
    if line_spacing is not None:
        p.line_spacing = line_spacing
    if text:
        run = p.add_run()
        run.text = text
        f = run.font
        f.name = font
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
    return box


def _add_paragraph(tf, text: str, *, font=FONT_BODY, size=18, bold=False,
                   italic=False, color=INK, align="left",
                   space_before: float | None = None,
                   space_after: float | None = None,
                   line_spacing: float | None = None,
                   bullet: bool = False, indent_level: int = 0):
    p = tf.add_paragraph()
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT
    if space_before is not None:
        p.space_before = Pt(space_before)
    if space_after is not None:
        p.space_after = Pt(space_after)
    if line_spacing is not None:
        p.line_spacing = line_spacing
    p.level = indent_level
    if bullet:
        # Add a real PPT bullet via XML (python-pptx has no first-class API for it).
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        # Remove any existing bullet props.
        for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
            existing = pPr.find(qn(tag))
            if existing is not None:
                pPr.remove(existing)
        bu = etree.SubElement(pPr, qn("a:buChar"))
        bu.set("char", "■")
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    return p


def _add_bulleted_list(slide, x, y, w, h, items: Iterable[str], *,
                       font=FONT_BODY, size=18, color=INK,
                       line_spacing: float = 1.2,
                       space_after: float = 6.0,
                       fill: RGBColor | None = None):
    box = slide.shapes.add_textbox(x, y, w, h)
    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    else:
        box.fill.background()
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    items = list(items)
    if not items:
        return box
    # First paragraph — set text on tf.paragraphs[0]; we'll re-decorate it as bullet.
    first = items[0]
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    p0.line_spacing = line_spacing
    p0.space_after = Pt(space_after)
    pPr0 = p0._p.get_or_add_pPr()
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
        existing = pPr0.find(qn(tag))
        if existing is not None:
            pPr0.remove(existing)
    bu0 = etree.SubElement(pPr0, qn("a:buChar"))
    bu0.set("char", "■")
    pPr0.set("indent", "-228600")  # ~0.25"
    pPr0.set("marL", "228600")
    run0 = p0.add_run()
    run0.text = first
    f0 = run0.font
    f0.name = font
    f0.size = Pt(size)
    f0.color.rgb = color
    for item in items[1:]:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        pPr = p._p.get_or_add_pPr()
        bu = etree.SubElement(pPr, qn("a:buChar"))
        bu.set("char", "■")
        pPr.set("indent", "-228600")
        pPr.set("marL", "228600")
        run = p.add_run()
        run.text = item
        f = run.font
        f.name = font
        f.size = Pt(size)
        f.color.rgb = color
    return box


def _set_notes(slide, paragraphs: list[str]) -> None:
    """Set the slide's native speaker-notes pane content."""
    notes_tf = slide.notes_slide.notes_text_frame
    # Reset
    notes_tf.clear()
    if not paragraphs:
        notes_tf.text = ""
        return
    notes_tf.text = paragraphs[0]
    for p_text in paragraphs[1:]:
        np = notes_tf.add_paragraph()
        np.text = p_text
    # Style the notes uniformly
    for p in notes_tf.paragraphs:
        for r in p.runs:
            r.font.name = FONT_BODY
            r.font.size = Pt(11)


def _footer(slide, left: str, right: str, *,
            on_dark: bool = False, accent: RGBColor = SCARLET):
    """Top scarlet hairline + bottom footer line."""
    # Top scarlet hairline (shows over both light and dark slides; matches HTML).
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), accent)
    # Bottom footer band (subtle on light, transparent on dark).
    band_y = SLIDE_H - Inches(0.45)
    if not on_dark:
        _add_rect(slide, 0, band_y, SLIDE_W, Inches(0.45), PAPER_WARM)
    color = NEAR_WHITE if on_dark else INK_SOFT
    _add_text(slide, Inches(0.6), band_y, Inches(8.0), Inches(0.4),
              text=left, font=FONT_BODY, size=10, color=color, anchor="middle")
    _add_text(slide, SLIDE_W - Inches(2.6), band_y, Inches(2.0), Inches(0.4),
              text=right, font=FONT_BODY, size=10, color=color, anchor="middle",
              align="right")


# ─────────────────────────────────────────────────────────────────────────────
# Presentation setup
# ─────────────────────────────────────────────────────────────────────────────
def _new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    # Core props
    prs.core_properties.title = "Week 4 — Advanced Workshop"
    prs.core_properties.subject = "Expert-Driven Development · Week 4"
    prs.core_properties.author = "EDD Course Team"
    return prs


def _blank_slide(prs: Presentation, fill: RGBColor = PAPER):
    layout = prs.slide_layouts[6]  # 6 = blank
    slide = prs.slides.add_slide(layout)
    # Force a background colour by drawing a full-bleed rectangle behind everything.
    bg = _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill)
    # Send to back so subsequent shapes sit above.
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return slide


# ─────────────────────────────────────────────────────────────────────────────
# Layout builders
# ─────────────────────────────────────────────────────────────────────────────
def build_cover(prs, *, eyebrow, title, tagline, meta, footer_left, page_no, notes):
    slide = _blank_slide(prs, fill=INK)
    # Scarlet vertical bar
    _add_rect(slide, 0, 0, Inches(0.20), SLIDE_H, SCARLET)
    # Eyebrow
    _add_text(slide, Inches(1.0), Inches(0.7), Inches(11), Inches(0.5),
              text=eyebrow, font=FONT_DISPLAY, size=18, bold=True, color=GOLD)
    # Title
    _add_text(slide, Inches(1.0), Inches(1.3), Inches(11), Inches(2.4),
              text=title, font=FONT_DISPLAY, size=72, bold=True, color=WHITE,
              line_spacing=1.0)
    # Tagline
    _add_text(slide, Inches(1.0), Inches(4.0), Inches(11), Inches(1.6),
              text=tagline, font=FONT_LIGHT, size=22, color=NEAR_WHITE,
              line_spacing=1.25)
    # Meta row
    cols = len(meta) if meta else 0
    if cols:
        col_w = Inches(11.0 / cols)
        for i, (label, value) in enumerate(meta):
            x = Inches(1.0) + col_w * i
            _add_text(slide, x, Inches(5.7), col_w - Inches(0.2), Inches(0.35),
                      text=label.upper(), font=FONT_DISPLAY, size=11, bold=True,
                      color=GOLD)
            _add_text(slide, x, Inches(6.05), col_w - Inches(0.2), Inches(0.7),
                      text=value, font=FONT_BODY, size=14, color=NEAR_WHITE,
                      line_spacing=1.25)
    _footer(slide, footer_left, str(page_no), on_dark=True)
    _set_notes(slide, notes)


def build_section(prs, *, module_tag, title, timing, blurb, footer_left, page_no, notes):
    slide = _blank_slide(prs, fill=SCARLET_DARK)
    # Module tag
    _add_text(slide, Inches(1.0), Inches(0.9), Inches(11), Inches(0.5),
              text=module_tag.upper(), font=FONT_DISPLAY, size=18, bold=True, color=GOLD)
    # Title
    _add_text(slide, Inches(1.0), Inches(1.5), Inches(11), Inches(2.6),
              text=title, font=FONT_DISPLAY, size=64, bold=True, color=WHITE,
              line_spacing=1.0)
    # Timing pill
    pill_w = Inches(2.6)
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(1.0), Inches(4.4), pill_w, Inches(0.6))
    pill.shadow.inherit = False
    pill.fill.solid()
    pill.fill.fore_color.rgb = GOLD
    pill.line.fill.background()
    tf = pill.text_frame
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = timing
    r.font.name = FONT_DISPLAY
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = INK
    # Blurb
    _add_text(slide, Inches(1.0), Inches(5.3), Inches(11.2), Inches(1.5),
              text=blurb, font=FONT_LIGHT, size=20, color=RGBColor(0xFF, 0xF5, 0xF5),
              line_spacing=1.3)
    _footer(slide, footer_left, str(page_no), on_dark=True)
    _set_notes(slide, notes)


def build_content_two_col(prs, *, title, subhead, left_h, left_items,
                          right_h, right_items, footer_left, page_no, notes,
                          tail: str | None = None):
    slide = _blank_slide(prs, fill=PAPER)
    # Title — size box and downstream subhead/body from actual line count.
    title_lines = _estimate_visible_lines(title, 12.2, 34)
    title_h_in = max(0.85, (34 / 72.0) * 1.05 * title_lines + 0.20)
    _add_text(slide, Inches(0.6), Inches(0.45), Inches(12.2), Inches(title_h_in),
              text=title, font=FONT_DISPLAY, size=34, bold=True, color=INK,
              line_spacing=1.05)
    sub_y = Inches(0.45 + title_h_in + 0.05)
    if subhead:
        _add_text(slide, Inches(0.6), sub_y, Inches(12.2), Inches(0.4),
                  text=subhead, font=FONT_LIGHT, size=16, italic=True, color=INK_MUTED)
        body_top = sub_y + Inches(0.55)
    else:
        body_top = sub_y + Inches(0.15)
    body_h = Inches(7.0) - body_top - Inches(0.5)
    col_w = Inches(6.0)
    # Left column header
    _add_text(slide, Inches(0.6), body_top, col_w, Inches(0.45),
              text=left_h, font=FONT_DISPLAY, size=18, bold=True, color=SCARLET)
    _add_bulleted_list(slide, Inches(0.6), body_top + Inches(0.5), col_w, body_h - Inches(0.5),
                       items=left_items, size=15, color=INK, line_spacing=1.2,
                       space_after=4)
    # Right column header
    _add_text(slide, Inches(6.9), body_top, col_w, Inches(0.45),
              text=right_h, font=FONT_DISPLAY, size=18, bold=True, color=SCARLET)
    _add_bulleted_list(slide, Inches(6.9), body_top + Inches(0.5), col_w, body_h - Inches(0.5),
                       items=right_items, size=15, color=INK, line_spacing=1.2,
                       space_after=4)
    if tail:
        _add_text(slide, Inches(0.6), Inches(6.55), Inches(12.2), Inches(0.4),
                  text=tail, font=FONT_LIGHT, size=13, italic=True, color=INK_SOFT)
    _footer(slide, footer_left, str(page_no))
    _set_notes(slide, notes)


def build_content_bullets(prs, *, title, subhead, items, footer_left, page_no,
                          notes, tail: str | None = None):
    slide = _blank_slide(prs, fill=PAPER)
    title_lines = _estimate_visible_lines(title, 12.2, 34)
    title_h_in = max(0.85, (34 / 72.0) * 1.05 * title_lines + 0.20)
    _add_text(slide, Inches(0.6), Inches(0.45), Inches(12.2), Inches(title_h_in),
              text=title, font=FONT_DISPLAY, size=34, bold=True, color=INK,
              line_spacing=1.05)
    sub_y = Inches(0.45 + title_h_in + 0.05)
    if subhead:
        _add_text(slide, Inches(0.6), sub_y, Inches(12.2), Inches(0.4),
                  text=subhead, font=FONT_LIGHT, size=16, italic=True, color=INK_MUTED)
        body_top = sub_y + Inches(0.55)
    else:
        body_top = sub_y + Inches(0.10)
    _add_bulleted_list(slide, Inches(0.6), body_top, Inches(12.2),
                       Inches(6.35) - body_top,
                       items=items, size=18, color=INK, line_spacing=1.3,
                       space_after=8)
    if tail:
        _add_text(slide, Inches(0.6), Inches(6.45), Inches(12.2), Inches(0.55),
                  text=tail, font=FONT_LIGHT, size=14, italic=True, color=INK_SOFT,
                  line_spacing=1.3)
    _footer(slide, footer_left, str(page_no))
    _set_notes(slide, notes)


def build_content_pullquote(prs, *, title, quote, attrib, footer_left, page_no, notes,
                            tail: str | None = None):
    slide = _blank_slide(prs, fill=PAPER)
    _add_text(slide, Inches(0.6), Inches(0.45), Inches(12.2), Inches(0.95),
              text=title, font=FONT_DISPLAY, size=34, bold=True, color=INK)
    # Quote box
    box_top = Inches(2.0)
    box_h = Inches(3.6)
    _add_rect(slide, Inches(0.6), box_top, Inches(12.2), box_h, PAPER_WARM, line=None)
    # Scarlet left border
    _add_rect(slide, Inches(0.6), box_top, Inches(0.12), box_h, SCARLET)
    _add_text(slide, Inches(1.0), box_top + Inches(0.2),
              Inches(11.6), box_h - Inches(0.4),
              text=quote, font=FONT_LIGHT, size=22, italic=False, color=INK,
              line_spacing=1.3, anchor="middle")
    if attrib:
        _add_text(slide, Inches(1.0), box_top + box_h - Inches(0.6),
                  Inches(11.6), Inches(0.4),
                  text="— " + attrib, font=FONT_DISPLAY, size=14, bold=True,
                  color=SCARLET_DARK, align="right")
    if tail:
        _add_text(slide, Inches(0.6), Inches(6.0), Inches(12.2), Inches(0.55),
                  text=tail, font=FONT_LIGHT, size=15, italic=True, color=INK_SOFT)
    _footer(slide, footer_left, str(page_no))
    _set_notes(slide, notes)


def build_table_slide(prs, *, title, subhead, headers, rows, footer_left,
                      page_no, notes, col_widths: list[float] | None = None,
                      first_col_accent: bool = True, header_band: bool = True):
    slide = _blank_slide(prs, fill=PAPER)
    _add_text(slide, Inches(0.6), Inches(0.45), Inches(12.2), Inches(0.7),
              text=title, font=FONT_DISPLAY, size=30, bold=True, color=INK)
    if subhead:
        _add_text(slide, Inches(0.6), Inches(1.15), Inches(12.2), Inches(0.4),
                  text=subhead, font=FONT_LIGHT, size=14, italic=True, color=INK_MUTED)
    # Table
    n_cols = len(headers) if headers else len(rows[0])
    n_rows = (1 if headers else 0) + len(rows)
    tbl_top = Inches(1.65) if subhead else Inches(1.45)
    tbl_left = Inches(0.6)
    tbl_width = Inches(12.2)
    tbl_height = Inches(7.5 - 0.5) - tbl_top - Inches(0.2)
    table_shape = slide.shapes.add_table(n_rows, n_cols, tbl_left, tbl_top,
                                         tbl_width, tbl_height)
    table = table_shape.table
    # Column widths
    total = sum(col_widths) if col_widths else n_cols
    if not col_widths:
        col_widths = [1.0] * n_cols
    for i, ratio in enumerate(col_widths):
        table.columns[i].width = int(tbl_width * (ratio / total))
    # Header row
    row_idx = 0
    if headers:
        for ci, htxt in enumerate(headers):
            cell = table.cell(0, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = INK if header_band else PAPER_WARM
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = Inches(0.08)
            tf.margin_right = Inches(0.08)
            tf.margin_top = Inches(0.04)
            tf.margin_bottom = Inches(0.04)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = htxt
            r.font.name = FONT_DISPLAY
            r.font.size = Pt(13)
            r.font.bold = True
            r.font.color.rgb = WHITE if header_band else INK
        row_idx = 1
    # Body rows
    for ri, row in enumerate(rows):
        rr = row_idx + ri
        zebra = (ri % 2 == 1)
        for ci, val in enumerate(row):
            cell = table.cell(rr, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PAPER_WARM if zebra else PAPER
            cell.text = ""
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.08)
            tf.margin_right = Inches(0.08)
            tf.margin_top = Inches(0.04)
            tf.margin_bottom = Inches(0.04)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = val
            r.font.name = FONT_BODY
            r.font.size = Pt(11)
            if first_col_accent and ci == 0:
                r.font.bold = True
                r.font.color.rgb = SCARLET_DARK
            else:
                r.font.color.rgb = INK
    _footer(slide, footer_left, str(page_no))
    _set_notes(slide, notes)


def build_workshop(prs, *, tag, title, timebox, prompt, what_good,
                   footer_left, page_no, notes):
    slide = _blank_slide(prs, fill=GOLD_SOFT)
    # Left gold bar
    _add_rect(slide, 0, 0, Inches(0.25), SLIDE_H, GOLD_DARK)
    # Activity tag pill
    pill_w = Inches(3.6)
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(0.6), Inches(0.45), pill_w, Inches(0.45))
    pill.shadow.inherit = False
    pill.fill.solid()
    pill.fill.fore_color.rgb = INK
    pill.line.fill.background()
    tf = pill.text_frame
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = tag.upper()
    r.font.name = FONT_DISPLAY; r.font.size = Pt(13); r.font.bold = True
    r.font.color.rgb = GOLD
    # Title — size box and downstream timebox/prompt from actual line count.
    title_lines = _estimate_visible_lines(title, 12.0, 44)
    title_h_in = max(1.0, (44 / 72.0) * 1.05 * title_lines + 0.20)
    _add_text(slide, Inches(0.6), Inches(1.05), Inches(12.0), Inches(title_h_in),
              text=title, font=FONT_DISPLAY, size=44, bold=True, color=INK,
              line_spacing=1.05)
    # Timebox
    timebox_y = Inches(1.05 + title_h_in + 0.05)
    _add_text(slide, Inches(0.6), timebox_y, Inches(12.0), Inches(0.4),
              text=timebox, font=FONT_DISPLAY, size=18, bold=True, color=SCARLET)
    # Prompt box
    pb_top = timebox_y + Inches(0.6); pb_h = Inches(2.1)
    _add_rect(slide, Inches(0.6), pb_top, Inches(12.2), pb_h, PAPER, line=INK)
    _add_text(slide, Inches(0.85), pb_top + Inches(0.1), Inches(11.6), Inches(0.35),
              text="PROMPT", font=FONT_DISPLAY, size=11, bold=True, color=SCARLET)
    _add_text(slide, Inches(0.85), pb_top + Inches(0.45),
              Inches(11.6), pb_h - Inches(0.55),
              text=prompt, font=FONT_BODY, size=14, color=INK, line_spacing=1.35)
    # What good
    wg_top = Inches(5.4); wg_h = Inches(1.4)
    _add_rect(slide, Inches(0.6), wg_top, Inches(12.2), wg_h,
              RGBColor(0xFF, 0xFC, 0xE6), line=SCARLET)
    _add_text(slide, Inches(0.85), wg_top + Inches(0.12),
              Inches(11.6), wg_h - Inches(0.24),
              text=what_good, font=FONT_BODY, size=13, color=INK, line_spacing=1.35)
    _footer(slide, footer_left, str(page_no))
    _set_notes(slide, notes)


def build_cue(prs, *, eyebrow, title, body_paragraphs,
              footer_left, page_no, notes):
    slide = _blank_slide(prs, fill=INK)
    # Gold ring
    rx, ry, rw, rh = Inches(1.4), Inches(1.0), Inches(10.5), Inches(5.5)
    ring = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, rx, ry, rw, rh)
    ring.shadow.inherit = False
    ring.fill.background()
    ring.line.color.rgb = GOLD
    ring.line.width = Pt(5)
    # Eyebrow
    _add_text(slide, rx + Inches(0.4), ry + Inches(0.4), rw - Inches(0.8), Inches(0.5),
              text=eyebrow.upper(), font=FONT_DISPLAY, size=14, bold=True,
              color=SCARLET, align="center")
    # Title
    _add_text(slide, rx + Inches(0.3), ry + Inches(1.0), rw - Inches(0.6), Inches(1.4),
              text=title, font=FONT_DISPLAY, size=48, bold=True, color=GOLD,
              align="center", line_spacing=1.05)
    # Body paragraphs
    box = slide.shapes.add_textbox(rx + Inches(0.5), ry + Inches(2.6),
                                   rw - Inches(1.0), rh - Inches(3.0))
    box.fill.background(); box.line.fill.background()
    tf = box.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    if body_paragraphs:
        first = body_paragraphs[0]
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        p0.line_spacing = 1.4
        r0 = p0.add_run(); r0.text = first
        r0.font.name = FONT_BODY; r0.font.size = Pt(20); r0.font.color.rgb = NEAR_WHITE
        for txt in body_paragraphs[1:]:
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            p.line_spacing = 1.4
            p.space_before = Pt(8)
            r = p.add_run(); r.text = txt
            r.font.name = FONT_BODY; r.font.size = Pt(20); r.font.color.rgb = NEAR_WHITE
    _footer(slide, footer_left, str(page_no), on_dark=True)
    _set_notes(slide, notes)


def build_break(prs, *, footer_left, page_no, notes,
                title="BREAK", subtitle="10 minutes · back at the time on screen"):
    slide = _blank_slide(prs, fill=GOLD)
    # Title huge
    _add_text(slide, Inches(0), Inches(2.3), SLIDE_W, Inches(2.6),
              text=title, font=FONT_DISPLAY, size=180, bold=True, color=INK,
              align="center", line_spacing=1.0)
    _add_text(slide, Inches(0), Inches(5.0), SLIDE_W, Inches(0.7),
              text=subtitle, font=FONT_DISPLAY, size=28, bold=True, color=INK,
              align="center")
    _footer(slide, footer_left, str(page_no))
    _set_notes(slide, notes)


def build_ref_grid(prs, *, title, subhead, cards, footer_left, page_no, notes,
                   wide_last: bool = False):
    """cards: list of dicts {num, h3, body}. 4-card grid, optional wide last spanning."""
    slide = _blank_slide(prs, fill=PAPER)
    _add_text(slide, Inches(0.6), Inches(0.45), Inches(12.2), Inches(0.7),
              text=title, font=FONT_DISPLAY, size=30, bold=True, color=INK)
    if subhead:
        _add_text(slide, Inches(0.6), Inches(1.15), Inches(12.2), Inches(0.4),
                  text=subhead, font=FONT_LIGHT, size=14, italic=True, color=INK_MUTED)
    grid_top = Inches(1.65)
    grid_h = Inches(5.2)
    n = len(cards)
    if wide_last and n == 5:
        # 2-col grid: 4 small (2x2) + wide bottom
        cell_w = Inches(6.0); gap = Inches(0.2)
        cell_h = (grid_h - Inches(0.4) - gap * 2) / 3  # 2 rows of small + 1 wide
        # Place first 4 in 2x2
        small_h = (grid_h - gap - Inches(2.0)) / 2  # leave room for wide row
        small_h = Inches(1.5)
        wide_h = grid_h - small_h * 2 - gap * 2
        for i, c in enumerate(cards[:4]):
            row = i // 2; col = i % 2
            x = Inches(0.6) + (cell_w + gap) * col
            y = grid_top + (small_h + gap) * row
            _ref_card(slide, x, y, cell_w, small_h, c)
        # Wide last
        wx = Inches(0.6); wy = grid_top + (small_h + gap) * 2
        ww = cell_w * 2 + gap; wh = wide_h
        _ref_card(slide, wx, wy, ww, wh, cards[4])
    else:
        # Plain 2-column grid
        cell_w = Inches(6.0); gap = Inches(0.2)
        rows = (n + 1) // 2
        cell_h = (grid_h - gap * (rows - 1)) / rows
        for i, c in enumerate(cards):
            row = i // 2; col = i % 2
            x = Inches(0.6) + (cell_w + gap) * col
            y = grid_top + (cell_h + gap) * row
            _ref_card(slide, x, y, cell_w, cell_h, c)
    _footer(slide, footer_left, str(page_no))
    _set_notes(slide, notes)


def _ref_card(slide, x, y, w, h, card):
    _add_rect(slide, x, y, w, h, PAPER_WARM)
    # Scarlet left border
    _add_rect(slide, x, y, Inches(0.10), h, SCARLET)
    _add_text(slide, x + Inches(0.25), y + Inches(0.12),
              w - Inches(0.4), Inches(0.32),
              text=card["num"].upper(), font=FONT_DISPLAY, size=11, bold=True,
              color=SCARLET)
    _add_text(slide, x + Inches(0.25), y + Inches(0.42),
              w - Inches(0.4), Inches(0.45),
              text=card["h3"], font=FONT_DISPLAY, size=18, bold=True, color=INK)
    _add_text(slide, x + Inches(0.25), y + Inches(0.92),
              w - Inches(0.4), h - Inches(1.05),
              text=card["body"], font=FONT_BODY, size=12, color=INK_SOFT,
              line_spacing=1.3)


def build_facilitation(prs, *, title, subhead, panels, footer_left, page_no, notes):
    """panels: list of {h3, items}"""
    slide = _blank_slide(prs, fill=PAPER)
    _add_text(slide, Inches(0.6), Inches(0.45), Inches(12.2), Inches(0.7),
              text=title, font=FONT_DISPLAY, size=30, bold=True, color=INK)
    if subhead:
        _add_text(slide, Inches(0.6), Inches(1.15), Inches(12.2), Inches(0.4),
                  text=subhead, font=FONT_LIGHT, size=14, italic=True, color=INK_MUTED)
    n = len(panels)
    grid_top = Inches(1.7); grid_h = Inches(5.2)
    gap = Inches(0.2)
    panel_w = (Inches(12.2) - gap * (n - 1)) / n
    for i, panel in enumerate(panels):
        x = Inches(0.6) + (panel_w + gap) * i
        y = grid_top
        _add_rect(slide, x, y, panel_w, grid_h, PAPER_WARM)
        # Top scarlet rule
        _add_rect(slide, x, y, panel_w, Inches(0.08), SCARLET)
        _add_text(slide, x + Inches(0.2), y + Inches(0.18),
                  panel_w - Inches(0.4), Inches(0.5),
                  text=panel["h3"].upper(), font=FONT_DISPLAY, size=13, bold=True,
                  color=SCARLET)
        _add_bulleted_list(slide, x + Inches(0.2), y + Inches(0.7),
                           panel_w - Inches(0.4), grid_h - Inches(0.85),
                           items=panel["items"], size=11, color=INK,
                           line_spacing=1.25, space_after=4)
    _footer(slide, footer_left, str(page_no))
    _set_notes(slide, notes)


def build_shareout(prs, *, title, subhead, left_h, left_items, right_h, right_items,
                   right_tail, footer_left, page_no, notes):
    slide = _blank_slide(prs, fill=PAPER)
    _add_text(slide, Inches(0.6), Inches(0.45), Inches(12.2), Inches(1.0),
              text=title, font=FONT_DISPLAY, size=38, bold=True, color=INK,
              line_spacing=1.05)
    if subhead:
        _add_text(slide, Inches(0.6), Inches(1.45), Inches(12.2), Inches(0.45),
                  text=subhead, font=FONT_DISPLAY, size=16, bold=True, color=SCARLET)
    panel_top = Inches(2.05)
    panel_h = Inches(4.6)
    gap = Inches(0.3)
    panel_w = (Inches(12.2) - gap) / 2
    # Left panel
    _add_rect(slide, Inches(0.6), panel_top, panel_w, panel_h, PAPER_WARM)
    _add_text(slide, Inches(0.85), panel_top + Inches(0.2), panel_w - Inches(0.5), Inches(0.5),
              text=left_h, font=FONT_DISPLAY, size=18, bold=True, color=SCARLET)
    _add_bulleted_list(slide, Inches(0.85), panel_top + Inches(0.75),
                       panel_w - Inches(0.5), panel_h - Inches(0.95),
                       items=left_items, size=14, color=INK,
                       line_spacing=1.35, space_after=6)
    # Right panel
    rx = Inches(0.6) + panel_w + gap
    _add_rect(slide, rx, panel_top, panel_w, panel_h, PAPER_WARM)
    _add_text(slide, rx + Inches(0.25), panel_top + Inches(0.2),
              panel_w - Inches(0.5), Inches(0.5),
              text=right_h, font=FONT_DISPLAY, size=18, bold=True, color=SCARLET)
    _add_bulleted_list(slide, rx + Inches(0.25), panel_top + Inches(0.75),
                       panel_w - Inches(0.5), panel_h - Inches(1.7),
                       items=right_items, size=14, color=INK,
                       line_spacing=1.35, space_after=6)
    if right_tail:
        _add_text(slide, rx + Inches(0.25),
                  panel_top + panel_h - Inches(0.85),
                  panel_w - Inches(0.5), Inches(0.7),
                  text=right_tail, font=FONT_DISPLAY, size=13, bold=True,
                  color=SCARLET_DARK, line_spacing=1.3)
    _footer(slide, footer_left, str(page_no))
    _set_notes(slide, notes)


def build_recap_strip(prs, *, title, subhead, cards, tail, footer_left,
                       page_no, notes):
    """cards: list of {week, h3, body}."""
    slide = _blank_slide(prs, fill=PAPER)
    _add_text(slide, Inches(0.6), Inches(0.45), Inches(12.2), Inches(0.95),
              text=title, font=FONT_DISPLAY, size=34, bold=True, color=INK,
              line_spacing=1.05)
    if subhead:
        _add_text(slide, Inches(0.6), Inches(1.45), Inches(12.2), Inches(0.4),
                  text=subhead, font=FONT_LIGHT, size=16, italic=True, color=INK_MUTED)
    grid_top = Inches(2.0); grid_h = Inches(4.0)
    gap = Inches(0.2)
    n = len(cards)
    panel_w = (Inches(12.2) - gap * (n - 1)) / n
    for i, c in enumerate(cards):
        x = Inches(0.6) + (panel_w + gap) * i
        _add_rect(slide, x, grid_top, panel_w, grid_h, PAPER_WARM)
        _add_rect(slide, x, grid_top, Inches(0.10), grid_h, SCARLET)
        _add_text(slide, x + Inches(0.25), grid_top + Inches(0.2),
                  panel_w - Inches(0.4), Inches(0.4),
                  text=c["week"].upper(), font=FONT_DISPLAY, size=13, bold=True,
                  color=SCARLET)
        _add_text(slide, x + Inches(0.25), grid_top + Inches(0.65),
                  panel_w - Inches(0.4), Inches(0.85),
                  text=c["h3"], font=FONT_DISPLAY, size=18, bold=True, color=INK,
                  line_spacing=1.1)
        _add_text(slide, x + Inches(0.25), grid_top + Inches(1.65),
                  panel_w - Inches(0.4), grid_h - Inches(1.85),
                  text=c["body"], font=FONT_BODY, size=13, color=INK_SOFT,
                  line_spacing=1.35)
    if tail:
        _add_text(slide, Inches(0.6), Inches(6.2), Inches(12.2), Inches(0.6),
                  text=tail, font=FONT_BODY, size=15, color=INK_SOFT,
                  line_spacing=1.3)
    _footer(slide, footer_left, str(page_no))
    _set_notes(slide, notes)


def build_phases(prs, *, title, phases, footer_left, page_no, notes):
    """phases: list of {pnum, h3, mode, why, items, checkpoint}"""
    slide = _blank_slide(prs, fill=PAPER)
    _add_text(slide, Inches(0.6), Inches(0.45), Inches(12.2), Inches(0.85),
              text=title, font=FONT_DISPLAY, size=32, bold=True, color=INK)
    grid_top = Inches(1.5); grid_h = Inches(5.4)
    gap = Inches(0.2)
    n = len(phases)
    pw = (Inches(12.2) - gap * (n - 1)) / n
    for i, p in enumerate(phases):
        x = Inches(0.6) + (pw + gap) * i
        _add_rect(slide, x, grid_top, pw, grid_h, PAPER_WARM)
        _add_rect(slide, x, grid_top, pw, Inches(0.08), SCARLET)
        _add_text(slide, x + Inches(0.25), grid_top + Inches(0.18),
                  pw - Inches(0.5), Inches(0.4),
                  text=p["pnum"].upper(), font=FONT_DISPLAY, size=12, bold=True,
                  color=SCARLET)
        _add_text(slide, x + Inches(0.25), grid_top + Inches(0.55),
                  pw - Inches(0.5), Inches(0.55),
                  text=p["h3"], font=FONT_DISPLAY, size=20, bold=True, color=INK)
        # Mode pill
        mode_label = p["mode"]
        is_cyborg = "cyborg" in mode_label.lower()
        pill_w = Inches(1.4)
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      x + Inches(0.25), grid_top + Inches(1.15),
                                      pill_w, Inches(0.32))
        pill.shadow.inherit = False
        pill.fill.solid()
        pill.fill.fore_color.rgb = SCARLET if is_cyborg else INK
        pill.line.fill.background()
        ptf = pill.text_frame
        ptf.margin_top = Inches(0.02); ptf.margin_bottom = Inches(0.02)
        pp = ptf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
        pr = pp.add_run(); pr.text = mode_label.upper()
        pr.font.name = FONT_DISPLAY; pr.font.size = Pt(10); pr.font.bold = True
        pr.font.color.rgb = WHITE if is_cyborg else GOLD
        # Why
        _add_text(slide, x + Inches(0.25), grid_top + Inches(1.55),
                  pw - Inches(0.5), Inches(0.7),
                  text=p["why"], font=FONT_BODY, size=12, color=INK,
                  line_spacing=1.3)
        # Items
        _add_bulleted_list(slide, x + Inches(0.25), grid_top + Inches(2.35),
                           pw - Inches(0.5), Inches(2.0),
                           items=p["items"], size=11, color=INK,
                           line_spacing=1.25, space_after=3)
        # Checkpoint
        _add_text(slide, x + Inches(0.25),
                  grid_top + grid_h - Inches(0.7),
                  pw - Inches(0.5), Inches(0.6),
                  text=p["checkpoint"], font=FONT_DISPLAY, size=11, bold=True,
                  color=SCARLET_DARK, line_spacing=1.3)
    _footer(slide, footer_left, str(page_no))
    _set_notes(slide, notes)


def build_agenda(prs, *, title, subhead, rows, footer_left, page_no, notes):
    """rows: list of (time, module, dur, mode, is_workshop)."""
    slide = _blank_slide(prs, fill=PAPER)
    _add_text(slide, Inches(0.6), Inches(0.45), Inches(12.2), Inches(0.7),
              text=title, font=FONT_DISPLAY, size=30, bold=True, color=INK)
    if subhead:
        _add_text(slide, Inches(0.6), Inches(1.15), Inches(12.2), Inches(0.4),
                  text=subhead, font=FONT_LIGHT, size=14, italic=True, color=INK_MUTED)
    n_rows = len(rows) + 1
    tbl_top = Inches(1.65); tbl_left = Inches(0.6); tbl_w = Inches(12.2)
    tbl_h = SLIDE_H - tbl_top - Inches(0.7)
    table_shape = slide.shapes.add_table(n_rows, 4, tbl_left, tbl_top, tbl_w, tbl_h)
    table = table_shape.table
    weights = [1.6, 6.0, 1.0, 1.4]
    total = sum(weights)
    for i, w in enumerate(weights):
        table.columns[i].width = int(tbl_w * (w / total))
    # Header
    for ci, htxt in enumerate(["Time", "Module", "Duration", "Mode"]):
        cell = table.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = INK
        tf = cell.text_frame
        tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = htxt
        r.font.name = FONT_DISPLAY; r.font.size = Pt(13); r.font.bold = True
        r.font.color.rgb = WHITE
    # Body
    for ri, row in enumerate(rows):
        time_, module, dur, mode, is_workshop = row
        bg = GOLD_SOFT if is_workshop else PAPER
        for ci, val in enumerate([time_, module, dur, mode]):
            cell = table.cell(ri + 1, ci)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            tf = cell.text_frame
            tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
            tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = val
            r.font.name = FONT_BODY; r.font.size = Pt(12)
            if ci == 0:
                r.font.bold = True; r.font.color.rgb = SCARLET_DARK
            elif ci == 1:
                r.font.bold = True; r.font.color.rgb = INK
            else:
                r.font.color.rgb = INK_MUTED
    _footer(slide, footer_left, str(page_no))
    _set_notes(slide, notes)


# ─────────────────────────────────────────────────────────────────────────────
# HTML parsing — extract content from the source deck
# ─────────────────────────────────────────────────────────────────────────────
def _normalize(s: str) -> str:
    return re.sub(r"\s+", " ", s.replace("\xa0", " ")).strip()


def _li_texts(parent: Tag) -> list[str]:
    out = []
    for li in parent.find_all("li", recursive=True):
        out.append(_normalize(li.get_text(" ", strip=True)))
    return out


def _direct_li_texts(ul: Tag | None) -> list[str]:
    if ul is None:
        return []
    return [_normalize(li.get_text(" ", strip=True))
            for li in ul.find_all("li", recursive=False)]


def _para_texts(parent: Tag) -> list[str]:
    return [_normalize(p.get_text(" ", strip=True))
            for p in parent.find_all("p", recursive=True)
            if _normalize(p.get_text())]


def parse_speaker_notes(soup: BeautifulSoup) -> list[list[str]]:
    """Return a list (slide-indexed from 0) of paragraph lists."""
    notes = []
    notes_root = soup.find("aside", id="notes") or soup.find("aside", class_="notes")
    if not notes_root:
        return notes
    for nc in notes_root.find_all("div", class_="note-content"):
        paragraphs = []
        for p in nc.find_all("p"):
            paragraphs.append(_normalize(p.get_text(" ", strip=True)))
        notes.append(paragraphs)
    return notes


def parse_deck(html: str):
    soup = BeautifulSoup(html, "lxml")
    slides = soup.select(".stage > section.slide")
    notes_list = parse_speaker_notes(soup)
    parsed = []
    for i, sec in enumerate(slides):
        classes = sec.get("class", [])
        kind = next((c.split("--", 1)[1] for c in classes if c.startswith("slide--")), "content")
        sid = sec.get("data-id", str(i + 1))
        notes = notes_list[i] if i < len(notes_list) else []
        # Footer left (course tag)
        footer_left = ""
        foot = sec.find("footer", class_="slide__footer")
        if foot:
            l = foot.find("div", class_="left")
            if l:
                footer_left = _normalize(l.get_text(" ", strip=True))
        parsed.append({
            "id": int(sid),
            "kind": kind,
            "section": sec,
            "notes": notes,
            "footer_left": footer_left,
        })
    return parsed


# ─────────────────────────────────────────────────────────────────────────────
# Slide-by-slide assembly (manual mapping — Week 4 has 38 slides)
# ─────────────────────────────────────────────────────────────────────────────
def _h(sec: Tag, name: str) -> str:
    el = sec.find(name)
    return _normalize(el.get_text(" ", strip=True)) if el else ""


def _subhead(sec: Tag) -> str:
    el = sec.find("p", class_="subhead")
    return _normalize(el.get_text(" ", strip=True)) if el else ""


def _cols(sec: Tag) -> list[Tag]:
    cols2 = sec.find("div", class_="cols-2")
    if not cols2:
        return []
    return cols2.find_all("div", class_="col", recursive=False)


def assemble(prs: Presentation, slides: list[dict]) -> None:
    by_id = {s["id"]: s for s in slides}

    # ── Slide 1 — Cover
    s = by_id[1]
    sec = s["section"]
    eyebrow = _normalize(sec.find("div", class_="eyebrow").get_text(" ", strip=True))
    title = _normalize(sec.find("h1").get_text(" ", strip=True))
    tagline = _normalize(sec.find("p", class_="tagline").get_text(" ", strip=True))
    meta_pairs = []
    for d in sec.find("div", class_="meta").find_all("div", recursive=False):
        # <strong>Label</strong>Value
        strong = d.find("strong")
        label = _normalize(strong.get_text(" ", strip=True)) if strong else ""
        # remove the strong tag's text from the parent text
        val = d.get_text(" ", strip=True)
        if label and val.startswith(label):
            val = val[len(label):]
        meta_pairs.append((label, _normalize(val)))
    build_cover(prs, eyebrow=eyebrow, title=title, tagline=tagline,
                meta=meta_pairs, footer_left="EDD · Week 4 Advanced Workshop · Cover",
                page_no=1, notes=s["notes"])

    # ── Slide 2 — Welcome (two-col)
    s = by_id[2]; sec = s["section"]
    cols = _cols(sec)
    build_content_two_col(
        prs, title=_h(sec, "h2"), subhead=_subhead(sec),
        left_h=_h(cols[0], "h3"), left_items=_li_texts(cols[0]),
        right_h=_h(cols[1], "h3"), right_items=_li_texts(cols[1]),
        footer_left="Module 0 · Welcome & Ground Rules", page_no=2, notes=s["notes"],
    )

    # ── Slide 3 — Recap
    s = by_id[3]; sec = s["section"]
    cards = []
    for c in sec.find("div", class_="recap-strip").find_all("div", class_="recap-card", recursive=False):
        cards.append({
            "week": _normalize(c.find("div", class_="week").get_text(" ", strip=True)),
            "h3": _normalize(c.find("h3").get_text(" ", strip=True)),
            "body": _normalize(c.find("p").get_text(" ", strip=True)),
        })
    tail = _normalize(sec.find_all("p")[-1].get_text(" ", strip=True))
    build_recap_strip(prs, title=_h(sec, "h2"), subhead=_subhead(sec),
                      cards=cards, tail=tail,
                      footer_left="Where We've Been", page_no=3, notes=s["notes"])

    # ── Slide 4 — What today is for
    s = by_id[4]; sec = s["section"]; cols = _cols(sec)
    build_content_two_col(
        prs, title=_h(sec, "h2"), subhead=_subhead(sec),
        left_h=_h(cols[0], "h3"), left_items=_li_texts(cols[0]),
        right_h=_h(cols[1], "h3"), right_items=_li_texts(cols[1]),
        footer_left="Why This Course Exists", page_no=4, notes=s["notes"],
    )

    # ── Slide 5 — Agenda table
    s = by_id[5]; sec = s["section"]
    rows = []
    for tr in sec.find("table", class_="agenda-table").find("tbody").find_all("tr"):
        is_workshop = "workshop" in (tr.get("class") or [])
        cells = tr.find_all("td")
        row = [_normalize(c.get_text(" ", strip=True)) for c in cells]
        rows.append((row[0], row[1], row[2], row[3], is_workshop))
    build_agenda(prs, title=_h(sec, "h2"), subhead=_subhead(sec),
                 rows=rows, footer_left="Agenda", page_no=5, notes=s["notes"])

    # ── Slide 6 — Module 1 divider
    s = by_id[6]; sec = s["section"]
    build_section(
        prs,
        module_tag=_normalize(sec.find("div", class_="module-tag").get_text()),
        title=_h(sec, "h1"),
        timing=_normalize(sec.find("div", class_="timing").get_text()),
        blurb=_normalize(sec.find("p", class_="blurb").get_text()),
        footer_left="Module 1 · Frontier Mapping", page_no=6, notes=s["notes"],
    )

    # ── Slide 7 — Why frontier mapping (stat + bullets)
    s = by_id[7]; sec = s["section"]
    cols = _cols(sec)
    # Build a custom split: left big stat + source; right h3 + bullets
    stat_big = cols[0].find("div", class_="stat-big")
    num_div = stat_big.find("div", class_="num")
    # Pull off any inline-styled suffix span (e.g. "<span style='font-size:120px'>pp</span>")
    # and render it at a reduced size so it doesn't read as part of the giant number.
    suffix_span = num_div.find("span")
    stat_suffix = ""
    if suffix_span is not None:
        stat_suffix = _normalize(suffix_span.get_text(" ", strip=True))
        suffix_span.extract()
    num = _normalize(num_div.get_text(" ", strip=True))
    label = _normalize(stat_big.find("div", class_="label").get_text(" ", strip=True))
    source = _normalize(stat_big.find("div", class_="source").get_text(" ", strip=True))
    right_h = _h(cols[1], "h3")
    right_items = _li_texts(cols[1])
    _build_stat_split(prs, title=_h(sec, "h2"),
                      stat_num=num, stat_suffix=stat_suffix,
                      stat_label=label, stat_source=source,
                      right_h=right_h, right_items=right_items,
                      footer_left="Module 1 · Why It Matters", page_no=7,
                      notes=s["notes"])

    # ── Slide 8 — Worked example (5-col map table)
    s = by_id[8]; sec = s["section"]
    table = sec.find("table", class_="map-table")
    headers = [_normalize(th.get_text(" ", strip=True))
               for th in table.find("thead").find_all("th")]
    body_rows = []
    for tr in table.find("tbody").find_all("tr"):
        body_rows.append([_normalize(td.get_text(" ", strip=True))
                          for td in tr.find_all("td")])
    build_table_slide(prs, title=_h(sec, "h2"), subhead=_subhead(sec),
                      headers=headers, rows=body_rows,
                      col_widths=[1.4, 3.0, 3.0, 2.6],
                      footer_left="Module 1 · Worked Example", page_no=8,
                      notes=s["notes"])

    # ── Slide 9 — What good looks like (custom 4-quadrant: 2 cols × 2 sub-headings each)
    s = by_id[9]; sec = s["section"]
    cols = _cols(sec)
    quads = []
    for col in cols:
        h3s = col.find_all("h3")
        uls = col.find_all("ul")
        for h3, ul in zip(h3s, uls):
            quads.append({
                "h3": _normalize(h3.get_text(" ", strip=True)),
                "items": _direct_li_texts(ul),
            })
    _build_quad(prs, title=_h(sec, "h2"), subhead=_subhead(sec), quads=quads,
                footer_left="Module 1 · What Good Looks Like", page_no=9,
                notes=s["notes"])

    # ── Slide 10 — Workshop: Build map
    s = by_id[10]; sec = s["section"]
    build_workshop(
        prs,
        tag=_normalize(sec.find("span", class_="activity-tag").get_text()),
        title=_h(sec, "h2"),
        timebox=_normalize(sec.find("div", class_="timebox").get_text()),
        prompt=_normalize(sec.find("div", class_="prompt-box").get_text(" ", strip=True))
            .replace("Prompt ", ""),
        what_good=_normalize(sec.find("div", class_="what-good").get_text(" ", strip=True)),
        footer_left="Module 1 · Workshop", page_no=10, notes=s["notes"],
    )

    # ── Slide 11 — Share-out
    s = by_id[11]; sec = s["section"]
    panels = sec.find("div", class_="grid").find_all("div", class_="panel", recursive=False)
    left_items = _direct_li_texts(panels[0].find("ul"))
    right_items = _direct_li_texts(panels[1].find("ul"))
    right_tail_p = panels[1].find_all("p")
    right_tail = _normalize(right_tail_p[-1].get_text(" ", strip=True)) if right_tail_p else ""
    build_shareout(
        prs, title=_h(sec, "h2"), subhead=_subhead(sec),
        left_h=_h(panels[0], "h3"), left_items=left_items,
        right_h=_h(panels[1], "h3"), right_items=right_items,
        right_tail=right_tail,
        footer_left="Module 1 · Share-Out", page_no=11, notes=s["notes"],
    )

    # ── Slide 12 — Module 1 close (pullquote)
    s = by_id[12]; sec = s["section"]
    pq = sec.find("div", class_="pullquote")
    attrib_el = pq.find("span", class_="attrib")
    attrib = _normalize(attrib_el.get_text(" ", strip=True)) if attrib_el else ""
    quote = _normalize(pq.get_text(" ", strip=True))
    if attrib:
        quote = quote.replace(attrib, "").strip()
    tail_ps = sec.find_all("p")
    tail = _normalize(tail_ps[-1].get_text(" ", strip=True)) if tail_ps else ""
    build_content_pullquote(
        prs, title=_h(sec, "h2"), quote=quote, attrib=attrib,
        footer_left="Module 1 · Takeaway", page_no=12, notes=s["notes"], tail=tail,
    )

    # ── Slide 13 — Module 2 divider
    s = by_id[13]; sec = s["section"]
    build_section(
        prs,
        module_tag=_normalize(sec.find("div", class_="module-tag").get_text()),
        title=_h(sec, "h1"),
        timing=_normalize(sec.find("div", class_="timing").get_text()),
        blurb=_normalize(sec.find("p", class_="blurb").get_text()),
        footer_left="Module 2 · Complex Build", page_no=13, notes=s["notes"],
    )

    # ── Slide 14 — Pre-build framing (pullquote + bullets)
    s = by_id[14]; sec = s["section"]
    pq = sec.find("div", class_="pullquote")
    quote = _normalize(pq.get_text(" ", strip=True)) if pq else ""
    items = _direct_li_texts(sec.find("ul"))
    _build_quote_bullets(prs, title=_h(sec, "h2"), subhead=_subhead(sec),
                          quote=quote, items=items,
                          footer_left="Module 2 · Pre-Build Framing", page_no=14,
                          notes=s["notes"])

    # ── Slide 15 — Build goal (two-col with tail)
    s = by_id[15]; sec = s["section"]; cols = _cols(sec)
    left_lists = cols[0].find_all("ul")
    left_h = _normalize(cols[0].find("h3").get_text(" ", strip=True))
    # Combine "Inputs" + "Outputs" headings into one column with sub-labels
    left_items = []
    for h3, ul in zip(cols[0].find_all("h3"), left_lists):
        left_items.append(_normalize(h3.get_text(" ", strip=True)).upper())
        left_items.extend(_direct_li_texts(ul))
    right_h = _normalize(cols[1].find("h3").get_text(" ", strip=True))
    right_items = _direct_li_texts(cols[1].find("ul"))
    right_tail_el = cols[1].find("p")
    right_tail = _normalize(right_tail_el.get_text(" ", strip=True)) if right_tail_el else ""
    build_content_two_col(
        prs, title=_h(sec, "h2"), subhead=_subhead(sec),
        left_h=left_h, left_items=left_items,
        right_h=right_h, right_items=right_items + ([right_tail] if right_tail else []),
        footer_left="Module 2 · Build Goal", page_no=15, notes=s["notes"],
    )

    # ── Slide 16 — Three phases
    s = by_id[16]; sec = s["section"]
    phases = []
    for ph in sec.find("div", class_="phases").find_all("div", class_="phase", recursive=False):
        phases.append({
            "pnum": _normalize(ph.find("div", class_="pnum").get_text(" ", strip=True)),
            "h3": _normalize(ph.find("h3").get_text(" ", strip=True)),
            "mode": _normalize(ph.find("span", class_="mode").get_text(" ", strip=True)),
            "why": _normalize(ph.find("p").get_text(" ", strip=True)),
            "items": _direct_li_texts(ph.find("ul")),
            "checkpoint": _normalize(ph.find_all("p")[-1].get_text(" ", strip=True)),
        })
    build_phases(prs, title=_h(sec, "h2"), phases=phases,
                 footer_left="Module 2 · Three Phases", page_no=16, notes=s["notes"])

    # ── Slide 17 — Cue (live build)
    s = by_id[17]; sec = s["section"]
    ring = sec.find("div", class_="ring")
    eyebrow = _normalize(ring.find("div", class_="eyebrow").get_text(" ", strip=True))
    title = _normalize(ring.find("h2").get_text(" ", strip=True))
    body_paras = [_normalize(p.get_text(" ", strip=True))
                  for p in ring.find_all("p")]
    build_cue(prs, eyebrow=eyebrow, title=title, body_paragraphs=body_paras,
              footer_left="Module 2 · Live Build — 45 min", page_no=17,
              notes=s["notes"])

    # ── Slide 18 — Welcome back debrief (bullets + pullquote)
    s = by_id[18]; sec = s["section"]
    items = _direct_li_texts(sec.find("ul"))
    pq = sec.find("div", class_="pullquote")
    quote = _normalize(pq.get_text(" ", strip=True)) if pq else ""
    _build_bullets_quote(prs, title=_h(sec, "h2"), subhead=_subhead(sec),
                         items=items, quote=quote,
                         footer_left="Module 2 · Debrief", page_no=18,
                         notes=s["notes"])

    # ── Slide 19 — Module 2 close
    s = by_id[19]; sec = s["section"]
    pq = sec.find("div", class_="pullquote")
    quote = _normalize(pq.get_text(" ", strip=True)) if pq else ""
    tail_ps = sec.find_all("p")
    tail = _normalize(tail_ps[-1].get_text(" ", strip=True)) if tail_ps else ""
    build_content_pullquote(prs, title=_h(sec, "h2"), quote=quote, attrib="",
                            footer_left="Module 2 · Takeaway", page_no=19,
                            notes=s["notes"], tail=tail)

    # ── Slide 20 — BREAK (yellow full-bleed)
    s = by_id[20]; sec = s["section"]
    title = _normalize(sec.find("h1").get_text(" ", strip=True))
    subtitle = _normalize(sec.find("div", class_="returns").get_text(" ", strip=True))
    build_break(prs, footer_left="Break", page_no=20, notes=s["notes"],
                title=title, subtitle=subtitle)

    # ── Slide 21 — Module 3 divider
    s = by_id[21]; sec = s["section"]
    build_section(
        prs,
        module_tag=_normalize(sec.find("div", class_="module-tag").get_text()),
        title=_h(sec, "h1"),
        timing=_normalize(sec.find("div", class_="timing").get_text()),
        blurb=_normalize(sec.find("p", class_="blurb").get_text()),
        footer_left="Module 3 · Debugging", page_no=21, notes=s["notes"],
    )

    # ── Slide 22 — Debugging clinic protocol (4-card ref)
    s = by_id[22]; sec = s["section"]
    cards = []
    for c in sec.find("div", class_="ref-grid").find_all("div", class_="ref-card", recursive=False):
        cards.append({
            "num": _normalize(c.find("div", class_="num").get_text(" ", strip=True)),
            "h3": _normalize(c.find("h3").get_text(" ", strip=True)),
            "body": _normalize(c.find("p").get_text(" ", strip=True)),
        })
    build_ref_grid(prs, title=_h(sec, "h2"), subhead=_subhead(sec), cards=cards,
                   footer_left="Module 3 · Clinic Protocol", page_no=22,
                   notes=s["notes"])

    # ── Slide 23 — Facilitation rules (3 panels)
    s = by_id[23]; sec = s["section"]
    panels = []
    for fp in sec.find("div", class_="facil-grid").find_all("div", class_="facil-panel", recursive=False):
        panels.append({
            "h3": _normalize(fp.find("h3").get_text(" ", strip=True)),
            "items": _direct_li_texts(fp.find("ul")),
        })
    build_facilitation(prs, title=_h(sec, "h2"), subhead=_subhead(sec),
                       panels=panels,
                       footer_left="Module 3 · Facilitation", page_no=23,
                       notes=s["notes"])

    # ── Slide 24 — Backup scenarios (2x2 col layout)
    s = by_id[24]; sec = s["section"]
    quads = []
    for cols2 in sec.find_all("div", class_="cols-2"):
        for col in cols2.find_all("div", class_="col", recursive=False):
            h3 = col.find("h3")
            ps = col.find_all("p")
            body = " ".join(_normalize(p.get_text(" ", strip=True)) for p in ps)
            quads.append({
                "h3": _normalize(h3.get_text(" ", strip=True)) if h3 else "",
                "items": [body],
            })
    _build_quad(prs, title=_h(sec, "h2"), subhead=_subhead(sec), quads=quads,
                footer_left="Module 3 · Backup Scenarios", page_no=24,
                notes=s["notes"], small=True)

    # ── Slide 25 — Synthesis (bullets + pullquote)
    s = by_id[25]; sec = s["section"]
    items = _direct_li_texts(sec.find("ul"))
    pq = sec.find("div", class_="pullquote")
    quote = _normalize(pq.get_text(" ", strip=True)) if pq else ""
    _build_bullets_quote(prs, title=_h(sec, "h2"), subhead=_subhead(sec),
                         items=items, quote=quote,
                         footer_left="Module 3 · Synthesis", page_no=25,
                         notes=s["notes"])

    # ── Slide 26 — Module 4 divider
    s = by_id[26]; sec = s["section"]
    build_section(
        prs,
        module_tag=_normalize(sec.find("div", class_="module-tag").get_text()),
        title=_h(sec, "h1"),
        timing=_normalize(sec.find("div", class_="timing").get_text()),
        blurb=_normalize(sec.find("p", class_="blurb").get_text()),
        footer_left="Module 4 · Verification & QA", page_no=26, notes=s["notes"],
    )

    # ── Slide 27 — QA reference card (5 cards, last spans 2)
    s = by_id[27]; sec = s["section"]
    cards = []
    for c in sec.find("div", class_="ref-grid").find_all("div", class_="ref-card", recursive=False):
        cards.append({
            "num": _normalize(c.find("div", class_="num").get_text(" ", strip=True)),
            "h3": _normalize(c.find("h3").get_text(" ", strip=True)),
            "body": _normalize(c.find("p").get_text(" ", strip=True)),
        })
    build_ref_grid(prs, title=_h(sec, "h2"), subhead=_subhead(sec), cards=cards,
                   footer_left="Module 4 · QA Reference", page_no=27,
                   notes=s["notes"], wide_last=(len(cards) == 5))

    # ── Slide 28 — Workshop: Find errors
    s = by_id[28]; sec = s["section"]
    build_workshop(
        prs,
        tag=_normalize(sec.find("span", class_="activity-tag").get_text()),
        title=_h(sec, "h2"),
        timebox=_normalize(sec.find("div", class_="timebox").get_text()),
        prompt=_normalize(sec.find("div", class_="prompt-box").get_text(" ", strip=True))
            .replace("Prompt ", ""),
        what_good=_normalize(sec.find("div", class_="what-good").get_text(" ", strip=True)),
        footer_left="Module 4 · Timed QA Drill", page_no=28, notes=s["notes"],
    )

    # ── Slide 29 — Module 5 divider
    s = by_id[29]; sec = s["section"]
    build_section(
        prs,
        module_tag=_normalize(sec.find("div", class_="module-tag").get_text()),
        title=_h(sec, "h1"),
        timing=_normalize(sec.find("div", class_="timing").get_text()),
        blurb=_normalize(sec.find("p", class_="blurb").get_text()),
        footer_left="Module 5 · Teaching Others", page_no=29, notes=s["notes"],
    )

    # ── Slide 30 — What you owe forward
    s = by_id[30]; sec = s["section"]; cols = _cols(sec)
    build_content_two_col(
        prs, title=_h(sec, "h2"), subhead=_subhead(sec),
        left_h=_h(cols[0], "h3"), left_items=_li_texts(cols[0]),
        right_h=_h(cols[1], "h3"), right_items=_li_texts(cols[1]),
        footer_left="Module 5 · The Owe", page_no=30, notes=s["notes"],
    )

    # ── Slide 31 — Apprentice problem
    s = by_id[31]; sec = s["section"]; cols = _cols(sec)
    build_content_two_col(
        prs, title=_h(sec, "h2"), subhead=_subhead(sec),
        left_h=_h(cols[0], "h3"), left_items=_li_texts(cols[0]),
        right_h=_h(cols[1], "h3"), right_items=_li_texts(cols[1]),
        footer_left="Module 5 · Apprentice Problem", page_no=31, notes=s["notes"],
    )

    # ── Slide 32 — Workshop: Teach-back
    s = by_id[32]; sec = s["section"]
    build_workshop(
        prs,
        tag=_normalize(sec.find("span", class_="activity-tag").get_text()),
        title=_h(sec, "h2"),
        timebox=_normalize(sec.find("div", class_="timebox").get_text()),
        prompt=_normalize(sec.find("div", class_="prompt-box").get_text(" ", strip=True))
            .replace("Prompt ", ""),
        what_good=_normalize(sec.find("div", class_="what-good").get_text(" ", strip=True)),
        footer_left="Module 5 · Teach-Back", page_no=32, notes=s["notes"],
    )

    # ── Slide 33 — Module 6 divider
    s = by_id[33]; sec = s["section"]
    build_section(
        prs,
        module_tag=_normalize(sec.find("div", class_="module-tag").get_text()),
        title=_h(sec, "h1"),
        timing=_normalize(sec.find("div", class_="timing").get_text()),
        blurb=_normalize(sec.find("p", class_="blurb").get_text()),
        footer_left="Module 6 · Playbook", page_no=33, notes=s["notes"],
    )

    # ── Slide 34 — Workflow playbook example (key/value table)
    s = by_id[34]; sec = s["section"]
    rows = []
    for tr in sec.find("table", class_="map-table").find("tbody").find_all("tr"):
        tds = tr.find_all("td")
        # Replace internal <br> with " · " for readability in PPTX cells.
        # Get text honoring line breaks via separator '\n', then collapse.
        label = _normalize(tds[0].get_text(" ", strip=True))
        # For the steps cell, join lines with " | "
        val_html = tds[1]
        val_text = val_html.get_text("\n", strip=True)
        # Normalize whitespace within lines, preserve line breaks
        val_lines = [re.sub(r"\s+", " ", l.strip()) for l in val_text.split("\n") if l.strip()]
        val = "\n".join(val_lines)
        rows.append([label, val])
    build_table_slide(
        prs, title=_h(sec, "h2"), subhead=_subhead(sec),
        headers=["Field", "Value"], rows=rows, col_widths=[2.0, 8.0],
        first_col_accent=True,
        footer_left="Module 6 · Example Playbook", page_no=34, notes=s["notes"],
    )

    # ── Slide 35 — Workshop: Build playbook
    s = by_id[35]; sec = s["section"]
    build_workshop(
        prs,
        tag=_normalize(sec.find("span", class_="activity-tag").get_text()),
        title=_h(sec, "h2"),
        timebox=_normalize(sec.find("div", class_="timebox").get_text()),
        prompt=_normalize(sec.find("div", class_="prompt-box").get_text(" ", strip=True))
            .replace("Prompt ", ""),
        what_good=_normalize(sec.find("div", class_="what-good").get_text(" ", strip=True)),
        footer_left="Module 6 · Build Playbook", page_no=35, notes=s["notes"],
    )

    # ── Slide 36 — Reflection (bullets)
    s = by_id[36]; sec = s["section"]
    items = _direct_li_texts(sec.find("ul"))
    tail_ps = sec.find_all("p")
    tail = _normalize(tail_ps[-1].get_text(" ", strip=True)) if tail_ps else ""
    build_content_bullets(
        prs, title=_h(sec, "h2"), subhead=_subhead(sec), items=items,
        footer_left="Reflection", page_no=36, notes=s["notes"], tail=tail,
    )

    # ── Slide 37 — Certification path
    s = by_id[37]; sec = s["section"]; cols = _cols(sec)
    tail_ps = sec.find_all("p")
    tail = _normalize(tail_ps[-1].get_text(" ", strip=True)) if tail_ps else ""
    build_content_two_col(
        prs, title=_h(sec, "h2"), subhead=_subhead(sec),
        left_h=_h(cols[0], "h3"), left_items=_li_texts(cols[0]),
        right_h=_h(cols[1], "h3"), right_items=_li_texts(cols[1]),
        footer_left="Certification", page_no=37, notes=s["notes"], tail=tail,
    )

    # ── Slide 38 — Week 5 preview
    s = by_id[38]; sec = s["section"]; cols = _cols(sec)
    # Right column has a closing line as a <p>
    right_close = cols[1].find("p")
    right_items = _li_texts(cols[1])
    if right_close:
        right_items.append(_normalize(right_close.get_text(" ", strip=True)))
    build_content_two_col(
        prs, title=_h(sec, "h2"), subhead=_subhead(sec),
        left_h=_h(cols[0], "h3"), left_items=_li_texts(cols[0]),
        right_h=_h(cols[1], "h3"), right_items=right_items,
        footer_left="Week 5 Preview", page_no=38, notes=s["notes"],
    )


# ─────────────────────────────────────────────────────────────────────────────
# Custom layout helpers used by assemble()
# ─────────────────────────────────────────────────────────────────────────────
def _build_stat_split(prs, *, title, stat_num, stat_label, stat_source,
                      right_h, right_items, footer_left, page_no, notes,
                      stat_suffix: str = ""):
    slide = _blank_slide(prs, fill=PAPER)
    title_lines = _estimate_visible_lines(title, 12.2, 32)
    title_h_in = max(0.85, (32 / 72.0) * 1.05 * title_lines + 0.20)
    _add_text(slide, Inches(0.6), Inches(0.45), Inches(12.2), Inches(title_h_in),
              text=title, font=FONT_DISPLAY, size=32, bold=True, color=INK,
              line_spacing=1.05)
    # Left: big stat. If a suffix span was present in the source HTML
    # (e.g. "<span style='font-size:120px'>pp</span>"), render the number
    # and suffix as two runs in the same line at distinct sizes so the
    # suffix reads as a unit, not as an extra digit pair.
    _add_rect(slide, Inches(0.6), Inches(1.85), Inches(6.0), Inches(4.8), PAPER_WARM)
    if stat_suffix:
        box = slide.shapes.add_textbox(Inches(0.85), Inches(2.0),
                                       Inches(5.5), Inches(2.2))
        box.fill.background(); box.line.fill.background()
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 0.95
        r1 = p.add_run(); r1.text = stat_num
        r1.font.name = FONT_DISPLAY; r1.font.size = Pt(120)
        r1.font.bold = True; r1.font.color.rgb = SCARLET
        r2 = p.add_run(); r2.text = stat_suffix
        r2.font.name = FONT_DISPLAY; r2.font.size = Pt(56)
        r2.font.bold = True; r2.font.color.rgb = SCARLET
    else:
        _add_text(slide, Inches(0.85), Inches(2.0), Inches(5.5), Inches(2.2),
                  text=stat_num, font=FONT_DISPLAY, size=120, bold=True,
                  color=SCARLET, line_spacing=0.95)
    _add_text(slide, Inches(0.85), Inches(4.4), Inches(5.5), Inches(1.6),
              text=stat_label, font=FONT_BODY, size=14, color=INK, line_spacing=1.35)
    _add_text(slide, Inches(0.85), Inches(6.05), Inches(5.5), Inches(0.4),
              text=stat_source, font=FONT_LIGHT, size=11, italic=True, color=INK_MUTED)
    # Right: heading + bullets
    _add_text(slide, Inches(6.9), Inches(1.85), Inches(6.0), Inches(0.5),
              text=right_h, font=FONT_DISPLAY, size=20, bold=True, color=SCARLET)
    _add_bulleted_list(slide, Inches(6.9), Inches(2.45), Inches(6.0), Inches(4.2),
                       items=right_items, size=15, color=INK,
                       line_spacing=1.3, space_after=8)
    _footer(slide, footer_left, str(page_no))
    _set_notes(slide, notes)


def _build_quad(prs, *, title, subhead, quads, footer_left, page_no, notes,
                small: bool = False):
    """2×2 grid; each cell has h3 + bullets/text."""
    slide = _blank_slide(prs, fill=PAPER)
    title_lines = _estimate_visible_lines(title, 12.2, 30)
    title_h_in = max(0.7, (30 / 72.0) * 1.05 * title_lines + 0.18)
    _add_text(slide, Inches(0.6), Inches(0.45), Inches(12.2), Inches(title_h_in),
              text=title, font=FONT_DISPLAY, size=30, bold=True, color=INK,
              line_spacing=1.05)
    sub_y = Inches(0.45 + title_h_in + 0.02)
    if subhead:
        _add_text(slide, Inches(0.6), sub_y, Inches(12.2), Inches(0.4),
                  text=subhead, font=FONT_LIGHT, size=14, italic=True, color=INK_MUTED)
        grid_top = sub_y + Inches(0.50)
    else:
        grid_top = sub_y + Inches(0.10)
    grid_h = Inches(7.0) - grid_top - Inches(0.65)
    gap = Inches(0.2)
    cell_w = (Inches(12.2) - gap) / 2
    n = len(quads)
    rows = (n + 1) // 2
    cell_h = (grid_h - gap * (rows - 1)) / rows
    item_size = 11 if small else 13
    for i, q in enumerate(quads):
        row = i // 2; col = i % 2
        x = Inches(0.6) + (cell_w + gap) * col
        y = grid_top + (cell_h + gap) * row
        _add_rect(slide, x, y, cell_w, cell_h, PAPER_WARM)
        _add_rect(slide, x, y, cell_w, Inches(0.06), SCARLET)
        _add_text(slide, x + Inches(0.25), y + Inches(0.18),
                  cell_w - Inches(0.5), Inches(0.5),
                  text=q["h3"], font=FONT_DISPLAY, size=16, bold=True, color=INK)
        if len(q["items"]) == 1 and len(q["items"][0]) > 80 and not q["items"][0].lstrip().startswith(("•", "-")):
            # Treat as paragraph
            _add_text(slide, x + Inches(0.25), y + Inches(0.7),
                      cell_w - Inches(0.5), cell_h - Inches(0.85),
                      text=q["items"][0], font=FONT_BODY, size=item_size,
                      color=INK_SOFT, line_spacing=1.35)
        else:
            _add_bulleted_list(slide, x + Inches(0.25), y + Inches(0.7),
                               cell_w - Inches(0.5), cell_h - Inches(0.85),
                               items=q["items"], size=item_size, color=INK,
                               line_spacing=1.3, space_after=4)
    _footer(slide, footer_left, str(page_no))
    _set_notes(slide, notes)


def _build_quote_bullets(prs, *, title, subhead, quote, items, footer_left,
                          page_no, notes):
    slide = _blank_slide(prs, fill=PAPER)
    title_lines = _estimate_visible_lines(title, 12.2, 32)
    title_h_in = max(0.85, (32 / 72.0) * 1.05 * title_lines + 0.18)
    _add_text(slide, Inches(0.6), Inches(0.45), Inches(12.2), Inches(title_h_in),
              text=title, font=FONT_DISPLAY, size=32, bold=True, color=INK,
              line_spacing=1.05)
    sub_y = Inches(0.45 + title_h_in + 0.02)
    if subhead:
        _add_text(slide, Inches(0.6), sub_y, Inches(12.2), Inches(0.4),
                  text=subhead, font=FONT_LIGHT, size=14, italic=True, color=INK_MUTED)
        qy = sub_y + Inches(0.55)
    else:
        qy = sub_y + Inches(0.10)
    qh = Inches(2.0)
    _add_rect(slide, Inches(0.6), qy, Inches(12.2), qh, PAPER_WARM)
    _add_rect(slide, Inches(0.6), qy, Inches(0.12), qh, SCARLET)
    _add_text(slide, Inches(1.0), qy + Inches(0.25), Inches(11.6), qh - Inches(0.5),
              text=quote, font=FONT_LIGHT, size=20, italic=False, color=INK,
              anchor="middle", line_spacing=1.3)
    _add_bulleted_list(slide, Inches(0.6), Inches(4.1), Inches(12.2), Inches(2.7),
                       items=items, size=16, color=INK, line_spacing=1.3,
                       space_after=8)
    _footer(slide, footer_left, str(page_no))
    _set_notes(slide, notes)


def _build_bullets_quote(prs, *, title, subhead, items, quote, footer_left,
                          page_no, notes):
    slide = _blank_slide(prs, fill=PAPER)
    title_lines = _estimate_visible_lines(title, 12.2, 32)
    title_h_in = max(0.85, (32 / 72.0) * 1.05 * title_lines + 0.18)
    _add_text(slide, Inches(0.6), Inches(0.45), Inches(12.2), Inches(title_h_in),
              text=title, font=FONT_DISPLAY, size=32, bold=True, color=INK,
              line_spacing=1.05)
    sub_y = Inches(0.45 + title_h_in + 0.02)
    if subhead:
        _add_text(slide, Inches(0.6), sub_y, Inches(12.2), Inches(0.4),
                  text=subhead, font=FONT_LIGHT, size=14, italic=True, color=INK_MUTED)
        bullets_y = sub_y + Inches(0.55)
    else:
        bullets_y = sub_y + Inches(0.10)
    _add_bulleted_list(slide, Inches(0.6), bullets_y, Inches(12.2),
                       Inches(4.55) - bullets_y,
                       items=items, size=16, color=INK, line_spacing=1.3,
                       space_after=6)
    qy = Inches(4.65); qh = Inches(2.1)
    _add_rect(slide, Inches(0.6), qy, Inches(12.2), qh, PAPER_WARM)
    _add_rect(slide, Inches(0.6), qy, Inches(0.12), qh, SCARLET)
    _add_text(slide, Inches(1.0), qy + Inches(0.2), Inches(11.6), qh - Inches(0.4),
              text=quote, font=FONT_LIGHT, size=18, italic=False, color=INK,
              anchor="middle", line_spacing=1.3)
    _footer(slide, footer_left, str(page_no))
    _set_notes(slide, notes)


# ─────────────────────────────────────────────────────────────────────────────
# Verification
# ─────────────────────────────────────────────────────────────────────────────
def verify(out_path: Path, expected_count: int, expected_notes: list[list[str]]) -> None:
    prs = Presentation(out_path)
    actual = len(prs.slides)
    assert actual == expected_count, (
        f"Slide count mismatch: expected {expected_count}, got {actual}")
    # Slide size
    assert prs.slide_width == Inches(13.333), f"Slide width wrong: {prs.slide_width}"
    assert prs.slide_height == Inches(7.5), f"Slide height wrong: {prs.slide_height}"
    # Notes populated and editable text frames present
    populated = 0
    text_frames = 0
    for i, sl in enumerate(prs.slides):
        ntf = sl.notes_slide.notes_text_frame
        ntext = ntf.text.strip()
        if ntext:
            populated += 1
        for shp in sl.shapes:
            if shp.has_text_frame and shp.text_frame.text.strip():
                text_frames += 1
    expected_notes_n = sum(1 for n in expected_notes if n)
    assert populated >= expected_notes_n, (
        f"Notes populated {populated}/{expected_notes_n} expected")
    assert text_frames > 0, "No editable text frames found"
    print(f"  ✓ slides: {actual}")
    print(f"  ✓ size:   13.333 × 7.5 in (16:9)")
    print(f"  ✓ notes:  {populated}/{actual} slides")
    print(f"  ✓ shapes: {text_frames} editable text frames across deck")


# ─────────────────────────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────────────────────────
def main():
    if not SRC_HTML.exists():
        print(f"ERROR: source HTML not found at {SRC_HTML}", file=sys.stderr)
        sys.exit(1)
    html = SRC_HTML.read_text(encoding="utf-8")
    parsed = parse_deck(html)
    expected = 38
    if len(parsed) != expected:
        print(f"ERROR: expected {expected} slides in source, got {len(parsed)}",
              file=sys.stderr)
        sys.exit(1)
    print(f"Parsed {len(parsed)} slides from {SRC_HTML.relative_to(ROOT)}")

    prs = _new_presentation()
    assemble(prs, parsed)
    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PPTX)
    print(f"Wrote {OUT_PPTX.relative_to(ROOT)}")
    verify(OUT_PPTX, expected, [s["notes"] for s in parsed])
    print("Done.")


if __name__ == "__main__":
    main()
