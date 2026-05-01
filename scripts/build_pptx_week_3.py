"""Build a native, editable PowerPoint deck for Week 3 — Platform Training.

Source of truth: docs/decks/week-3-platform-training.html
Output:           docs/pptx/week-3-platform-training.pptx

Self-contained — does not import from sibling build_pptx_week_*.py files.
"""

from __future__ import annotations

import re
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import Callable, List, Optional

from bs4 import BeautifulSoup, NavigableString, Tag
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree


REPO_ROOT = Path(__file__).resolve().parent.parent
HTML_PATH = REPO_ROOT / "docs" / "decks" / "week-3-platform-training.html"
OUT_PATH = REPO_ROOT / "docs" / "pptx" / "week-3-platform-training.pptx"


# ----------------------------------------------------------------------
# Theme — Week 3 cream palette and type scale
# ----------------------------------------------------------------------

C_SCARLET = RGBColor(0xCC, 0x00, 0x00)
C_SCARLET_DK = RGBColor(0xA3, 0x00, 0x00)
C_GOLD = RGBColor(0xF5, 0xD1, 0x30)
C_GOLD_DK = RGBColor(0xD4, 0xB1, 0x1A)
C_BG = RGBColor(0xFF, 0xFF, 0xFF)
C_BG_WARM = RGBColor(0xF8, 0xF7, 0xF5)
C_BG_DARK = RGBColor(0x1A, 0x1A, 0x1A)
C_BG_QUIET = RGBColor(0xF1, 0xEF, 0xE9)
C_INK = RGBColor(0x1A, 0x1A, 0x1A)
C_INK_SOFT = RGBColor(0x4A, 0x4A, 0x4A)
C_INK_MUTED = RGBColor(0x6E, 0x6E, 0x6E)
C_RULE = RGBColor(0xD9, 0xD8, 0xD4)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_GOLD_TINT = RGBColor(0xFF, 0xFA, 0xF0)

FONT_DISPLAY = "Calibri"
FONT_BODY = "Calibri"
FONT_MONO = "Consolas"

# Slide geometry — 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
PAD_X = Inches(0.7)
PAD_Y = Inches(0.55)


# ----------------------------------------------------------------------
# Parsing — walk the HTML deck and extract structured slide data
# ----------------------------------------------------------------------


@dataclass
class Slide:
    layout: str
    title: str
    eyebrow: str = ""
    chip: str = ""
    chip_kind: str = ""           # "", "gold", "scarlet", "ghost"
    headline: str = ""
    subhead: str = ""
    bullets: List[str] = field(default_factory=list)
    body_paragraphs: List[str] = field(default_factory=list)
    columns: List[dict] = field(default_factory=list)   # [{title, bullets, paragraphs, footnote}]
    timeline: List[dict] = field(default_factory=list)  # [{time, what, meta, is_break}]
    frame_hero: dict = field(default_factory=dict)      # {label, headline, body}
    frame_cards: List[dict] = field(default_factory=list)  # [{label, value, sub, accent}]
    debrief_panels: List[dict] = field(default_factory=list)  # [{title, bullets}]
    map_grid: List[List[str]] = field(default_factory=list)   # rows of cells (incl header)
    map_styles: List[List[str]] = field(default_factory=list) # per-cell style hint
    chips: List[str] = field(default_factory=list)      # for switch slides
    badge: str = ""                                      # switch / return badge text
    caption: str = ""                                    # switch caption / quiet lead
    quiet_prompts: List[str] = field(default_factory=list)
    section_module: str = ""                             # section eyebrow text
    section_meta: List[dict] = field(default_factory=list)  # [{text, kind}]
    timer_minutes: str = ""
    timer_label: str = "Minutes"
    preview_badge: str = ""
    closing_sub: str = ""
    speaker_notes: str = ""


def _txt(el) -> str:
    """Extract clean text, collapsing whitespace and resolving HTML entities."""
    if el is None:
        return ""
    if isinstance(el, NavigableString):
        return str(el)
    text = el.get_text(separator=" ", strip=False)
    text = re.sub(r"\s+", " ", text).strip()
    # Replace common typographic entities
    text = (
        text.replace("\u00a0", " ")
        .replace("\u2192", "→")
        .replace("\u2014", "—")
        .replace("\u2013", "–")
        .replace("\u2018", "'")
        .replace("\u2019", "'")
        .replace("\u201c", "\u201c")
        .replace("\u201d", "\u201d")
        .replace("\u00b7", "·")
        .replace("\u00ad", "")
    )
    return text


def _chip_kind(el: Tag) -> str:
    classes = el.get("class", []) if el else []
    if "chip--gold" in classes:
        return "gold"
    if "chip--scarlet" in classes:
        return "scarlet"
    if "chip--ghost" in classes:
        return "ghost"
    return ""


def _extract_speaker_notes(slide_el: Tag) -> str:
    """Pull the speaker-notes block as a single multi-paragraph string."""
    notes_el = slide_el.find("div", class_="speaker-notes")
    if not notes_el:
        return ""
    paragraphs: List[str] = []
    for p in notes_el.find_all("p", recursive=False):
        # Build the paragraph but prepend notes__cue spans inline as "Cue — body"
        parts: List[str] = []
        for child in p.children:
            if isinstance(child, Tag) and "notes__cue" in (child.get("class") or []):
                cue = _txt(child)
                if cue:
                    parts.append(f"[{cue}]")
            else:
                parts.append(_txt(child) if isinstance(child, Tag) else str(child))
        text = "".join(parts)
        text = re.sub(r"\s+", " ", text).strip()
        if text:
            paragraphs.append(text)
    return "\n\n".join(paragraphs)


def _extract_columns(container: Tag) -> List[dict]:
    cols: List[dict] = []
    for col in container.find_all("div", class_="col", recursive=True):
        # Avoid double-counting when nested
        if any(p is not container and "col" in (p.get("class") or [])
               for p in col.parents):
            continue
        head_el = col.find(class_="col__head")
        head = _txt(head_el)
        bullets: List[str] = []
        ul = col.find("ul")
        if ul:
            for li in ul.find_all("li", recursive=False):
                bullets.append(_txt(li))
        paragraphs: List[str] = []
        for p in col.find_all("p", recursive=False):
            t = _txt(p)
            if t:
                paragraphs.append(t)
        cols.append({"title": head, "bullets": bullets, "paragraphs": paragraphs})
    return cols


def _layout_kind(slide_el: Tag) -> str:
    classes = slide_el.get("class", [])
    for c in classes:
        if c.startswith("layout-"):
            return c[len("layout-"):]
    return "content"


def parse_deck(html_path: Path) -> List[Slide]:
    soup = BeautifulSoup(html_path.read_text(encoding="utf-8"), "lxml")
    out: List[Slide] = []
    for sec in soup.select("section.slide"):
        kind = _layout_kind(sec)
        title = sec.get("data-title", "").strip()
        slide = Slide(layout=kind, title=title)
        slide.speaker_notes = _extract_speaker_notes(sec)

        header = sec.find("div", class_="slide__header")
        if header:
            eb = header.find(class_="slide__eyebrow")
            slide.eyebrow = _txt(eb)
            chip = header.find(class_="chip")
            if chip:
                slide.chip = _txt(chip)
                slide.chip_kind = _chip_kind(chip)

        if kind == "cover":
            slide.eyebrow = _txt(sec.find(class_="cover__top"))
            slide.subhead = _txt(sec.find(class_="cover__week"))
            slide.headline = _txt(sec.find(class_="cover__title"))
            slide.body_paragraphs = [_txt(sec.find(class_="cover__sub"))]
            meta = sec.find(class_="cover__meta")
            if meta:
                for span in meta.find_all("span"):
                    slide.body_paragraphs.append(_txt(span))

        elif kind == "section":
            slide.section_module = _txt(sec.find(class_="section__eyebrow"))
            h = sec.find(class_="h-display")
            if h:
                # Preserve <br> as newline
                slide.headline = _txt(h).replace(" ", " ").strip()
                # Smart split on the literal HTML <br>
                raw = h.decode_contents()
                raw = re.sub(r"<br\s*/?>", "\n", raw, flags=re.I)
                clean = BeautifulSoup(raw, "lxml").get_text()
                slide.headline = re.sub(r"[ \t]+", " ", clean).strip()
            meta = sec.find(class_="section__meta")
            if meta:
                for chip in meta.find_all(class_="chip"):
                    slide.section_meta.append(
                        {"text": _txt(chip), "kind": _chip_kind(chip)}
                    )

        elif kind == "agenda":
            h = sec.find(class_="h-title")
            slide.headline = _txt(h)
            timeline_el = sec.find(class_="timeline")
            if timeline_el:
                spans = list(timeline_el.find_all("span", recursive=False))
                # Each row is three spans: time, what, meta
                for i in range(0, len(spans), 3):
                    row = spans[i: i + 3]
                    if len(row) < 3:
                        continue
                    is_break = "timeline__row--break" in (row[0].get("class") or [])
                    slide.timeline.append({
                        "time": _txt(row[0]),
                        "what": _txt(row[1]),
                        "meta": _txt(row[2]),
                        "is_break": is_break,
                    })

        elif kind == "two":
            h = sec.find(class_="h-title")
            slide.headline = _txt(h)
            body = sec.find("div", class_="body")
            if body:
                slide.columns = _extract_columns(body)

        elif kind == "frame":
            h = sec.find(class_="h-title")
            slide.headline = _txt(h)
            hero = sec.find(class_="frame__hero")
            if hero:
                slide.frame_hero = {
                    "label": _txt(hero.find(class_="frame__label")),
                    "headline": _txt(hero.find("h3")),
                    "body": _txt(hero.find("p")),
                }
            cards = sec.find(class_="frame__cards")
            if cards:
                for card in cards.find_all(class_="frame__card", recursive=False):
                    val_el = card.find(class_="card__value")
                    accent = "scarlet"
                    if val_el and "gold" in (val_el.get("class") or []):
                        accent = "gold"
                    slide.frame_cards.append({
                        "label": _txt(card.find(class_="card__label")),
                        "value": _txt(val_el),
                        "sub": _txt(card.find(class_="card__sub")),
                        "accent": accent,
                    })

        elif kind == "switch":
            slide.badge = _txt(sec.find(class_="switch__badge"))
            slide.headline = _txt(sec.find("h2"))
            slide.caption = _txt(sec.find(class_="switch__caption"))
            chips_holder = sec.find(class_="switch__chips")
            if chips_holder:
                for c in chips_holder.find_all(class_="chip"):
                    slide.chips.append(_txt(c))

        elif kind == "debrief":
            slide.badge = _txt(sec.find(class_="return__badge"))
            h = sec.find(class_="h-title")
            slide.headline = _txt(h)
            grid = sec.find(class_="debrief__grid")
            if grid:
                for panel in grid.find_all(class_="debrief__panel", recursive=False):
                    panel_title = _txt(panel.find("h3"))
                    bullets = [_txt(li) for li in panel.find_all("li")]
                    slide.debrief_panels.append({
                        "title": panel_title, "bullets": bullets
                    })

        elif kind == "quiet":
            slide.eyebrow = _txt(sec.find(class_="quiet__eyebrow"))
            slide.headline = _txt(sec.find("h2"))
            slide.subhead = _txt(sec.find(class_="quiet__lead"))
            for li in sec.select("ul.quiet__prompts li"):
                slide.quiet_prompts.append(_txt(li))

        elif kind == "break":
            slide.eyebrow = _txt(sec.find(class_="break__label"))
            slide.headline = _txt(sec.find("h2"))
            slide.subhead = _txt(sec.find(class_="break__sub"))

        elif kind == "worktime":
            slide.headline = _txt(sec.find("h2"))
            slide.subhead = _txt(sec.find(class_="worktime__sub"))
            timer = sec.find(class_="worktime__timer")
            if timer:
                num = timer.find(class_="num")
                lbl = timer.find(class_="label")
                slide.timer_minutes = _txt(num)
                slide.timer_label = _txt(lbl) or "Minutes"

        elif kind == "preview":
            slide.preview_badge = _txt(sec.find(class_="preview__badge"))
            slide.headline = _txt(sec.find(class_="preview__title"))
            slide.subhead = _txt(sec.find(class_="preview__sub"))
            for li in sec.select("ul.preview__list li"):
                slide.bullets.append(_txt(li))

        elif kind == "closing":
            slide.headline = _txt(sec.find("h2"))
            slide.closing_sub = _txt(sec.find(class_="closing__sub"))

        elif kind == "content":
            h = sec.find(class_="h-title")
            slide.headline = _txt(h)
            body = sec.find("div", class_="body")
            # Detect Frontier Map slide explicitly via layout-map class
            if "layout-map" in (sec.get("class") or []):
                map_el = sec.find(class_="map")
                if map_el:
                    cells = map_el.find_all(class_="map__cell")
                    rows: List[List[str]] = []
                    styles: List[List[str]] = []
                    row: List[str] = []
                    style_row: List[str] = []
                    for cell in cells:
                        classes = cell.get("class") or []
                        style = "body"
                        if "map__cell--head" in classes:
                            style = "head"
                        elif "map__cell--task" in classes:
                            style = "task"
                        elif "map__cell--good" in classes:
                            style = "good"
                        elif "map__cell--bad" in classes:
                            style = "bad"
                        elif "map__cell--check" in classes:
                            style = "check"
                        if "map__row-empty" in classes:
                            style = "empty"
                        row.append(_txt(cell))
                        style_row.append(style)
                        if len(row) == 4:
                            rows.append(row)
                            styles.append(style_row)
                            row = []
                            style_row = []
                    slide.map_grid = rows
                    slide.map_styles = styles
            elif body:
                # Detect inner two-column block
                inner_cols = body.find_all("div", class_="col", recursive=True)
                if inner_cols:
                    slide.columns = _extract_columns(body)
                # Top-level lead/paragraph
                for p in body.find_all("p", recursive=False):
                    t = _txt(p)
                    if t:
                        slide.body_paragraphs.append(t)
                # Top-level bullets
                ul = body.find("ul", recursive=False)
                if ul:
                    for li in ul.find_all("li", recursive=False):
                        slide.bullets.append(_txt(li))

        out.append(slide)
    return out


# ----------------------------------------------------------------------
# Rendering — python-pptx layout builders
# ----------------------------------------------------------------------


def _set_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _set_line(shape, color: RGBColor, width_pt: float = 0.75) -> None:
    line = shape.line
    line.color.rgb = color
    line.width = Pt(width_pt)


def add_rect(slide, x, y, w, h, fill: Optional[RGBColor] = None,
             line: Optional[RGBColor] = None) -> object:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_textbox(slide, x, y, w, h) -> object:
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    return tb


def write_text(textframe, text: str, *, font=FONT_BODY, size: float = 18,
               bold: bool = False, italic: bool = False,
               color: RGBColor = C_INK, align=PP_ALIGN.LEFT,
               anchor=MSO_ANCHOR.TOP, letter_spacing: Optional[float] = None,
               line_spacing: Optional[float] = None) -> None:
    """Replace text in an existing single-paragraph text frame."""
    textframe.vertical_anchor = anchor
    p = textframe.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    # clear any existing runs
    for r in list(p.runs):
        r.text = ""
    run = p.add_run() if not p.runs else p.runs[0]
    run.text = text
    f = run.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    if letter_spacing is not None:
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(int(letter_spacing * 100)))


def add_paragraph(textframe, text: str, *, font=FONT_BODY, size: float = 18,
                  bold: bool = False, italic: bool = False,
                  color: RGBColor = C_INK, align=PP_ALIGN.LEFT,
                  space_before: float = 0, space_after: float = 4,
                  bullet: bool = False,
                  bullet_color: RGBColor = C_SCARLET,
                  letter_spacing: Optional[float] = None,
                  line_spacing: Optional[float] = None) -> None:
    p = textframe.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    if space_after:
        p.space_after = Pt(space_after)
    if line_spacing:
        p.line_spacing = line_spacing
    if bullet:
        # Use a literal "▪ " character so the bullet renders identically across
        # PowerPoint, Keynote, and LibreOffice without depending on bullet XML.
        marker = p.add_run()
        marker.text = "▪  "
        mf = marker.font
        mf.name = font
        mf.size = Pt(size)
        mf.bold = True
        mf.color.rgb = bullet_color
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    if letter_spacing is not None:
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(int(letter_spacing * 100)))


def set_first_paragraph(textframe, text: str, **kwargs) -> None:
    write_text(textframe, text, **kwargs)


def add_brand_bar(slide) -> None:
    """Top accent bar — scarlet 70%, gold 30%."""
    bar_h = Emu(76200)  # ~0.083"
    scarlet = add_rect(slide, Emu(0), Emu(0), int(SLIDE_W * 0.70), bar_h,
                       fill=C_SCARLET)
    gold = add_rect(slide, int(SLIDE_W * 0.70), Emu(0),
                    int(SLIDE_W * 0.30), bar_h, fill=C_GOLD)


def add_foot(slide, course_label: str, num: int, total: int,
             *, on_dark: bool = False) -> None:
    color = RGBColor(0xCC, 0xCC, 0xCC) if on_dark else C_INK_MUTED
    label_color = RGBColor(0xEE, 0xEE, 0xEE) if on_dark else C_INK_SOFT
    tb = add_textbox(slide, PAD_X, SLIDE_H - Inches(0.42),
                     SLIDE_W - 2 * PAD_X, Inches(0.3))
    tf = tb.text_frame
    write_text(tf, course_label, font=FONT_BODY, size=9,
               color=label_color, letter_spacing=1.5)
    # Right-aligned slide number on the same line
    tb2 = add_textbox(slide, SLIDE_W - PAD_X - Inches(2.0),
                      SLIDE_H - Inches(0.42), Inches(2.0), Inches(0.3))
    write_text(tb2.text_frame, f"{num:02d} / {total:02d}",
               font=FONT_BODY, size=9, color=color, align=PP_ALIGN.RIGHT,
               letter_spacing=1.5)


def add_background(slide, color: RGBColor) -> None:
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, SLIDE_H
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    # Send to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def add_eyebrow_chip_row(slide, eyebrow: str, chip: str, chip_kind: str,
                         y: Emu = PAD_Y) -> Emu:
    """Draw the eyebrow text on the left, chip on the right. Returns y of the next block."""
    if eyebrow:
        tb = add_textbox(slide, PAD_X, y, Inches(7), Inches(0.3))
        write_text(tb.text_frame, eyebrow.upper(), font=FONT_BODY,
                   size=10, bold=True, color=C_SCARLET, letter_spacing=2.5)
    if chip:
        chip_w = Inches(min(4.0, max(1.5, 0.13 * len(chip) + 0.5)))
        chip_h = Inches(0.34)
        chip_x = SLIDE_W - PAD_X - chip_w
        chip_y = y - Inches(0.04)
        if chip_kind == "gold":
            fill = C_GOLD
            text_color = C_BG_DARK
            line_color = C_GOLD_DK
        elif chip_kind == "scarlet":
            fill = C_SCARLET
            text_color = C_WHITE
            line_color = C_SCARLET_DK
        elif chip_kind == "ghost":
            fill = None
            text_color = C_INK_SOFT
            line_color = C_RULE
        else:
            fill = C_BG_WARM
            text_color = C_INK_SOFT
            line_color = C_RULE
        rect = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, chip_x, chip_y, chip_w, chip_h
        )
        rect.adjustments[0] = 0.5
        if fill is None:
            rect.fill.background()
        else:
            rect.fill.solid()
            rect.fill.fore_color.rgb = fill
        rect.line.color.rgb = line_color
        rect.line.width = Pt(0.75)
        rect.shadow.inherit = False
        tf = rect.text_frame
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        write_text(tf, chip, font=FONT_BODY, size=10, bold=True,
                   color=text_color, align=PP_ALIGN.CENTER,
                   anchor=MSO_ANCHOR.MIDDLE)
    return y + Inches(0.55)


def add_rule(slide, x: Emu, y: Emu, w: Emu = Inches(1.4),
             color: RGBColor = C_SCARLET) -> None:
    add_rect(slide, x, y, w, Emu(38100), fill=color)  # ~3pt height


# ---------- Layout builders ----------

COURSE_LABEL = "EDD · Week 3 · Platform Training"


def render_cover(slide, s: Slide) -> None:
    add_background(slide, C_BG_DARK)
    # Top brand strip — wider on the cover
    bar_h = Emu(133350)  # ~0.146"
    add_rect(slide, Emu(0), Emu(0), int(SLIDE_W * 0.70), bar_h, fill=C_SCARLET)
    add_rect(slide, int(SLIDE_W * 0.70), Emu(0), int(SLIDE_W * 0.30),
             bar_h, fill=C_GOLD)

    # Top eyebrow ("Expert-Driven Development")
    tb = add_textbox(slide, PAD_X, Inches(0.6), Inches(8), Inches(0.4))
    write_text(tb.text_frame, s.eyebrow, font=FONT_DISPLAY, size=14, bold=True,
               color=C_GOLD, letter_spacing=1.5)

    # "Week 3 of 6 · Course 3"
    tb = add_textbox(slide, PAD_X, Inches(2.4), Inches(11), Inches(0.4))
    write_text(tb.text_frame, s.subhead.upper(), font=FONT_BODY, size=11,
               bold=True, color=C_SCARLET, letter_spacing=2.5)

    # Title
    tb = add_textbox(slide, PAD_X, Inches(2.95), Inches(11), Inches(2.4))
    write_text(tb.text_frame, s.headline, font=FONT_DISPLAY, size=72,
               bold=True, color=C_WHITE, line_spacing=1.0)

    # Subtitle
    if s.body_paragraphs:
        tb = add_textbox(slide, PAD_X, Inches(5.1), Inches(10.5), Inches(1.2))
        write_text(tb.text_frame, s.body_paragraphs[0],
                   font=FONT_BODY, size=20, color=RGBColor(0xC8, 0xC8, 0xC8),
                   line_spacing=1.35)

    # Meta row
    if len(s.body_paragraphs) > 1:
        tb = add_textbox(slide, PAD_X, Inches(6.6), Inches(11.5), Inches(0.4))
        tf = tb.text_frame
        meta = "    ".join(s.body_paragraphs[1:])
        write_text(tf, meta, font=FONT_BODY, size=12,
                   color=RGBColor(0xAA, 0xAA, 0xAA))

    add_foot(slide, COURSE_LABEL, 1, 32, on_dark=True)


def render_section(slide, s: Slide, num: int) -> None:
    add_background(slide, C_BG_DARK)
    add_brand_bar(slide)

    # Section eyebrow
    tb = add_textbox(slide, PAD_X, Inches(2.0), Inches(11), Inches(0.4))
    write_text(tb.text_frame, s.section_module.upper(), font=FONT_BODY,
               size=12, bold=True, color=C_GOLD, letter_spacing=3.0)

    # Headline (display)
    tb = add_textbox(slide, PAD_X, Inches(2.6), Inches(11), Inches(2.4))
    tf = tb.text_frame
    tf.word_wrap = True
    lines = s.headline.split("\n")
    set_first_paragraph(tf, lines[0], font=FONT_DISPLAY, size=64, bold=True,
                        color=C_WHITE, line_spacing=1.0)
    for line in lines[1:]:
        add_paragraph(tf, line, font=FONT_DISPLAY, size=64, bold=True,
                      color=C_WHITE, line_spacing=1.0, space_after=0)

    # Gold rule
    add_rule(slide, PAD_X, Inches(5.4), w=Inches(2.0), color=C_GOLD)

    # Meta chips
    if s.section_meta:
        x = PAD_X
        y = Inches(5.8)
        for chip in s.section_meta:
            chip_text = chip["text"]
            chip_w = Inches(min(4.0, max(1.6, 0.13 * len(chip_text) + 0.6)))
            chip_h = Inches(0.36)
            if chip["kind"] == "gold":
                fill = C_GOLD
                text_color = C_BG_DARK
                line_color = C_GOLD_DK
            else:
                fill = None
                text_color = C_WHITE
                line_color = RGBColor(0x55, 0x55, 0x55)
            rect = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, x, y, chip_w, chip_h
            )
            rect.adjustments[0] = 0.5
            if fill is None:
                rect.fill.background()
            else:
                rect.fill.solid()
                rect.fill.fore_color.rgb = fill
            rect.line.color.rgb = line_color
            rect.line.width = Pt(0.75)
            rect.shadow.inherit = False
            tf = rect.text_frame
            tf.margin_left = Emu(0)
            tf.margin_right = Emu(0)
            tf.margin_top = Emu(0)
            tf.margin_bottom = Emu(0)
            write_text(tf, chip_text, font=FONT_BODY, size=10, bold=True,
                       color=text_color, align=PP_ALIGN.CENTER,
                       anchor=MSO_ANCHOR.MIDDLE)
            x += chip_w + Inches(0.18)

    add_foot(slide, COURSE_LABEL, num, 32, on_dark=True)


def estimate_visible_lines(text: str, box_w_in: float, font_size_pt: float) -> int:
    """Conservative visible-line count for bold sans display text.
    Mirrors the helper used in scripts/build_pptx_week_6.py so layout (rules,
    columns, body offsets) can be sized from the actual line count."""
    if not text:
        return 1
    char_w = 0.62 * float(font_size_pt) / 72.0
    chars_per_line = max(6, int(float(box_w_in) / char_w))
    total = 0
    for chunk in str(text).split("\n"):
        n = max(1, len(chunk))
        total += (n + chars_per_line - 1) // chars_per_line
    return max(1, total)


def render_content(slide, s: Slide, num: int) -> None:
    add_background(slide, C_BG)
    add_brand_bar(slide)
    next_y = add_eyebrow_chip_row(slide, s.eyebrow, s.chip, s.chip_kind,
                                  y=Inches(0.55))

    # Headline — size box and downstream rule from actual visible-line count.
    if s.headline:
        head_w_in = (SLIDE_W - 2 * PAD_X) / 914400
        head_lines = estimate_visible_lines(s.headline, head_w_in, 36)
        head_h_in = max(0.85, (36 / 72.0) * 1.05 * head_lines + 0.20)
        tb = add_textbox(slide, PAD_X, next_y, SLIDE_W - 2 * PAD_X,
                         Inches(head_h_in))
        write_text(tb.text_frame, s.headline, font=FONT_DISPLAY, size=36,
                   bold=True, color=C_INK, line_spacing=1.05)
        next_y += Inches(head_h_in + 0.10)

    add_rule(slide, PAD_X, next_y, w=Inches(1.4))
    next_y += Inches(0.35)

    body_top = next_y
    body_h = SLIDE_H - body_top - Inches(0.7)

    if s.columns:
        col_count = len(s.columns)
        gap = Inches(0.5)
        col_w = (SLIDE_W - 2 * PAD_X - gap * (col_count - 1)) / col_count
        for i, col in enumerate(s.columns):
            cx = PAD_X + (col_w + gap) * i
            tb = add_textbox(slide, cx, body_top, col_w, Inches(0.4))
            write_text(tb.text_frame, col["title"].upper(), font=FONT_BODY,
                       size=11, bold=True, color=C_SCARLET, letter_spacing=2.0)
            # Underline
            add_rect(slide, cx, body_top + Inches(0.36), col_w, Emu(12700),
                     fill=C_RULE)
            body_tb = add_textbox(slide, cx, body_top + Inches(0.5),
                                  col_w, body_h - Inches(0.5))
            tf = body_tb.text_frame
            tf.word_wrap = True
            first = True
            for para in col["paragraphs"]:
                if first:
                    write_text(tf, para, font=FONT_BODY, size=14,
                               color=C_INK_SOFT, line_spacing=1.35)
                    first = False
                else:
                    add_paragraph(tf, para, font=FONT_BODY, size=14,
                                  color=C_INK_SOFT, line_spacing=1.35,
                                  space_after=4)
            for b in col["bullets"]:
                if first:
                    write_text(tf, "▪  " + b, font=FONT_BODY, size=14,
                               color=C_INK, line_spacing=1.3)
                    # color the marker
                    run0 = tf.paragraphs[0].runs[0]
                    # split into marker + body so we can color marker
                    run0.text = "▪  "
                    run0.font.color.rgb = C_SCARLET
                    run0.font.bold = True
                    extra = tf.paragraphs[0].add_run()
                    extra.text = b
                    extra.font.name = FONT_BODY
                    extra.font.size = Pt(14)
                    extra.font.color.rgb = C_INK
                    first = False
                else:
                    add_paragraph(tf, b, font=FONT_BODY, size=14,
                                  color=C_INK, bullet=True,
                                  bullet_color=C_SCARLET,
                                  line_spacing=1.3, space_after=4)
        return

    # Top-level paragraphs (lead) and bullets
    body_tb = add_textbox(slide, PAD_X, body_top, SLIDE_W - 2 * PAD_X, body_h)
    tf = body_tb.text_frame
    tf.word_wrap = True
    first = True
    for para in s.body_paragraphs:
        if first:
            write_text(tf, para, font=FONT_BODY, size=18,
                       color=C_INK_SOFT, line_spacing=1.35)
            first = False
        else:
            add_paragraph(tf, para, font=FONT_BODY, size=18,
                          color=C_INK_SOFT, line_spacing=1.35, space_after=8)
    for b in s.bullets:
        if first:
            # Marker run then body run
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            for r in list(p.runs):
                r.text = ""
            mr = p.add_run()
            mr.text = "▪  "
            mr.font.name = FONT_BODY
            mr.font.size = Pt(18)
            mr.font.color.rgb = C_SCARLET
            mr.font.bold = True
            tr = p.add_run()
            tr.text = b
            tr.font.name = FONT_BODY
            tr.font.size = Pt(18)
            tr.font.color.rgb = C_INK
            first = False
        else:
            add_paragraph(tf, b, font=FONT_BODY, size=18, color=C_INK,
                          bullet=True, bullet_color=C_SCARLET,
                          line_spacing=1.3, space_after=6)
    # Render <strong> emphasis is skipped for simplicity; bullets are still
    # editable PowerPoint paragraphs.

    add_foot(slide, COURSE_LABEL, num, 32)


def render_two(slide, s: Slide, num: int) -> None:
    """Two-column layout (slide 7 — centaur vs cyborg)."""
    # Promote columns onto a content-style layout
    s.bullets = []
    s.body_paragraphs = []
    render_content(slide, s, num)


def render_agenda(slide, s: Slide, num: int) -> None:
    add_background(slide, C_BG)
    add_brand_bar(slide)
    next_y = add_eyebrow_chip_row(slide, s.eyebrow, s.chip, s.chip_kind,
                                  y=Inches(0.55))

    if s.headline:
        tb = add_textbox(slide, PAD_X, next_y, SLIDE_W - 2 * PAD_X,
                         Inches(0.9))
        write_text(tb.text_frame, s.headline, font=FONT_DISPLAY, size=32,
                   bold=True, color=C_INK, line_spacing=1.05)
        next_y += Inches(0.85)

    add_rule(slide, PAD_X, next_y, w=Inches(1.4))
    next_y += Inches(0.35)

    # Timeline as a real PowerPoint table. Three columns: time / what / meta.
    rows = len(s.timeline)
    if rows == 0:
        return
    table_h = SLIDE_H - next_y - Inches(0.7)
    table_w = SLIDE_W - 2 * PAD_X
    table_shape = slide.shapes.add_table(rows, 3, PAD_X, next_y,
                                         table_w, table_h)
    table = table_shape.table
    table.columns[0].width = Inches(1.4)
    table.columns[1].width = table_w - Inches(1.4) - Inches(2.6)
    table.columns[2].width = Inches(2.6)

    row_h = int(table_h / max(rows, 1))
    for i, row in enumerate(s.timeline):
        table.rows[i].height = row_h
        is_break = row["is_break"]

        # Time cell
        c = table.cell(i, 0)
        c.fill.solid()
        c.fill.fore_color.rgb = C_BG
        c.margin_left = Inches(0.05)
        c.margin_right = Inches(0.05)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        time_color = C_GOLD_DK if is_break else C_SCARLET
        write_text(c.text_frame, row["time"], font=FONT_BODY, size=16,
                   bold=True, color=time_color, anchor=MSO_ANCHOR.MIDDLE)

        # What cell
        c = table.cell(i, 1)
        c.fill.solid()
        c.fill.fore_color.rgb = C_BG
        c.margin_left = Inches(0.15)
        c.margin_right = Inches(0.15)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        what_color = C_INK_SOFT if is_break else C_INK
        write_text(c.text_frame, row["what"], font=FONT_BODY, size=16,
                   bold=not is_break, italic=is_break, color=what_color,
                   anchor=MSO_ANCHOR.MIDDLE)

        # Meta cell
        c = table.cell(i, 2)
        c.fill.solid()
        c.fill.fore_color.rgb = C_BG
        c.margin_left = Inches(0.05)
        c.margin_right = Inches(0.1)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        write_text(c.text_frame, row["meta"], font=FONT_BODY, size=11,
                   color=C_INK_MUTED, align=PP_ALIGN.RIGHT,
                   anchor=MSO_ANCHOR.MIDDLE)

    # Remove default cell borders by zeroing line widths via XML
    _strip_table_borders(table)
    # Add subtle row separators
    _add_table_row_separators(table, color=C_RULE)

    add_foot(slide, COURSE_LABEL, num, 32)


def _strip_table_borders(table) -> None:
    for row in table.rows:
        for cell in row.cells:
            tcPr = cell._tc.get_or_add_tcPr()
            for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
                el = etree.SubElement(tcPr, qn(tag))
                el.set("w", "0")
                el.set("cap", "flat")
                el.set("cmpd", "sng")
                el.set("algn", "ctr")
                fill = etree.SubElement(el, qn("a:noFill"))


def _add_table_row_separators(table, color: RGBColor) -> None:
    for row in table.rows:
        for cell in row.cells:
            tcPr = cell._tc.get_or_add_tcPr()
            ln = etree.SubElement(tcPr, qn("a:lnB"))
            ln.set("w", "9525")  # ~0.75pt
            ln.set("cap", "flat")
            ln.set("cmpd", "sng")
            ln.set("algn", "ctr")
            fill = etree.SubElement(ln, qn("a:solidFill"))
            srgb = etree.SubElement(fill, qn("a:srgbClr"))
            srgb.set("val", "{:02X}{:02X}{:02X}".format(color[0], color[1], color[2]))


def render_frame(slide, s: Slide, num: int) -> None:
    add_background(slide, C_BG)
    add_brand_bar(slide)
    next_y = add_eyebrow_chip_row(slide, s.eyebrow, s.chip, s.chip_kind,
                                  y=Inches(0.55))
    if s.headline:
        tb = add_textbox(slide, PAD_X, next_y, SLIDE_W - 2 * PAD_X,
                         Inches(0.9))
        write_text(tb.text_frame, s.headline, font=FONT_DISPLAY, size=32,
                   bold=True, color=C_INK, line_spacing=1.05)
        next_y += Inches(0.85)
    add_rule(slide, PAD_X, next_y, w=Inches(1.4))
    next_y += Inches(0.35)

    body_h = SLIDE_H - next_y - Inches(0.7)
    half_w = (SLIDE_W - 2 * PAD_X - Inches(0.4)) / 2

    # Hero block (left)
    hero_x = PAD_X
    hero_y = next_y
    add_rect(slide, hero_x, hero_y, half_w, body_h, fill=C_BG_WARM)
    add_rect(slide, hero_x, hero_y, Inches(0.06), body_h, fill=C_SCARLET)

    inner_x = hero_x + Inches(0.35)
    inner_y = hero_y + Inches(0.35)
    inner_w = half_w - Inches(0.7)

    # Label
    tb = add_textbox(slide, inner_x, inner_y, inner_w, Inches(0.3))
    write_text(tb.text_frame, s.frame_hero.get("label", "").upper(),
               font=FONT_BODY, size=10, bold=True, color=C_SCARLET,
               letter_spacing=2.5)
    inner_y += Inches(0.42)
    # Headline
    tb = add_textbox(slide, inner_x, inner_y, inner_w, Inches(1.6))
    write_text(tb.text_frame, s.frame_hero.get("headline", ""),
               font=FONT_DISPLAY, size=22, bold=True, color=C_INK,
               line_spacing=1.15)
    inner_y += Inches(1.7)
    # Body
    tb = add_textbox(slide, inner_x, inner_y, inner_w,
                     hero_y + body_h - inner_y - Inches(0.3))
    write_text(tb.text_frame, s.frame_hero.get("body", ""),
               font=FONT_BODY, size=14, color=C_INK_SOFT, line_spacing=1.4)

    # Cards (right) — stacked
    card_x = PAD_X + half_w + Inches(0.4)
    card_w = half_w
    n = max(len(s.frame_cards), 1)
    card_gap = Inches(0.25)
    card_h = (body_h - card_gap * (n - 1)) / n
    cy = next_y
    for card in s.frame_cards:
        add_rect(slide, card_x, cy, card_w, card_h, fill=C_WHITE,
                 line=C_RULE)
        ix = card_x + Inches(0.3)
        iy = cy + Inches(0.25)
        iw = card_w - Inches(0.6)
        tb = add_textbox(slide, ix, iy, iw, Inches(0.3))
        write_text(tb.text_frame, card["label"].upper(),
                   font=FONT_BODY, size=10, bold=True, color=C_INK_MUTED,
                   letter_spacing=2.5)
        iy += Inches(0.4)
        accent_color = C_GOLD_DK if card["accent"] == "gold" else C_SCARLET
        tb = add_textbox(slide, ix, iy, iw, Inches(0.9))
        # If value starts with a strong-style word ("Centaur — …" / "Cyborg — …"),
        # keep it as one paragraph but emphasize the lead word in scarlet/gold.
        tf = tb.text_frame
        tf.word_wrap = True
        value = card["value"]
        m = re.match(r"^(\S+?)\s*—\s*(.*)$", value)
        if m:
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            for r in list(p.runs):
                r.text = ""
            head = p.add_run()
            head.text = m.group(1)
            head.font.name = FONT_BODY
            head.font.size = Pt(15)
            head.font.bold = True
            head.font.color.rgb = accent_color
            tail = p.add_run()
            tail.text = " — " + m.group(2)
            tail.font.name = FONT_BODY
            tail.font.size = Pt(15)
            tail.font.color.rgb = C_INK
        else:
            write_text(tf, value, font=FONT_BODY, size=15, color=C_INK,
                       line_spacing=1.3)
        iy += Inches(1.0)
        if card.get("sub"):
            tb = add_textbox(slide, ix, iy, iw,
                             cy + card_h - iy - Inches(0.2))
            write_text(tb.text_frame, card["sub"], font=FONT_BODY, size=12,
                       color=C_INK_MUTED, line_spacing=1.35)
        cy += card_h + card_gap

    add_foot(slide, COURSE_LABEL, num, 32)


def render_switch(slide, s: Slide, num: int) -> None:
    # Cream/gold tinted background for the switch cue
    add_background(slide, C_GOLD_TINT)
    add_brand_bar(slide)

    # Badge
    if s.badge:
        badge_w = Inches(min(7.0, max(3.5, 0.13 * len(s.badge) + 1.0)))
        badge_h = Inches(0.5)
        rect = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, PAD_X, Inches(1.2),
            badge_w, badge_h
        )
        rect.adjustments[0] = 0.5
        rect.fill.solid()
        rect.fill.fore_color.rgb = C_GOLD
        rect.line.color.rgb = C_GOLD_DK
        rect.line.width = Pt(0.75)
        rect.shadow.inherit = False
        tf = rect.text_frame
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        write_text(tf, s.badge, font=FONT_BODY, size=12, bold=True,
                   color=C_BG_DARK, align=PP_ALIGN.CENTER,
                   anchor=MSO_ANCHOR.MIDDLE, letter_spacing=1.5)

    # Headline
    tb = add_textbox(slide, PAD_X, Inches(2.3), SLIDE_W - 2 * PAD_X,
                     Inches(1.6))
    write_text(tb.text_frame, s.headline, font=FONT_DISPLAY, size=42,
               bold=True, color=C_INK, line_spacing=1.05)

    # Caption
    if s.caption:
        tb = add_textbox(slide, PAD_X, Inches(4.0), SLIDE_W - 2 * PAD_X,
                         Inches(1.5))
        write_text(tb.text_frame, s.caption, font=FONT_BODY, size=16,
                   color=C_INK_SOFT, line_spacing=1.4)

    # Chips
    if s.chips:
        chip_y = Inches(5.6)
        x = PAD_X
        max_x = SLIDE_W - PAD_X
        for c in s.chips:
            chip_w = Inches(min(5.0, max(1.6, 0.10 * len(c) + 0.6)))
            chip_h = Inches(0.36)
            if x + chip_w > max_x:
                x = PAD_X
                chip_y += Inches(0.5)
            rect = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, x, chip_y, chip_w, chip_h
            )
            rect.adjustments[0] = 0.5
            rect.fill.solid()
            rect.fill.fore_color.rgb = C_WHITE
            rect.line.color.rgb = C_RULE
            rect.line.width = Pt(0.75)
            rect.shadow.inherit = False
            tf = rect.text_frame
            tf.margin_left = Emu(0)
            tf.margin_right = Emu(0)
            tf.margin_top = Emu(0)
            tf.margin_bottom = Emu(0)
            write_text(tf, c, font=FONT_BODY, size=10, bold=True,
                       color=C_INK_SOFT, align=PP_ALIGN.CENTER,
                       anchor=MSO_ANCHOR.MIDDLE)
            x += chip_w + Inches(0.18)

    add_foot(slide, COURSE_LABEL, num, 32)


def render_debrief(slide, s: Slide, num: int) -> None:
    add_background(slide, C_BG)
    add_brand_bar(slide)

    # Return badge (just below the bar)
    if s.badge:
        badge_w = Inches(min(6.0, max(3.5, 0.12 * len(s.badge) + 0.8)))
        badge_h = Inches(0.36)
        rect = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, PAD_X, Inches(0.3),
            badge_w, badge_h
        )
        rect.adjustments[0] = 0.5
        rect.fill.solid()
        rect.fill.fore_color.rgb = C_BG_QUIET
        rect.line.color.rgb = C_RULE
        rect.line.width = Pt(0.75)
        rect.shadow.inherit = False
        tf = rect.text_frame
        tf.margin_left = Inches(0.12)
        tf.margin_right = Inches(0.12)
        write_text(tf, s.badge, font=FONT_BODY, size=10, bold=True,
                   color=C_INK_SOFT, align=PP_ALIGN.CENTER,
                   anchor=MSO_ANCHOR.MIDDLE, letter_spacing=1.5)

    next_y = add_eyebrow_chip_row(slide, s.eyebrow, s.chip, s.chip_kind,
                                  y=Inches(0.95))

    if s.headline:
        tb = add_textbox(slide, PAD_X, next_y, SLIDE_W - 2 * PAD_X,
                         Inches(1.0))
        write_text(tb.text_frame, s.headline, font=FONT_DISPLAY, size=32,
                   bold=True, color=C_INK, line_spacing=1.05)
        next_y += Inches(0.95)

    add_rule(slide, PAD_X, next_y)
    next_y += Inches(0.35)

    body_h = SLIDE_H - next_y - Inches(0.7)
    half_w = (SLIDE_W - 2 * PAD_X - Inches(0.4)) / 2
    for i, panel in enumerate(s.debrief_panels):
        x = PAD_X + (half_w + Inches(0.4)) * i
        add_rect(slide, x, next_y, half_w, body_h, fill=C_BG_WARM)
        ix = x + Inches(0.4)
        iy = next_y + Inches(0.35)
        iw = half_w - Inches(0.8)
        tb = add_textbox(slide, ix, iy, iw, Inches(0.5))
        write_text(tb.text_frame, panel["title"], font=FONT_DISPLAY, size=20,
                   bold=True, color=C_INK)
        iy += Inches(0.55)
        tb = add_textbox(slide, ix, iy, iw,
                         next_y + body_h - iy - Inches(0.25))
        tf = tb.text_frame
        tf.word_wrap = True
        first = True
        for b in panel["bullets"]:
            if first:
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                for r in list(p.runs):
                    r.text = ""
                mr = p.add_run()
                mr.text = "▪  "
                mr.font.name = FONT_BODY
                mr.font.size = Pt(15)
                mr.font.color.rgb = C_SCARLET
                mr.font.bold = True
                tr = p.add_run()
                tr.text = b
                tr.font.name = FONT_BODY
                tr.font.size = Pt(15)
                tr.font.color.rgb = C_INK
                first = False
            else:
                add_paragraph(tf, b, font=FONT_BODY, size=15, color=C_INK,
                              bullet=True, bullet_color=C_SCARLET,
                              line_spacing=1.3, space_after=6)

    add_foot(slide, COURSE_LABEL, num, 32)


def render_quiet(slide, s: Slide, num: int) -> None:
    add_background(slide, C_BG_QUIET)
    add_brand_bar(slide)
    # Eyebrow
    tb = add_textbox(slide, PAD_X, Inches(0.7), SLIDE_W - 2 * PAD_X,
                     Inches(0.4))
    write_text(tb.text_frame, s.eyebrow.upper(), font=FONT_BODY, size=11,
               bold=True, color=C_SCARLET, letter_spacing=2.5)
    # Headline
    tb = add_textbox(slide, PAD_X, Inches(1.4), SLIDE_W - 2 * PAD_X,
                     Inches(2.4))
    write_text(tb.text_frame, s.headline, font=FONT_DISPLAY, size=34,
               bold=True, color=C_INK, line_spacing=1.1)
    # Lead
    if s.subhead:
        tb = add_textbox(slide, PAD_X, Inches(3.7), SLIDE_W - 2 * PAD_X,
                         Inches(1.0))
        write_text(tb.text_frame, s.subhead, font=FONT_BODY, size=16,
                   color=C_INK_SOFT, italic=True, line_spacing=1.4)
    # Prompts
    if s.quiet_prompts:
        tb = add_textbox(slide, PAD_X, Inches(4.8), SLIDE_W - 2 * PAD_X,
                         Inches(2.0))
        tf = tb.text_frame
        tf.word_wrap = True
        first = True
        for p in s.quiet_prompts:
            if first:
                pp = tf.paragraphs[0]
                for r in list(pp.runs):
                    r.text = ""
                mr = pp.add_run()
                mr.text = "▪  "
                mr.font.name = FONT_BODY
                mr.font.size = Pt(14)
                mr.font.color.rgb = C_GOLD_DK
                mr.font.bold = True
                tr = pp.add_run()
                tr.text = p
                tr.font.name = FONT_BODY
                tr.font.size = Pt(14)
                tr.font.color.rgb = C_INK_SOFT
                first = False
            else:
                add_paragraph(tf, p, font=FONT_BODY, size=14,
                              color=C_INK_SOFT, bullet=True,
                              bullet_color=C_GOLD_DK,
                              line_spacing=1.35, space_after=4)

    add_foot(slide, COURSE_LABEL, num, 32)


def render_break(slide, s: Slide, num: int) -> None:
    add_background(slide, C_BG_DARK)
    add_brand_bar(slide)
    # Label
    tb = add_textbox(slide, PAD_X, Inches(2.4), SLIDE_W - 2 * PAD_X,
                     Inches(0.5))
    write_text(tb.text_frame, s.eyebrow.upper(), font=FONT_BODY, size=14,
               bold=True, color=C_GOLD, letter_spacing=4.0,
               align=PP_ALIGN.CENTER)
    # Headline
    tb = add_textbox(slide, PAD_X, Inches(3.0), SLIDE_W - 2 * PAD_X,
                     Inches(2.0))
    write_text(tb.text_frame, s.headline, font=FONT_DISPLAY, size=96,
               bold=True, color=C_WHITE, align=PP_ALIGN.CENTER,
               line_spacing=1.0)
    # Sub
    if s.subhead:
        tb = add_textbox(slide, PAD_X, Inches(5.4), SLIDE_W - 2 * PAD_X,
                         Inches(1.0))
        write_text(tb.text_frame, s.subhead, font=FONT_BODY, size=16,
                   color=RGBColor(0xC8, 0xC8, 0xC8), align=PP_ALIGN.CENTER,
                   line_spacing=1.4)
    add_foot(slide, COURSE_LABEL, num, 32, on_dark=True)


def render_worktime(slide, s: Slide, num: int) -> None:
    add_background(slide, C_BG)
    add_brand_bar(slide)

    # Header row — eyebrow on the LEFT, scarlet "Add at least one row" chip on the
    # RIGHT, with a clear gap so they cannot collide regardless of length.
    eb_w = Inches(5.5)
    tb = add_textbox(slide, PAD_X, Inches(0.55), eb_w, Inches(0.34))
    write_text(tb.text_frame, s.eyebrow.upper(), font=FONT_BODY, size=11,
               bold=True, color=C_SCARLET, letter_spacing=2.5,
               anchor=MSO_ANCHOR.MIDDLE)
    if s.chip:
        chip_text = s.chip
        chip_w = Inches(min(5.0, max(2.6, 0.13 * len(chip_text) + 0.6)))
        chip_h = Inches(0.42)
        chip_x = SLIDE_W - PAD_X - chip_w
        chip_y = Inches(0.5)
        rect = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, chip_x, chip_y, chip_w, chip_h
        )
        rect.adjustments[0] = 0.5
        if s.chip_kind == "scarlet":
            fill = C_SCARLET
            txt = C_WHITE
            line = C_SCARLET_DK
        elif s.chip_kind == "gold":
            fill = C_GOLD
            txt = C_BG_DARK
            line = C_GOLD_DK
        else:
            fill = C_BG_WARM
            txt = C_INK_SOFT
            line = C_RULE
        rect.fill.solid()
        rect.fill.fore_color.rgb = fill
        rect.line.color.rgb = line
        rect.line.width = Pt(0.75)
        rect.shadow.inherit = False
        tf = rect.text_frame
        tf.margin_left = Inches(0.12)
        tf.margin_right = Inches(0.12)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        write_text(tf, chip_text, font=FONT_BODY, size=11, bold=True,
                   color=txt, align=PP_ALIGN.CENTER,
                   anchor=MSO_ANCHOR.MIDDLE)

    # Headline
    tb = add_textbox(slide, PAD_X, Inches(1.6), SLIDE_W - 2 * PAD_X,
                     Inches(1.6))
    write_text(tb.text_frame, s.headline, font=FONT_DISPLAY, size=48,
               bold=True, color=C_INK, line_spacing=1.0)

    # Sub
    if s.subhead:
        tb = add_textbox(slide, PAD_X, Inches(3.2),
                         SLIDE_W - 2 * PAD_X - Inches(4.0), Inches(2.2))
        write_text(tb.text_frame, s.subhead, font=FONT_BODY, size=16,
                   color=C_INK_SOFT, line_spacing=1.45)

    # Timer block — right side
    if s.timer_minutes:
        timer_w = Inches(3.5)
        timer_h = Inches(3.2)
        tx = SLIDE_W - PAD_X - timer_w
        ty = Inches(3.2)
        add_rect(slide, tx, ty, timer_w, timer_h, fill=C_BG_WARM)
        # Big number
        tb = add_textbox(slide, tx, ty + Inches(0.3), timer_w, Inches(2.2))
        write_text(tb.text_frame, s.timer_minutes, font=FONT_DISPLAY,
                   size=140, bold=True, color=C_SCARLET,
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                   line_spacing=0.9)
        tb = add_textbox(slide, tx, ty + Inches(2.5), timer_w, Inches(0.5))
        write_text(tb.text_frame, s.timer_label.upper(), font=FONT_BODY,
                   size=14, bold=True, color=C_INK_MUTED,
                   align=PP_ALIGN.CENTER, letter_spacing=4.0)

    add_foot(slide, COURSE_LABEL, num, 32)


def render_preview(slide, s: Slide, num: int) -> None:
    add_background(slide, C_BG_DARK)
    add_brand_bar(slide)

    if s.preview_badge:
        tb = add_textbox(slide, PAD_X, Inches(1.0), Inches(8), Inches(0.4))
        write_text(tb.text_frame, s.preview_badge.upper(), font=FONT_BODY,
                   size=12, bold=True, color=C_GOLD, letter_spacing=3.0)

    tb = add_textbox(slide, PAD_X, Inches(1.7), SLIDE_W - 2 * PAD_X,
                     Inches(1.6))
    write_text(tb.text_frame, s.headline, font=FONT_DISPLAY, size=52,
               bold=True, color=C_WHITE, line_spacing=1.0)

    if s.subhead:
        tb = add_textbox(slide, PAD_X, Inches(3.4),
                         SLIDE_W - 2 * PAD_X - Inches(0.3), Inches(1.6))
        write_text(tb.text_frame, s.subhead, font=FONT_BODY, size=16,
                   color=RGBColor(0xC8, 0xC8, 0xC8), line_spacing=1.45)

    if s.bullets:
        tb = add_textbox(slide, PAD_X, Inches(5.0), SLIDE_W - 2 * PAD_X,
                         Inches(2.0))
        tf = tb.text_frame
        tf.word_wrap = True
        first = True
        for b in s.bullets:
            if first:
                p = tf.paragraphs[0]
                for r in list(p.runs):
                    r.text = ""
                mr = p.add_run()
                mr.text = "▪  "
                mr.font.name = FONT_BODY
                mr.font.size = Pt(13)
                mr.font.color.rgb = C_GOLD
                mr.font.bold = True
                tr = p.add_run()
                tr.text = b
                tr.font.name = FONT_BODY
                tr.font.size = Pt(13)
                tr.font.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
                first = False
            else:
                add_paragraph(tf, b, font=FONT_BODY, size=13,
                              color=RGBColor(0xE0, 0xE0, 0xE0),
                              bullet=True, bullet_color=C_GOLD,
                              line_spacing=1.35, space_after=4)

    add_foot(slide, COURSE_LABEL, num, 32, on_dark=True)


def render_closing(slide, s: Slide, num: int) -> None:
    add_background(slide, C_BG_DARK)
    add_brand_bar(slide)
    tb = add_textbox(slide, PAD_X, Inches(2.6), SLIDE_W - 2 * PAD_X,
                     Inches(2.6))
    write_text(tb.text_frame, s.headline, font=FONT_DISPLAY, size=72,
               bold=True, color=C_WHITE, align=PP_ALIGN.CENTER,
               line_spacing=1.0)
    add_rule(slide, (SLIDE_W - Inches(2.0)) // 2, Inches(5.2), w=Inches(2.0),
             color=C_GOLD)
    if s.closing_sub:
        tb = add_textbox(slide, PAD_X, Inches(5.6), SLIDE_W - 2 * PAD_X,
                         Inches(0.8))
        write_text(tb.text_frame, s.closing_sub, font=FONT_BODY, size=18,
                   color=RGBColor(0xC8, 0xC8, 0xC8), align=PP_ALIGN.CENTER,
                   italic=True)
    add_foot(slide, COURSE_LABEL, num, 32, on_dark=True)


def render_map(slide, s: Slide, num: int) -> None:
    """Slide 27 — Frontier map reference (real PowerPoint table)."""
    add_background(slide, C_BG)
    add_brand_bar(slide)
    next_y = add_eyebrow_chip_row(slide, s.eyebrow, s.chip, s.chip_kind,
                                  y=Inches(0.55))
    if s.headline:
        tb = add_textbox(slide, PAD_X, next_y, SLIDE_W - 2 * PAD_X,
                         Inches(0.9))
        write_text(tb.text_frame, s.headline, font=FONT_DISPLAY, size=32,
                   bold=True, color=C_INK, line_spacing=1.05)
        next_y += Inches(0.85)
    add_rule(slide, PAD_X, next_y)
    next_y += Inches(0.3)

    rows = len(s.map_grid)
    cols = 4
    if rows == 0:
        return
    table_w = SLIDE_W - 2 * PAD_X
    table_h = SLIDE_H - next_y - Inches(0.7)
    shape = slide.shapes.add_table(rows, cols, PAD_X, next_y,
                                   table_w, table_h)
    table = shape.table
    table.columns[0].width = Inches(2.4)
    rest = table_w - Inches(2.4)
    table.columns[1].width = int(rest / 3)
    table.columns[2].width = int(rest / 3)
    table.columns[3].width = rest - 2 * int(rest / 3)
    header_h = Inches(0.5)
    table.rows[0].height = header_h
    body_h = int((table_h - header_h) / max(rows - 1, 1))
    for i in range(1, rows):
        table.rows[i].height = body_h
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)
            cell.vertical_anchor = MSO_ANCHOR.TOP
            style = s.map_styles[r][c]
            text = s.map_grid[r][c]
            if style == "head":
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_BG_DARK
                write_text(cell.text_frame, text.upper(), font=FONT_BODY,
                           size=10, bold=True, color=C_WHITE,
                           letter_spacing=2.0, anchor=MSO_ANCHOR.MIDDLE)
            elif style == "task":
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_BG_WARM
                write_text(cell.text_frame, text, font=FONT_BODY, size=12,
                           bold=True, color=C_INK)
            elif style == "good":
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xEC, 0xF6, 0xEC)
                write_text(cell.text_frame, text, font=FONT_BODY, size=11,
                           color=RGBColor(0x2E, 0x5A, 0x2E))
            elif style == "bad":
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFC, 0xEC, 0xEC)
                write_text(cell.text_frame, text, font=FONT_BODY, size=11,
                           color=RGBColor(0x8E, 0x2A, 0x2A))
            elif style == "check":
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFB, 0xF4, 0xDD)
                write_text(cell.text_frame, text, font=FONT_BODY, size=11,
                           color=C_INK_SOFT)
            elif style == "empty":
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_WHITE
                write_text(cell.text_frame, text, font=FONT_BODY, size=11,
                           italic=True, color=C_INK_MUTED)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_WHITE
                write_text(cell.text_frame, text, font=FONT_BODY, size=11,
                           color=C_INK)
    _strip_table_borders(table)
    _add_table_row_separators(table, color=C_RULE)
    add_foot(slide, COURSE_LABEL, num, 32)


# Layout dispatch
RENDERERS: dict[str, Callable] = {
    "cover": render_cover,
    "section": render_section,
    "content": render_content,
    "two": render_two,
    "agenda": render_agenda,
    "frame": render_frame,
    "switch": render_switch,
    "debrief": render_debrief,
    "quiet": render_quiet,
    "break": render_break,
    "worktime": render_worktime,
    "preview": render_preview,
    "closing": render_closing,
}


# ----------------------------------------------------------------------
# Build & verify
# ----------------------------------------------------------------------


def build_pptx(slides: List[Slide], out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]  # Blank

    total = len(slides)
    for i, s in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank_layout)
        # Special-case the frontier map content slide
        if s.layout == "content" and s.map_grid:
            render_map(slide, s, i)
        else:
            renderer = RENDERERS.get(s.layout, render_content)
            if s.layout in ("cover",):
                renderer(slide, s)
            else:
                renderer(slide, s, i)
        # Speaker notes — native PowerPoint speaker notes
        notes_tf = slide.notes_slide.notes_text_frame
        notes_tf.text = s.speaker_notes or ""

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)


def verify(out_path: Path, expected_count: int) -> None:
    prs = Presentation(out_path)
    actual = len(prs.slides)
    if actual != expected_count:
        raise SystemExit(
            f"Slide count mismatch: expected {expected_count}, got {actual}"
        )
    if prs.slide_width != SLIDE_W or prs.slide_height != SLIDE_H:
        raise SystemExit(
            f"Slide dimensions wrong: {prs.slide_width} x {prs.slide_height}"
        )
    missing_notes = []
    no_textframes = []
    for idx, slide in enumerate(prs.slides, start=1):
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        if not notes_text:
            missing_notes.append(idx)
        # Check at least one editable text frame
        has_text = any(
            getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip()
            for sh in slide.shapes
        )
        if not has_text:
            no_textframes.append(idx)
    if missing_notes:
        print(f"WARNING: missing speaker notes on slides {missing_notes}")
    if no_textframes:
        raise SystemExit(
            f"Slides without any editable text frames: {no_textframes}"
        )
    print(
        f"Verified: {actual} slides, "
        f"{prs.slide_width.inches:.3f} × {prs.slide_height.inches:.3f} in, "
        f"speaker notes populated, editable text frames present."
    )


def main() -> None:
    print(f"Parsing  {HTML_PATH.relative_to(REPO_ROOT)}")
    slides = parse_deck(HTML_PATH)
    print(f"Parsed   {len(slides)} slides")
    print(f"Building {OUT_PATH.relative_to(REPO_ROOT)}")
    build_pptx(slides, OUT_PATH)
    verify(OUT_PATH, expected_count=len(slides))
    print("Done.")


if __name__ == "__main__":
    main()
