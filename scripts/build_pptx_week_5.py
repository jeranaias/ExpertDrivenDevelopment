#!/usr/bin/env python3
"""Build a native, editable PowerPoint deck for Week 5 (Supervisor Orientation).

Source of truth: docs/decks/week-5-supervisor.html.
Output: docs/pptx/week-5-supervisor.pptx.

Self-contained — does not import from sibling build_pptx_week_*.py scripts.
"""

from __future__ import annotations

import re
import sys
from dataclasses import dataclass
from pathlib import Path
from typing import Callable, Iterable

from bs4 import BeautifulSoup, NavigableString, Tag
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------

REPO_ROOT = Path(__file__).resolve().parent.parent
SOURCE_HTML = REPO_ROOT / "docs" / "decks" / "week-5-supervisor.html"
OUTPUT_PPTX = REPO_ROOT / "docs" / "pptx" / "week-5-supervisor.pptx"

# ---------------------------------------------------------------------------
# Theme — Week 5 wine-red palette + type
# ---------------------------------------------------------------------------

SCARLET = RGBColor(0xCC, 0x00, 0x00)
SCARLET_DARK = RGBColor(0xA3, 0x00, 0x00)
SCARLET_DEEP = RGBColor(0x5C, 0x00, 0x00)
GOLD = RGBColor(0xF5, 0xD1, 0x30)
GOLD_DARK = RGBColor(0xD4, 0xB1, 0x1A)

INK = RGBColor(0x1A, 0x1A, 0x1A)
INK_2 = RGBColor(0x4A, 0x4A, 0x4A)
INK_3 = RGBColor(0x6E, 0x6E, 0x6E)
PAPER = RGBColor(0xF8, 0xF7, 0xF5)
PAPER_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RULE_LIGHT = RGBColor(0xE8, 0xE7, 0xE3)

INFO_BG = RGBColor(0xE8, 0xF0, 0xFE)
INFO_EDGE = RGBColor(0x42, 0x85, 0xF4)
WARN_BG = RGBColor(0xFE, 0xF7, 0xE0)
WARN_EDGE = RGBColor(0xF9, 0xAB, 0x00)
WARN_TEXT = RGBColor(0x5F, 0x4B, 0x08)
GOOD_BG = RGBColor(0xE6, 0xF4, 0xEA)
GOOD_EDGE = RGBColor(0x34, 0xA8, 0x53)
GOOD_TEXT = RGBColor(0x0D, 0x65, 0x2D)

WINE_BG_DARK = RGBColor(0x1A, 0x1A, 0x1A)
WINE_BG_MID = RGBColor(0x2A, 0x06, 0x06)

WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_HEAD = "Calibri"  # PowerPoint-bundled; safe across PPT, Keynote, Impress
FONT_BODY = "Calibri"

# Slide dimensions — 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Standard padding
PAD_X = Inches(0.55)
PAD_Y = Inches(0.40)
CONTENT_W = SLIDE_W - 2 * PAD_X

DECK_TAG = "Week 5 · Supervisor Orientation"


# ---------------------------------------------------------------------------
# Shape helpers
# ---------------------------------------------------------------------------


def set_slide_background(slide, color: RGBColor) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill: RGBColor | None,
             line: RGBColor | None = None, line_width: float | None = None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        if line_width is not None:
            shp.line.width = Pt(line_width)
    return shp


def add_textbox(slide, left, top, width, height,
                text: str | None = None,
                font_size: float = 14,
                bold: bool = False,
                italic: bool = False,
                color: RGBColor = INK,
                font_name: str = FONT_BODY,
                align: int = PP_ALIGN.LEFT,
                anchor: int = MSO_ANCHOR.TOP,
                line_spacing: float | None = None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    if text is not None:
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb, tf


def style_run(run, *, font_size: float | None = None, bold: bool | None = None,
              italic: bool | None = None, color: RGBColor | None = None,
              font_name: str = FONT_BODY) -> None:
    run.font.name = font_name
    if font_size is not None:
        run.font.size = Pt(font_size)
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color


def add_paragraph_runs(tf, parts, *, font_size: float, color: RGBColor = INK,
                       bold: bool = False, italic: bool = False,
                       align: int = PP_ALIGN.LEFT,
                       line_spacing: float | None = None,
                       space_after: float | None = None,
                       bullet_marker: str | None = None,
                       new_paragraph: bool = True,
                       font_name: str = FONT_BODY) -> None:
    """Append a paragraph composed of styled runs.

    `parts` is an iterable of either strings or (text, style_overrides) tuples.
    """
    if new_paragraph and tf.paragraphs and (
        tf.paragraphs[0].runs or len(tf.paragraphs) > 1
        or (tf.paragraphs[0].text or "").strip()
    ):
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    if space_after is not None:
        p.space_after = Pt(space_after)
    if bullet_marker:
        marker_run = p.add_run()
        marker_run.text = bullet_marker
        style_run(marker_run, font_size=font_size, bold=True, color=SCARLET,
                  font_name=font_name)
    for part in parts:
        if isinstance(part, str):
            run = p.add_run()
            run.text = part
            style_run(run, font_size=font_size, bold=bold, italic=italic,
                      color=color, font_name=font_name)
        else:
            text, overrides = part
            run = p.add_run()
            run.text = text
            style_run(
                run,
                font_size=overrides.get("size", font_size),
                bold=overrides.get("bold", bold),
                italic=overrides.get("italic", italic),
                color=overrides.get("color", color),
                font_name=overrides.get("font", font_name),
            )


# ---------------------------------------------------------------------------
# Inline-styled text extraction from the source HTML
# ---------------------------------------------------------------------------


def collapse_text(s: str) -> str:
    return re.sub(r"\s+", " ", s).replace("\xa0", " ").strip()


def extract_runs(node: Tag | NavigableString,
                 inherit: dict | None = None) -> list[tuple[str, dict]]:
    """Walk a BeautifulSoup node and return a list of (text, style) runs.

    Style dict keys: bold, italic, color, size, font.
    """
    inherit = inherit or {}
    runs: list[tuple[str, dict]] = []

    def walk(n, style):
        if isinstance(n, NavigableString):
            text = str(n).replace("\xa0", " ")
            text = re.sub(r"[\t\r\n]+", " ", text)
            text = re.sub(r" +", " ", text)
            if text:
                runs.append((text, dict(style)))
            return
        if not isinstance(n, Tag):
            return
        local = dict(style)
        name = n.name.lower()
        if name in ("strong", "b"):
            local["bold"] = True
        if name in ("em", "i"):
            local["italic"] = True
        if name == "br":
            runs.append(("\n", dict(style)))
            return
        if name == "code":
            local["font"] = "Consolas"
        for child in n.children:
            walk(child, local)

    walk(node, inherit)
    # Merge adjacent runs that share styling
    merged: list[tuple[str, dict]] = []
    for text, style in runs:
        if merged and merged[-1][1] == style and text != "\n" and merged[-1][0] != "\n":
            merged[-1] = (merged[-1][0] + text, merged[-1][1])
        else:
            merged.append((text, style))
    # Trim leading whitespace on first run
    if merged:
        first_text, first_style = merged[0]
        merged[0] = (first_text.lstrip(), first_style)
    if merged:
        last_text, last_style = merged[-1]
        merged[-1] = (last_text.rstrip(), last_style)
    return [(t, s) for t, s in merged if t]


def runs_to_string(runs: list[tuple[str, dict]]) -> str:
    return "".join(t for t, _ in runs)


# ---------------------------------------------------------------------------
# Speaker notes extraction
# ---------------------------------------------------------------------------


TAG_LABELS = {
    "briefing": "BRIEFING — 30 MIN",
    "joint": "JOINT SESSION — EXTENDED",
    "transition": "TRANSITION",
}


def extract_notes_text(slide_section: Tag) -> str:
    template = slide_section.find("template", class_="speaker-notes")
    if not template:
        return ""
    inner = template.decode_contents()
    soup = BeautifulSoup(inner, "html.parser")
    chunks: list[str] = []
    for block in soup.find_all("div", class_="note-block"):
        classes = block.get("class", [])
        kind = next((c for c in classes if c in TAG_LABELS), None)
        if kind is None:
            # Could be 'transition' or unknown; handle 'transition' anyway
            for c in classes:
                if c == "transition":
                    kind = "transition"
                    break
        # Try to read explicit tag span text first
        tag_span = block.find("span", class_="tag")
        if tag_span:
            label = collapse_text(tag_span.get_text())
            label_upper = label.upper()
        elif kind:
            label_upper = TAG_LABELS.get(kind, kind.upper())
        else:
            label_upper = "NOTES"
        # Remove the tag span from the block before extracting body
        if tag_span:
            tag_span.extract()
        body_lines: list[str] = [f"[{label_upper}]"]
        # Walk top-level children for paragraphs / lists
        for child in block.children:
            if isinstance(child, NavigableString):
                txt = collapse_text(str(child))
                if txt:
                    body_lines.append(txt)
                continue
            if not isinstance(child, Tag):
                continue
            if child.name == "p":
                txt = collapse_text(child.get_text(" "))
                if txt:
                    body_lines.append(txt)
            elif child.name in ("ul", "ol"):
                for li in child.find_all("li", recursive=False):
                    txt = collapse_text(li.get_text(" "))
                    if txt:
                        body_lines.append(f"  • {txt}")
            else:
                txt = collapse_text(child.get_text(" "))
                if txt:
                    body_lines.append(txt)
        chunks.append("\n".join(body_lines))
    return "\n\n".join(chunks)


def set_speaker_notes(slide, text: str) -> None:
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.clear()
    if not text:
        return
    paragraphs = text.split("\n")
    p = notes_tf.paragraphs[0]
    first = True
    for line in paragraphs:
        if first:
            target = p
            first = False
        else:
            target = notes_tf.add_paragraph()
        run = target.add_run()
        run.text = line
        run.font.size = Pt(11)
        run.font.name = FONT_BODY


# ---------------------------------------------------------------------------
# Slide chrome — top brand bar + foot strip
# ---------------------------------------------------------------------------


def add_top_brand_bar(slide) -> None:
    """A red+gold brand bar at the very top (matches deck.css)."""
    bar_h = Inches(0.10)
    # Scarlet 70%, gold 30%
    width_scarlet = Inches(13.333 * 0.7)
    add_rect(slide, Emu(0), Emu(0), width_scarlet, bar_h, SCARLET)
    add_rect(slide, width_scarlet, Emu(0),
             SLIDE_W - width_scarlet, bar_h, GOLD)


def add_foot(slide, course_label: str, slide_num: int, total: int,
             on_dark: bool = False) -> None:
    color = RGBColor(0x9E, 0x9E, 0x9E) if on_dark else INK_3
    foot_top = SLIDE_H - Inches(0.45)
    # Scarlet tick
    tick_w = Inches(0.06)
    tick_h = Inches(0.18)
    add_rect(slide, PAD_X, foot_top + Inches(0.04), tick_w, tick_h, SCARLET)
    add_textbox(
        slide,
        PAD_X + tick_w + Inches(0.08),
        foot_top,
        Inches(8),
        Inches(0.32),
        text=course_label,
        font_size=10,
        color=color,
    )
    add_textbox(
        slide,
        SLIDE_W - PAD_X - Inches(2.0),
        foot_top,
        Inches(2.0),
        Inches(0.32),
        text=f"{slide_num} / {total}",
        font_size=10,
        color=color,
        align=PP_ALIGN.RIGHT,
    )


def _estimate_visible_lines(text: str, box_w_in: float, font_size_pt: float) -> int:
    """Conservative visible-line count for bold sans display text.
    Mirrors the helper in scripts/build_pptx_week_6.py so head layout (badge
    pill, rule, body_top) can be sized from the actual title line count."""
    if not text:
        return 1
    char_w = 0.62 * float(font_size_pt) / 72.0
    chars_per_in = max(6, int(float(box_w_in) / char_w))
    total = 0
    for chunk in str(text).split("\n"):
        n = max(1, len(chunk))
        total += (n + chars_per_in - 1) // chars_per_in
    return max(1, total)


def _title_plain_text(title) -> str:
    if isinstance(title, str):
        return title
    return "".join(text for text, _ in title)


def add_slide_head(slide, badge_text: str, badge_kind: str, title: str | list):
    """Render a content-slide head: pill badge + title row + thin rule."""
    # Badge
    pad = PAD_X
    head_top = PAD_Y + Inches(0.18)
    badge_h = Inches(0.34)
    badge_text = badge_text.upper()
    # Estimate badge width based on character count
    badge_w = Inches(max(1.4, 0.13 * len(badge_text) + 0.5))
    fill, fg, line = _badge_palette(badge_kind)
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, pad, head_top, badge_w, badge_h
    )
    badge.shadow.inherit = False
    badge.adjustments[0] = 0.5
    badge.fill.solid()
    badge.fill.fore_color.rgb = fill
    if line is not None:
        badge.line.color.rgb = line
        badge.line.width = Pt(0.75)
    else:
        badge.line.fill.background()
    btf = badge.text_frame
    btf.margin_left = Emu(0); btf.margin_right = Emu(0)
    btf.margin_top = Emu(0); btf.margin_bottom = Emu(0)
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    brun = bp.add_run()
    brun.text = badge_text
    brun.font.name = FONT_HEAD
    brun.font.size = Pt(10)
    brun.font.bold = True
    brun.font.color.rgb = fg

    # Title to the right of badge — size box and downstream rule/body_top
    # from actual visible-line count so a 2-line title doesn't punch through
    # the rule or the body region.
    title_left = pad + badge_w + Inches(0.30)
    title_w = SLIDE_W - title_left - PAD_X
    title_w_in = title_w / 914400
    title_plain = _title_plain_text(title)
    title_lines = _estimate_visible_lines(title_plain, title_w_in, 28)
    title_h_in = max(0.50, (28 / 72.0) * 1.10 * title_lines + 0.18)
    title_top = head_top - Inches(0.06)
    title_box, ttf = add_textbox(
        slide, title_left, title_top, title_w, Inches(title_h_in),
        font_size=28, bold=True, color=INK,
        font_name=FONT_HEAD,
        anchor=MSO_ANCHOR.MIDDLE if title_lines == 1 else MSO_ANCHOR.TOP,
    )
    if isinstance(title, str):
        run = ttf.paragraphs[0].add_run() if not ttf.paragraphs[0].runs else ttf.paragraphs[0].runs[0]
        run.text = title
        run.font.name = FONT_HEAD
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = INK
    else:
        # title is a list of (text, style) runs
        p = ttf.paragraphs[0]
        for text, style in title:
            r = p.add_run()
            r.text = text
            style_run(
                r,
                font_size=style.get("size", 28),
                bold=style.get("bold", True),
                italic=style.get("italic", False),
                color=style.get("color", INK),
                font_name=style.get("font", FONT_HEAD),
            )

    # Rule under head — position below whichever is taller (badge or title).
    head_bottom = max(head_top + badge_h, title_top + Inches(title_h_in))
    rule_top = head_bottom + Inches(0.12)
    add_rect(slide, pad, rule_top, CONTENT_W, Emu(12700), RULE_LIGHT)

    # Return the y-coordinate where body content can begin
    return rule_top + Inches(0.12)


def _badge_palette(kind: str) -> tuple[RGBColor, RGBColor, RGBColor | None]:
    if kind == "scarlet":
        return SCARLET, WHITE, None
    if kind == "gold":
        return GOLD, INK, None
    if kind == "ghost":
        return PAPER, INK_2, RULE_LIGHT
    if kind == "instr":
        return RGBColor(0xFC, 0xE8, 0xE8), SCARLET_DEEP, SCARLET
    # default
    return RULE_LIGHT, INK_2, None


# ---------------------------------------------------------------------------
# Content block helpers
# ---------------------------------------------------------------------------


def add_bullets(slide, left, top, width, height,
                items: list[list[tuple[str, dict]]],
                font_size: float = 14,
                color: RGBColor = INK,
                line_spacing: float = 1.25,
                space_after: float = 6.0) -> None:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    first = True
    for runs in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        # Marker
        m = p.add_run()
        m.text = "■  "
        style_run(m, font_size=font_size * 0.85, bold=True, color=SCARLET)
        for text, style in runs:
            r = p.add_run()
            r.text = text
            style_run(
                r,
                font_size=style.get("size", font_size),
                bold=style.get("bold", False),
                italic=style.get("italic", False),
                color=style.get("color", color),
                font_name=style.get("font", FONT_BODY),
            )


def add_paragraph_text(slide, left, top, width, height,
                       runs: list[tuple[str, dict]],
                       font_size: float = 14,
                       color: RGBColor = INK,
                       bold: bool = False,
                       italic: bool = False,
                       align: int = PP_ALIGN.LEFT,
                       line_spacing: float = 1.3) -> None:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for text, style in runs:
        # Handle newline marks
        if text == "\n":
            p = tf.add_paragraph()
            p.alignment = align
            p.line_spacing = line_spacing
            continue
        r = p.add_run()
        r.text = text
        style_run(
            r,
            font_size=style.get("size", font_size),
            bold=style.get("bold", bold),
            italic=style.get("italic", italic),
            color=style.get("color", color),
            font_name=style.get("font", FONT_BODY),
        )


def add_card(slide, left, top, width, height, *,
             fill: RGBColor = PAPER_WHITE,
             top_accent: RGBColor | None = None,
             left_accent: RGBColor | None = None,
             outline: RGBColor | None = RULE_LIGHT) -> None:
    add_rect(slide, left, top, width, height, fill,
             line=outline, line_width=0.75)
    if top_accent is not None:
        add_rect(slide, left, top, width, Inches(0.07), top_accent)
    if left_accent is not None:
        add_rect(slide, left, top, Inches(0.07), height, left_accent)


# ---------------------------------------------------------------------------
# Layout builders
# ---------------------------------------------------------------------------


def build_cover(slide, *, eyebrow: str, title_runs, subtitle: str,
                badges: list[tuple[str, str]], corner: str) -> None:
    set_slide_background(slide, WINE_BG_DARK)
    # Wine gradient feel — overlay a deep wine block on the right
    add_rect(slide, Inches(6), Emu(0), SLIDE_W - Inches(6), SLIDE_H,
             WINE_BG_MID)
    add_rect(slide, Inches(9), Emu(0), SLIDE_W - Inches(9), SLIDE_H,
             SCARLET_DEEP)

    # Corner badge
    tb, tf = add_textbox(slide, SLIDE_W - PAD_X - Inches(5.5),
                         Inches(0.35), Inches(5.5), Inches(0.35),
                         font_size=10, color=GOLD,
                         align=PP_ALIGN.RIGHT)
    p = tf.paragraphs[0]
    parts = corner.split("·")
    if len(parts) == 2:
        r1 = p.add_run()
        r1.text = parts[0].strip()
        style_run(r1, font_size=10, bold=True, color=GOLD)
        r2 = p.add_run()
        r2.text = "  ·  " + parts[1].strip()
        style_run(r2, font_size=10, color=RGBColor(0xCC, 0xCC, 0xCC))
    else:
        r = p.add_run(); r.text = corner
        style_run(r, font_size=10, color=GOLD)
    # uppercase via font
    for run in p.runs:
        run.font.name = FONT_HEAD

    # Eyebrow
    add_textbox(slide, PAD_X, Inches(1.6), Inches(11), Inches(0.35),
                text=eyebrow.upper(),
                font_size=12, bold=True, color=GOLD)

    # Title
    tb, tf = add_textbox(slide, PAD_X, Inches(2.0),
                         Inches(11.5), Inches(3.0),
                         font_size=44, bold=True, color=WHITE,
                         font_name=FONT_HEAD,
                         line_spacing=1.05)
    p = tf.paragraphs[0]
    for text, style in title_runs:
        r = p.add_run()
        r.text = text
        style_run(r, font_size=style.get("size", 44),
                  bold=style.get("bold", True),
                  italic=style.get("italic", False),
                  color=style.get("color", WHITE),
                  font_name=FONT_HEAD)

    # Subtitle
    add_textbox(slide, PAD_X, Inches(5.0),
                Inches(10.5), Inches(1.2),
                text=subtitle,
                font_size=18, color=RGBColor(0xDD, 0xDD, 0xDD),
                line_spacing=1.35)

    # Badge row
    bx = PAD_X
    by = Inches(6.3)
    for label, kind in badges:
        bw = Inches(0.18 * len(label) + 0.6)
        bh = Inches(0.36)
        if kind == "gold":
            fill = GOLD; fg = INK; line = None
        else:
            fill = RGBColor(0x33, 0x16, 0x16); fg = WHITE; line = RGBColor(0x66, 0x33, 0x33)
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, bw, bh)
        shp.shadow.inherit = False
        shp.adjustments[0] = 0.5
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
        if line is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = line
            shp.line.width = Pt(0.75)
        tf = shp.text_frame
        tf.margin_left = Emu(0); tf.margin_right = Emu(0)
        tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = label.upper()
        style_run(run, font_size=10, bold=True, color=fg)
        bx += bw + Inches(0.12)

    # Gold + scarlet bottom bars
    add_rect(slide, Emu(0), SLIDE_H - Inches(0.18),
             SLIDE_W, Inches(0.05), GOLD)
    add_rect(slide, Emu(0), SLIDE_H - Inches(0.13),
             SLIDE_W, Inches(0.13), SCARLET)


def build_divider(slide, *, module_label: str, title: str,
                  subtitle: str) -> None:
    set_slide_background(slide, WINE_BG_DARK)
    add_rect(slide, Inches(8), Emu(0), SLIDE_W - Inches(8), SLIDE_H,
             WINE_BG_MID)
    # Module label
    add_textbox(slide, PAD_X, Inches(2.4), Inches(11), Inches(0.4),
                text=module_label.upper(),
                font_size=14, bold=True, color=GOLD)
    # Title
    add_textbox(slide, PAD_X, Inches(2.95), Inches(11.5), Inches(2.2),
                text=title,
                font_size=52, bold=True, color=WHITE,
                font_name=FONT_HEAD, line_spacing=1.05)
    # Subtitle
    add_textbox(slide, PAD_X, Inches(5.5), Inches(11), Inches(1.0),
                text=subtitle,
                font_size=18, color=RGBColor(0xCC, 0xCC, 0xCC),
                line_spacing=1.35)
    # Bottom scarlet bar
    add_rect(slide, Emu(0), SLIDE_H - Inches(0.13),
             SLIDE_W, Inches(0.13), SCARLET)


def build_closing(slide, *, eyebrow: str, title: str,
                  ask_label: str, ask_text_runs) -> None:
    set_slide_background(slide, WINE_BG_DARK)
    add_rect(slide, Inches(8.5), Emu(0), SLIDE_W - Inches(8.5), SLIDE_H,
             WINE_BG_MID)
    add_textbox(slide, PAD_X, Inches(1.3), Inches(11), Inches(0.4),
                text=eyebrow.upper(),
                font_size=13, bold=True, color=GOLD)
    add_textbox(slide, PAD_X, Inches(1.85), Inches(12.0), Inches(3.5),
                text=title,
                font_size=40, bold=True, color=WHITE,
                font_name=FONT_HEAD, line_spacing=1.1)
    # Ask block — gold left bar + tinted gold panel
    ask_top = Inches(5.3)
    ask_h = Inches(1.5)
    add_rect(slide, PAD_X, ask_top, Inches(11), ask_h,
             RGBColor(0x3E, 0x33, 0x10), line=None)
    add_rect(slide, PAD_X, ask_top, Inches(0.08), ask_h, GOLD)
    add_textbox(slide, PAD_X + Inches(0.30), ask_top + Inches(0.18),
                Inches(10.5), Inches(0.35),
                text=ask_label.upper(),
                font_size=10, bold=True, color=GOLD)
    add_paragraph_text(slide,
                       PAD_X + Inches(0.30),
                       ask_top + Inches(0.55),
                       Inches(10.5), Inches(1.2),
                       ask_text_runs,
                       font_size=20, color=WHITE,
                       line_spacing=1.3)
    add_rect(slide, Emu(0), SLIDE_H - Inches(0.13),
             SLIDE_W, Inches(0.13), SCARLET)


# ---------------------------------------------------------------------------
# Main build pipeline — slide-by-slide
# ---------------------------------------------------------------------------


@dataclass
class SlideSpec:
    """Holds the parsed source <section> plus its index."""
    section: Tag
    index: int  # 1-based
    total: int
    course_label: str  # foot label from source


def parse_source() -> list[SlideSpec]:
    html = SOURCE_HTML.read_text(encoding="utf-8")
    soup = BeautifulSoup(html, "html.parser")
    deck = soup.find(id="deck")
    sections = deck.find_all("section", class_="slide", recursive=False)
    total = len(sections)
    specs: list[SlideSpec] = []
    for i, sec in enumerate(sections, start=1):
        foot = sec.find("div", class_="slide__foot")
        course_label = DECK_TAG
        if foot:
            label_span = foot.find("span")
            if label_span:
                # Strip the inner .scarlet-mark span
                mark = label_span.find("span", class_="scarlet-mark")
                if mark:
                    mark.extract()
                course_label = collapse_text(label_span.get_text(" "))
        specs.append(SlideSpec(section=sec, index=i, total=total,
                               course_label=course_label or DECK_TAG))
    return specs


# -- Per-slide builders ---------------------------------------------------


def render_slide_1_cover(slide, spec: SlideSpec) -> None:
    sec = spec.section
    eyebrow = collapse_text(sec.find("div", class_="eyebrow").get_text())
    title_runs = extract_runs(sec.find("h1"))
    subtitle = collapse_text(sec.find("p", class_="deck-sub").get_text())
    badges: list[tuple[str, str]] = []
    for b in sec.find("div", class_="cover-meta").find_all("span", class_="badge"):
        classes = b.get("class", [])
        kind = "gold" if "badge--gold" in classes else "default"
        badges.append((collapse_text(b.get_text()), kind))
    corner = collapse_text(sec.find("div", class_="corner").get_text())
    build_cover(slide,
                eyebrow=eyebrow,
                title_runs=title_runs,
                subtitle=subtitle,
                badges=badges,
                corner=corner)


def render_content_head(slide, sec: Tag):
    head = sec.find("div", class_="slide__head")
    badge = head.find("span", class_="badge")
    badge_classes = badge.get("class", [])
    if "badge--scarlet" in badge_classes:
        kind = "scarlet"
    elif "badge--gold" in badge_classes:
        kind = "gold"
    elif "badge--ghost" in badge_classes:
        kind = "ghost"
    elif "badge--instr" in badge_classes:
        kind = "instr"
    else:
        kind = "default"
    badge_text = collapse_text(badge.get_text())
    title_runs = extract_runs(head.find(["h2", "h3"]))
    body_top = add_slide_head(slide, badge_text, kind, title_runs)
    return body_top


# ---- Slide 2: Audience shift -------------------------------------------


def render_slide_2_audience(slide, spec: SlideSpec) -> None:
    sec = spec.section
    body_top = render_content_head(slide, sec)
    # Two panels
    panels = sec.find_all("div", class_="panel")
    panel_w = (CONTENT_W - Inches(0.30)) / 2
    panel_h = SLIDE_H - body_top - Inches(0.55)
    for i, panel in enumerate(panels):
        left = PAD_X + i * (panel_w + Inches(0.30))
        accent = SCARLET if "this-week" in panel.get("class", []) else INK_3
        add_card(slide, left, body_top, panel_w, panel_h,
                 fill=PAPER_WHITE, top_accent=accent)
        who = panel.find("span", class_="who")
        if who:
            add_textbox(slide, left + Inches(0.30),
                        body_top + Inches(0.30),
                        panel_w - Inches(0.60), Inches(0.3),
                        text=collapse_text(who.get_text()).upper(),
                        font_size=10, bold=True, color=INK_3)
        h3 = panel.find("h3")
        title_y = body_top + Inches(0.65)
        add_paragraph_text(
            slide, left + Inches(0.30), title_y,
            panel_w - Inches(0.60), Inches(1.2),
            extract_runs(h3),
            font_size=20, bold=True, color=INK,
            line_spacing=1.15,
        )
        # Bullets
        items = [extract_runs(li) for li in panel.find_all("li")]
        list_y = title_y + Inches(1.4)
        add_bullets(slide, left + Inches(0.30), list_y,
                    panel_w - Inches(0.60),
                    panel_h - Inches(2.1),
                    items, font_size=14, color=INK_2,
                    line_spacing=1.25, space_after=8.0)


# ---- Slide 3: Two delivery modes ---------------------------------------


def render_slide_3_modes(slide, spec: SlideSpec) -> None:
    sec = spec.section
    body_top = render_content_head(slide, sec)
    modes = sec.find_all("div", class_="mode")
    panel_w = (CONTENT_W - Inches(0.30)) / 2
    panel_h = SLIDE_H - body_top - Inches(0.55)
    for i, mode in enumerate(modes):
        left = PAD_X + i * (panel_w + Inches(0.30))
        accent = SCARLET if "long" in mode.get("class", []) else GOLD
        add_card(slide, left, body_top, panel_w, panel_h,
                 fill=PAPER_WHITE, top_accent=accent)
        time = mode.find("span", class_="mode__time")
        h3 = mode.find("h3")
        ps = mode.find_all("p", recursive=False)
        ul = mode.find("ul")
        cur_y = body_top + Inches(0.30)
        if time:
            add_textbox(slide, left + Inches(0.30), cur_y,
                        panel_w - Inches(0.60), Inches(0.3),
                        text=collapse_text(time.get_text()).upper(),
                        font_size=10, bold=True, color=INK_3)
            cur_y += Inches(0.35)
        if h3:
            add_paragraph_text(slide, left + Inches(0.30), cur_y,
                               panel_w - Inches(0.60), Inches(0.6),
                               extract_runs(h3),
                               font_size=22, bold=True, color=INK,
                               line_spacing=1.15)
            cur_y += Inches(0.65)
        for p_tag in ps:
            add_paragraph_text(slide, left + Inches(0.30), cur_y,
                               panel_w - Inches(0.60), Inches(1.0),
                               extract_runs(p_tag),
                               font_size=13, color=INK_2,
                               line_spacing=1.35)
            cur_y += Inches(0.85)
        if ul:
            items = [extract_runs(li) for li in ul.find_all("li")]
            add_bullets(slide, left + Inches(0.30), cur_y,
                        panel_w - Inches(0.60),
                        body_top + panel_h - cur_y - Inches(0.20),
                        items, font_size=12.5, color=INK_2,
                        line_spacing=1.25, space_after=6.0)


# ---- Slide 4: Agenda table ---------------------------------------------


def render_slide_4_agenda(slide, spec: SlideSpec) -> None:
    sec = spec.section
    body_top = render_content_head(slide, sec)
    table = sec.find("table", class_="deck-tbl")
    headers = [th.get_text(" ", strip=True)
               for th in table.find("thead").find_all("th")]
    rows = []
    for tr in table.find("tbody").find_all("tr"):
        cells = []
        for td in tr.find_all("td"):
            cells.append(extract_runs(td))
        rows.append(cells)
    cols = len(headers)
    rows_n = len(rows) + 1
    table_w = CONTENT_W
    table_h = Inches(4.2)
    table_top = body_top + Inches(0.15)
    shape = slide.shapes.add_table(rows_n, cols, PAD_X, table_top, table_w, table_h)
    tbl = shape.table
    # Column widths
    col_w = [Inches(2.0), Inches(3.6), CONTENT_W - Inches(2.0) - Inches(3.6)]
    for i, w in enumerate(col_w):
        tbl.columns[i].width = w
    # Header row
    tbl.rows[0].height = Inches(0.55)
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = INK
        cell.margin_left = Inches(0.18)
        cell.margin_right = Inches(0.18)
        cell.margin_top = Inches(0.10)
        cell.margin_bottom = Inches(0.10)
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = h
        style_run(r, font_size=12, bold=True, color=WHITE)
    # Body rows
    for ri, row in enumerate(rows, start=1):
        tbl.rows[ri].height = Inches(0.65)
        for ci, runs in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                RGBColor(0xFA, 0xFA, 0xF8) if ri % 2 == 0 else PAPER_WHITE
            )
            cell.margin_left = Inches(0.18)
            cell.margin_right = Inches(0.18)
            cell.margin_top = Inches(0.10)
            cell.margin_bottom = Inches(0.10)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            for text, style in runs:
                r = p.add_run(); r.text = text
                style_run(r,
                          font_size=style.get("size", 12),
                          bold=style.get("bold", False),
                          italic=style.get("italic", False),
                          color=style.get("color", INK))


# ---- Slide 5: Centerpiece ----------------------------------------------


def render_slide_5_core(slide, spec: SlideSpec) -> None:
    sec = spec.section
    body_top = render_content_head(slide, sec)
    cp = sec.find("div", class_="centerpiece")
    pull = cp.find("div", class_="pull")
    attribution = cp.find("div", class_="attribution")
    box_top = body_top + Inches(0.20)
    box_h = SLIDE_H - box_top - Inches(0.55)
    add_rect(slide, PAD_X, box_top, CONTENT_W, box_h, PAPER_WHITE,
             line=RULE_LIGHT, line_width=0.75)
    add_rect(slide, PAD_X, box_top, CONTENT_W, Inches(0.10), SCARLET)
    # Pull text
    runs = extract_runs(pull)
    add_paragraph_text(slide, PAD_X + Inches(0.5),
                       box_top + Inches(0.9),
                       CONTENT_W - Inches(1.0),
                       box_h - Inches(2.0),
                       runs,
                       font_size=44, bold=True, color=INK,
                       align=PP_ALIGN.CENTER, line_spacing=1.1)
    if attribution:
        add_textbox(slide, PAD_X, box_top + box_h - Inches(0.7),
                    CONTENT_W, Inches(0.5),
                    text=collapse_text(attribution.get_text()).upper(),
                    font_size=11, color=INK_3, align=PP_ALIGN.CENTER)


# ---- Slide 6, 9, 13, 19, 23 dividers -----------------------------------


def render_divider(slide, spec: SlideSpec) -> None:
    sec = spec.section
    module_num = collapse_text(sec.find("div", class_="module-num").get_text())
    h1 = collapse_text(sec.find("h1").get_text())
    meta = collapse_text(sec.find("div", class_="module-meta").get_text())
    build_divider(slide, module_label=module_num, title=h1, subtitle=meta)


# ---- Slide 7: Data wall (5 stat cards) ---------------------------------


def render_slide_7_data(slide, spec: SlideSpec) -> None:
    sec = spec.section
    body_top = render_content_head(slide, sec)
    stats = sec.find("div", class_="stat-row").find_all("div", class_="stat")
    n = len(stats)
    gap = Inches(0.18)
    card_w = (CONTENT_W - gap * (n - 1)) / n
    card_h = Inches(2.4)
    for i, s in enumerate(stats):
        left = PAD_X + i * (card_w + gap)
        add_card(slide, left, body_top, card_w, card_h,
                 fill=PAPER_WHITE, top_accent=SCARLET)
        num = s.find("div", class_="num")
        lbl = s.find("div", class_="lbl")
        add_textbox(slide, left + Inches(0.10),
                    body_top + Inches(0.30),
                    card_w - Inches(0.20), Inches(0.9),
                    text=collapse_text(num.get_text()),
                    font_size=30, bold=True, color=SCARLET,
                    align=PP_ALIGN.CENTER, line_spacing=1.0)
        # Label may include <em> emphasis and a <span class="micro"> source
        micro = lbl.find("span", class_="micro")
        micro_text = collapse_text(micro.get_text()) if micro else ""
        if micro:
            micro.extract()
        # Convert <br> to newline
        for br in lbl.find_all("br"):
            br.replace_with("\n")
        lbl_runs = extract_runs(lbl)
        add_paragraph_text(slide, left + Inches(0.10),
                           body_top + Inches(1.20),
                           card_w - Inches(0.20), Inches(0.9),
                           lbl_runs,
                           font_size=10.5, color=INK_2,
                           align=PP_ALIGN.CENTER, line_spacing=1.25)
        if micro_text:
            add_textbox(slide, left + Inches(0.10),
                        body_top + card_h - Inches(0.45),
                        card_w - Inches(0.20), Inches(0.3),
                        text=micro_text.upper(),
                        font_size=8.5, bold=True, color=INK_3,
                        align=PP_ALIGN.CENTER)
    # Lead and source
    lead = sec.find("p", class_="lead")
    if lead:
        runs = extract_runs(lead)
        add_paragraph_text(slide, PAD_X,
                           body_top + card_h + Inches(0.25),
                           CONTENT_W, Inches(0.7),
                           runs,
                           font_size=14, color=INK,
                           line_spacing=1.35)
    source = sec.find("p", class_="source")
    if source:
        add_textbox(slide, PAD_X,
                    SLIDE_H - Inches(0.85),
                    CONTENT_W, Inches(0.3),
                    text=collapse_text(source.get_text()),
                    font_size=9, italic=True, color=INK_3)


# ---- Slide 8: AI is a SITREP item (centerpiece + lead) -----------------


def render_slide_8_sitrep(slide, spec: SlideSpec) -> None:
    sec = spec.section
    body_top = render_content_head(slide, sec)
    cp = sec.find("div", class_="centerpiece")
    pull = cp.find("div", class_="pull")
    lead = cp.find("p", class_="lead")
    box_top = body_top + Inches(0.20)
    box_h = SLIDE_H - box_top - Inches(0.55)
    add_rect(slide, PAD_X, box_top, CONTENT_W, box_h, PAPER_WHITE,
             line=RULE_LIGHT, line_width=0.75)
    add_rect(slide, PAD_X, box_top, CONTENT_W, Inches(0.10), SCARLET)
    add_paragraph_text(slide, PAD_X + Inches(0.5),
                       box_top + Inches(0.7),
                       CONTENT_W - Inches(1.0),
                       Inches(1.5),
                       extract_runs(pull),
                       font_size=46, bold=True, color=INK,
                       align=PP_ALIGN.CENTER, line_spacing=1.1)
    if lead:
        add_paragraph_text(slide, PAD_X + Inches(1.0),
                           box_top + Inches(2.5),
                           CONTENT_W - Inches(2.0),
                           Inches(2.5),
                           extract_runs(lead),
                           font_size=18, color=INK_2,
                           align=PP_ALIGN.CENTER, line_spacing=1.4)


# ---- Slide 10: What yes / what kills (two cols) ------------------------


def render_two_cols_with_heads(slide, spec: SlideSpec) -> None:
    sec = spec.section
    body_top = render_content_head(slide, sec)
    cols = sec.find("div", class_="cols").find_all("div", class_="col",
                                                   recursive=False)
    panel_w = (CONTENT_W - Inches(0.40)) / 2
    panel_h = SLIDE_H - body_top - Inches(0.55)
    for i, col in enumerate(cols):
        left = PAD_X + i * (panel_w + Inches(0.40))
        head = col.find("div", class_="col__head")
        head_color = INK_2
        rule_color = RULE_LIGHT
        if head and "col__head--good" in head.get("class", []):
            head_color = GOOD_TEXT; rule_color = GOOD_EDGE
        elif head and "col__head--bad" in head.get("class", []):
            head_color = SCARLET_DEEP; rule_color = SCARLET
        cur_y = body_top
        if head:
            add_textbox(slide, left, cur_y, panel_w, Inches(0.4),
                        text=collapse_text(head.get_text()).upper(),
                        font_size=12, bold=True, color=head_color)
            cur_y += Inches(0.40)
            add_rect(slide, left, cur_y, panel_w, Emu(19050), rule_color)
            cur_y += Inches(0.20)
        ul = col.find("ul")
        if ul:
            items = [extract_runs(li) for li in ul.find_all("li")]
            add_bullets(slide, left, cur_y, panel_w,
                        panel_h - (cur_y - body_top),
                        items, font_size=14, color=INK,
                        line_spacing=1.3, space_after=8.0)


# ---- Slide 11: What NOT to do (numbered big-list inside warn callout) --


def render_big_list_in_card(slide, spec: SlideSpec, *,
                            card_fill: RGBColor = WARN_BG,
                            edge: RGBColor = WARN_EDGE,
                            text_color: RGBColor = INK) -> None:
    sec = spec.section
    body_top = render_content_head(slide, sec)
    card_top = body_top + Inches(0.20)
    card_h = SLIDE_H - card_top - Inches(0.55)
    add_rect(slide, PAD_X, card_top, CONTENT_W, card_h, card_fill,
             line=None)
    add_rect(slide, PAD_X, card_top, Inches(0.10), card_h, edge)
    items = sec.find_all("li")
    n = len(items)
    inner_left = PAD_X + Inches(0.40)
    inner_top = card_top + Inches(0.30)
    inner_w = CONTENT_W - Inches(0.70)
    row_h = (card_h - Inches(0.60)) / n
    for i, li in enumerate(items):
        row_y = inner_top + i * row_h
        # Numbered circle
        circle_d = Inches(0.55)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                        inner_left,
                                        row_y + Inches(0.05),
                                        circle_d, circle_d)
        circle.shadow.inherit = False
        circle.fill.solid(); circle.fill.fore_color.rgb = SCARLET
        circle.line.fill.background()
        ctf = circle.text_frame
        ctf.margin_left = Emu(0); ctf.margin_right = Emu(0)
        ctf.margin_top = Emu(0); ctf.margin_bottom = Emu(0)
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run(); cr.text = str(i + 1)
        style_run(cr, font_size=18, bold=True, color=WHITE)
        # Text — strong + paragraph
        text_left = inner_left + circle_d + Inches(0.30)
        strong = li.find("strong")
        para = li.find("p")
        if strong:
            add_paragraph_text(slide, text_left, row_y,
                               inner_w - circle_d - Inches(0.30),
                               Inches(0.5),
                               extract_runs(strong),
                               font_size=15, bold=True, color=text_color,
                               line_spacing=1.2)
        if para:
            add_paragraph_text(slide, text_left, row_y + Inches(0.42),
                               inner_w - circle_d - Inches(0.30),
                               row_h - Inches(0.45),
                               extract_runs(para),
                               font_size=12.5, color=text_color,
                               line_spacing=1.35)


# ---- Slide 12: Guard rails not roadblocks (two cols + callout) ---------


def render_slide_12_guardrails(slide, spec: SlideSpec) -> None:
    sec = spec.section
    body_top = render_content_head(slide, sec)
    cols = sec.find("div", class_="cols").find_all("div", class_="col",
                                                   recursive=False)
    panel_w = (CONTENT_W - Inches(0.40)) / 2
    panel_h = SLIDE_H - body_top - Inches(0.55)
    # Left column: heading + bullets
    left_col = cols[0]
    h3 = left_col.find("h3")
    if h3:
        add_paragraph_text(slide, PAD_X, body_top,
                           panel_w, Inches(0.7),
                           extract_runs(h3),
                           font_size=20, bold=True, color=INK,
                           line_spacing=1.2)
    ul = left_col.find("ul")
    if ul:
        items = [extract_runs(li) for li in ul.find_all("li")]
        add_bullets(slide, PAD_X, body_top + Inches(0.85),
                    panel_w, panel_h - Inches(0.85),
                    items, font_size=14, color=INK,
                    line_spacing=1.3, space_after=8.0)
    # Right column: green "Default posture" callout
    right_col = cols[1]
    callout = right_col.find("div", class_="callout")
    cl_left = PAD_X + panel_w + Inches(0.40)
    cl_top = body_top + Inches(0.7)
    cl_h = panel_h - Inches(1.4)
    add_rect(slide, cl_left, cl_top, panel_w, cl_h, GOOD_BG, line=None)
    add_rect(slide, cl_left, cl_top, Inches(0.10), cl_h, GOOD_EDGE)
    h3c = callout.find("h3")
    paras = callout.find_all("p")
    add_textbox(slide, cl_left + Inches(0.40), cl_top + Inches(0.30),
                panel_w - Inches(0.60), Inches(0.5),
                text=collapse_text(h3c.get_text()),
                font_size=18, bold=True, color=GOOD_TEXT)
    if paras:
        add_paragraph_text(slide, cl_left + Inches(0.40),
                           cl_top + Inches(0.95),
                           panel_w - Inches(0.60), Inches(1.4),
                           extract_runs(paras[0]),
                           font_size=20, bold=True, color=GOOD_TEXT,
                           line_spacing=1.3)
    if len(paras) > 1:
        add_paragraph_text(slide, cl_left + Inches(0.40),
                           cl_top + Inches(2.2),
                           panel_w - Inches(0.60), Inches(2.0),
                           extract_runs(paras[1]),
                           font_size=13, color=GOOD_TEXT,
                           line_spacing=1.35)


# ---- Slide 14: Four questions (big-list 4 on white) --------------------


def render_big_list_plain(slide, spec: SlideSpec, *,
                          tail_runs: list[tuple[str, dict]] | None = None,
                          tail_size: float = 13.5,
                          ) -> None:
    sec = spec.section
    body_top = render_content_head(slide, sec)
    items = sec.find_all("li")
    n = len(items)
    list_top = body_top + Inches(0.15)
    list_bottom = SLIDE_H - Inches(0.85)
    if tail_runs:
        list_bottom -= Inches(0.85)
    list_h = list_bottom - list_top
    row_h = list_h / n
    for i, li in enumerate(items):
        row_y = list_top + i * row_h
        circle_d = Inches(0.55)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                        PAD_X,
                                        row_y + Inches(0.05),
                                        circle_d, circle_d)
        circle.shadow.inherit = False
        circle.fill.solid(); circle.fill.fore_color.rgb = SCARLET
        circle.line.fill.background()
        ctf = circle.text_frame
        ctf.margin_left = Emu(0); ctf.margin_right = Emu(0)
        ctf.margin_top = Emu(0); ctf.margin_bottom = Emu(0)
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run(); cr.text = str(i + 1)
        style_run(cr, font_size=18, bold=True, color=WHITE)
        text_left = PAD_X + circle_d + Inches(0.30)
        text_w = CONTENT_W - circle_d - Inches(0.30)
        strong = li.find("strong")
        para = li.find("p")
        if strong:
            add_paragraph_text(slide, text_left, row_y,
                               text_w, Inches(0.5),
                               extract_runs(strong),
                               font_size=16, bold=True, color=INK,
                               line_spacing=1.2)
        if para:
            add_paragraph_text(slide, text_left, row_y + Inches(0.45),
                               text_w, row_h - Inches(0.5),
                               extract_runs(para),
                               font_size=13, color=INK_2,
                               line_spacing=1.35)
    if tail_runs:
        add_paragraph_text(slide, PAD_X, list_bottom + Inches(0.15),
                           CONTENT_W, Inches(0.7),
                           tail_runs,
                           font_size=tail_size, color=INK,
                           line_spacing=1.35)


# ---- Slides 15, 17, 22: Scenario card + prompt -------------------------


def render_scenario_slide(slide, spec: SlideSpec) -> None:
    sec = spec.section
    body_top = render_content_head(slide, sec)
    scenario = sec.find("div", class_="scenario")
    label = scenario.find("div", class_="scenario__label")
    paragraph = scenario.find("p")
    sc_top = body_top + Inches(0.20)
    sc_h = Inches(2.4)
    add_rect(slide, PAD_X, sc_top, CONTENT_W, sc_h, PAPER_WHITE,
             line=RULE_LIGHT, line_width=0.75)
    add_rect(slide, PAD_X, sc_top, Inches(0.10), sc_h, GOLD)
    add_textbox(slide, PAD_X + Inches(0.30),
                sc_top + Inches(0.25),
                CONTENT_W - Inches(0.60), Inches(0.35),
                text=collapse_text(label.get_text()).upper(),
                font_size=11, bold=True, color=GOLD_DARK)
    add_paragraph_text(slide, PAD_X + Inches(0.30),
                       sc_top + Inches(0.70),
                       CONTENT_W - Inches(0.60), Inches(1.6),
                       extract_runs(paragraph),
                       font_size=15, color=INK, line_spacing=1.4)
    # Prompt box (dark)
    prompt = sec.find("div", class_="prompt-box")
    pr_top = sc_top + sc_h + Inches(0.30)
    pr_h = Inches(1.4)
    add_rect(slide, PAD_X, pr_top, CONTENT_W, pr_h,
             RGBColor(0x1A, 0x1A, 0x1A), line=None)
    label_span = prompt.find("span", class_="label")
    label_text = collapse_text(label_span.get_text())
    label_span.extract()
    body_text = collapse_text(prompt.get_text(" "))
    add_textbox(slide, PAD_X + Inches(0.30),
                pr_top + Inches(0.25),
                CONTENT_W - Inches(0.60), Inches(0.35),
                text=label_text.upper(),
                font_size=10, bold=True, color=GOLD)
    add_textbox(slide, PAD_X + Inches(0.30),
                pr_top + Inches(0.65),
                CONTENT_W - Inches(0.60), Inches(0.7),
                text=body_text,
                font_size=14, color=WHITE, line_spacing=1.4)


# ---- Slides 16, 18: Debrief two-cols (bullets + callout) ---------------


def render_debrief_slide(slide, spec: SlideSpec) -> None:
    sec = spec.section
    body_top = render_content_head(slide, sec)
    cols = sec.find("div", class_="cols").find_all("div", class_="col",
                                                   recursive=False)
    panel_w = (CONTENT_W - Inches(0.40)) / 2
    panel_h = SLIDE_H - body_top - Inches(0.65)
    # Left: bullets
    left_col = cols[0]
    ul = left_col.find("ul")
    items = [extract_runs(li) for li in ul.find_all("li")]
    add_bullets(slide, PAD_X, body_top + Inches(0.10),
                panel_w, panel_h, items,
                font_size=13.5, color=INK,
                line_spacing=1.3, space_after=10.0)
    # Right: callout
    right_col = cols[1]
    callout = right_col.find("div", class_="callout")
    classes = callout.get("class", [])
    if "callout--info" in classes:
        bg, edge, fg = INFO_BG, INFO_EDGE, INK
    elif "callout--scarlet" in classes:
        bg, edge, fg = PAPER_WHITE, SCARLET, INK
    elif "callout--good" in classes:
        bg, edge, fg = GOOD_BG, GOOD_EDGE, GOOD_TEXT
    elif "callout--warn" in classes:
        bg, edge, fg = WARN_BG, WARN_EDGE, WARN_TEXT
    else:
        bg, edge, fg = PAPER_WHITE, RULE_LIGHT, INK
    cl_left = PAD_X + panel_w + Inches(0.40)
    cl_top = body_top + Inches(0.6)
    cl_h = panel_h - Inches(1.0)
    add_rect(slide, cl_left, cl_top, panel_w, cl_h, bg,
             line=RULE_LIGHT if bg == PAPER_WHITE else None,
             line_width=0.5 if bg == PAPER_WHITE else None)
    add_rect(slide, cl_left, cl_top, Inches(0.10), cl_h, edge)
    h4 = callout.find("h4")
    p = callout.find("p")
    if h4:
        add_textbox(slide, cl_left + Inches(0.40),
                    cl_top + Inches(0.30),
                    panel_w - Inches(0.60), Inches(0.5),
                    text=collapse_text(h4.get_text()),
                    font_size=18, bold=True, color=fg)
    if p:
        add_paragraph_text(slide, cl_left + Inches(0.40),
                           cl_top + Inches(0.95),
                           panel_w - Inches(0.60), cl_h - Inches(1.2),
                           extract_runs(p),
                           font_size=13.5, color=fg,
                           line_spacing=1.45)


# ---- Slide 20: Apprentice problem (quote + 3 stats) --------------------


def render_slide_20_apprentice(slide, spec: SlideSpec) -> None:
    sec = spec.section
    body_top = render_content_head(slide, sec)
    app = sec.find("div", class_="apprentice")
    quote = app.find("div", class_="quote")
    stats = app.find("div", class_="stats").find_all("div", class_="stat")
    panel_h = SLIDE_H - body_top - Inches(0.65)
    quote_w = (CONTENT_W - Inches(0.40)) * 0.55
    # Quote panel — dark background
    add_rect(slide, PAD_X, body_top + Inches(0.10),
             quote_w, panel_h - Inches(0.10),
             RGBColor(0x1A, 0x1A, 0x1A), line=None)
    # Big scarlet quote glyph
    add_textbox(slide, PAD_X + Inches(0.30),
                body_top - Inches(0.3),
                Inches(2.0), Inches(2.0),
                text="\u201C",
                font_size=110, bold=True, color=SCARLET,
                font_name="Georgia")
    q = quote.find("p", class_="q")
    by = quote.find("div", class_="by")
    add_paragraph_text(slide,
                       PAD_X + Inches(0.45),
                       body_top + Inches(0.85),
                       quote_w - Inches(0.90),
                       panel_h - Inches(2.0),
                       extract_runs(q),
                       font_size=24, bold=True, color=WHITE,
                       line_spacing=1.2)
    add_textbox(slide, PAD_X + Inches(0.45),
                body_top + panel_h - Inches(0.55),
                quote_w - Inches(0.90), Inches(0.4),
                text=collapse_text(by.get_text()).upper(),
                font_size=10, color=GOLD)
    # Stats column on right
    stats_left = PAD_X + quote_w + Inches(0.40)
    stats_w = CONTENT_W - quote_w - Inches(0.40)
    stat_h = (panel_h - Inches(0.30) - Inches(0.20) * (len(stats) - 1)) / len(stats)
    for i, s in enumerate(stats):
        s_top = body_top + Inches(0.10) + i * (stat_h + Inches(0.20))
        add_card(slide, stats_left, s_top, stats_w, stat_h,
                 fill=PAPER_WHITE, left_accent=SCARLET)
        num = s.find("div", class_="num")
        lbl = s.find("div", class_="lbl")
        add_textbox(slide, stats_left + Inches(0.30),
                    s_top + Inches(0.25),
                    Inches(2.4), Inches(0.9),
                    text=collapse_text(num.get_text()),
                    font_size=30, bold=True, color=SCARLET)
        add_paragraph_text(slide,
                           stats_left + Inches(0.30) + Inches(2.5),
                           s_top + Inches(0.30),
                           stats_w - Inches(2.9), stat_h - Inches(0.5),
                           extract_runs(lbl),
                           font_size=12.5, color=INK_2,
                           line_spacing=1.35)


# ---- Slide 22: Exercise C — scenario + 2 cols --------------------------


def render_slide_22_exercise_c(slide, spec: SlideSpec) -> None:
    sec = spec.section
    body_top = render_content_head(slide, sec)
    scenario = sec.find("div", class_="scenario")
    label = scenario.find("div", class_="scenario__label")
    paragraph = scenario.find("p")
    sc_top = body_top + Inches(0.15)
    sc_h = Inches(1.6)
    add_rect(slide, PAD_X, sc_top, CONTENT_W, sc_h, PAPER_WHITE,
             line=RULE_LIGHT, line_width=0.75)
    add_rect(slide, PAD_X, sc_top, Inches(0.10), sc_h, GOLD)
    add_textbox(slide, PAD_X + Inches(0.30),
                sc_top + Inches(0.18),
                CONTENT_W - Inches(0.60), Inches(0.30),
                text=collapse_text(label.get_text()).upper(),
                font_size=10, bold=True, color=GOLD_DARK)
    add_paragraph_text(slide, PAD_X + Inches(0.30),
                       sc_top + Inches(0.55),
                       CONTENT_W - Inches(0.60), Inches(1.0),
                       extract_runs(paragraph),
                       font_size=13, color=INK, line_spacing=1.4)
    # Two columns under
    cols = sec.find_all("div", class_="cols")[0].find_all(
        "div", class_="col", recursive=False)
    panel_w = (CONTENT_W - Inches(0.40)) / 2
    panel_top = sc_top + sc_h + Inches(0.30)
    panel_h = SLIDE_H - panel_top - Inches(0.55)
    head_colors = [SCARLET_DEEP, GOOD_TEXT]
    for i, col in enumerate(cols):
        left = PAD_X + i * (panel_w + Inches(0.40))
        h4 = col.find("h4")
        if h4:
            add_textbox(slide, left, panel_top, panel_w, Inches(0.4),
                        text=collapse_text(h4.get_text()),
                        font_size=14, bold=True, color=head_colors[i])
        ul = col.find("ul")
        items = [extract_runs(li) for li in ul.find_all("li")]
        add_bullets(slide, left, panel_top + Inches(0.45),
                    panel_w, panel_h - Inches(0.5),
                    items, font_size=12.5, color=INK,
                    line_spacing=1.3, space_after=6.0)


# ---- Slide 24: Quick reference — 3 cards -------------------------------


def render_slide_24_qref(slide, spec: SlideSpec) -> None:
    sec = spec.section
    body_top = render_content_head(slide, sec)
    cards = sec.find("div", class_="qref").find_all("div", class_="card")
    panel_h = SLIDE_H - body_top - Inches(1.10)
    n = len(cards)
    gap = Inches(0.20)
    card_w = (CONTENT_W - gap * (n - 1)) / n
    accent_for = {"scarlet": SCARLET, "good": GOOD_EDGE, "warn": WARN_EDGE}
    head_color_for = {"scarlet": SCARLET_DEEP, "good": GOOD_TEXT,
                      "warn": WARN_TEXT}
    for i, card in enumerate(cards):
        cls = card.get("class", [])
        kind = next((k for k in ("scarlet", "good", "warn") if k in cls),
                    "scarlet")
        left = PAD_X + i * (card_w + gap)
        add_card(slide, left, body_top, card_w, panel_h,
                 fill=PAPER_WHITE, top_accent=accent_for[kind])
        h4 = card.find("h4")
        add_textbox(slide, left + Inches(0.25),
                    body_top + Inches(0.30),
                    card_w - Inches(0.50), Inches(0.4),
                    text=collapse_text(h4.get_text()).upper(),
                    font_size=12, bold=True, color=head_color_for[kind])
        ul = card.find("ul")
        items = [extract_runs(li) for li in ul.find_all("li")]
        add_bullets(slide, left + Inches(0.25),
                    body_top + Inches(0.85),
                    card_w - Inches(0.50),
                    panel_h - Inches(1.0),
                    items, font_size=12, color=INK,
                    line_spacing=1.3, space_after=6.0)
    # Lead at bottom (default posture)
    lead = sec.find("p", class_="lead")
    if lead:
        add_paragraph_text(slide, PAD_X,
                           body_top + panel_h + Inches(0.20),
                           CONTENT_W, Inches(0.6),
                           extract_runs(lead),
                           font_size=15, color=INK,
                           align=PP_ALIGN.CENTER, line_spacing=1.3)


# ---- Slide 25: Approved tools — 3 cols + always callout ---------------


def render_slide_25_tools(slide, spec: SlideSpec) -> None:
    sec = spec.section
    body_top = render_content_head(slide, sec)
    cols = sec.find("div", class_="cols cols--3") or sec.find("div",
        class_="cols")
    cards = cols.find_all("div", class_="col", recursive=False)
    n = len(cards)
    gap = Inches(0.20)
    card_w = (CONTENT_W - gap * (n - 1)) / n
    panel_h = Inches(3.4)
    accent_for_callout = {
        "callout--scarlet": SCARLET,
        "callout--info": INFO_EDGE,
        "default": INK_3,
    }
    for i, col in enumerate(cards):
        callout = col.find("div", class_="callout")
        cls = callout.get("class", [])
        if "callout--scarlet" in cls:
            edge = SCARLET; fg = INK
        elif "callout--info" in cls:
            edge = INFO_EDGE; fg = INK
        else:
            edge = INK_3; fg = INK
        left = PAD_X + i * (card_w + gap)
        add_card(slide, left, body_top, card_w, panel_h,
                 fill=PAPER_WHITE, top_accent=edge)
        h4 = callout.find("h4")
        add_textbox(slide, left + Inches(0.25),
                    body_top + Inches(0.30),
                    card_w - Inches(0.50), Inches(0.4),
                    text=collapse_text(h4.get_text()),
                    font_size=15, bold=True, color=fg)
        p = callout.find("p")
        if p:
            add_paragraph_text(slide, left + Inches(0.25),
                               body_top + Inches(0.85),
                               card_w - Inches(0.50),
                               panel_h - Inches(1.0),
                               extract_runs(p),
                               font_size=12, color=INK_2,
                               line_spacing=1.4)
    # Bottom 'Always' warn callout
    always_callout = sec.find_all("div", class_="callout callout--warn")
    if always_callout:
        c = always_callout[-1]
        cl_top = body_top + panel_h + Inches(0.25)
        cl_h = SLIDE_H - cl_top - Inches(0.55)
        add_rect(slide, PAD_X, cl_top, CONTENT_W, cl_h, WARN_BG, line=None)
        add_rect(slide, PAD_X, cl_top, Inches(0.10), cl_h, WARN_EDGE)
        h4 = c.find("h4")
        p = c.find("p")
        if h4:
            add_textbox(slide, PAD_X + Inches(0.30),
                        cl_top + Inches(0.20),
                        CONTENT_W - Inches(0.60), Inches(0.4),
                        text=collapse_text(h4.get_text()),
                        font_size=14, bold=True, color=WARN_TEXT)
        if p:
            add_paragraph_text(slide, PAD_X + Inches(0.30),
                               cl_top + Inches(0.65),
                               CONTENT_W - Inches(0.60),
                               cl_h - Inches(0.8),
                               extract_runs(p),
                               font_size=12.5, color=WARN_TEXT,
                               line_spacing=1.4)


# ---- Slide 26: Further reading — 2 cols + lead -------------------------
# ---- Slide 27: Week 6 preview — 2 cols ---------------------------------


def render_two_cols_callouts(slide, spec: SlideSpec, *,
                             with_lead: bool = False) -> None:
    sec = spec.section
    body_top = render_content_head(slide, sec)
    cols = sec.find("div", class_="cols").find_all("div", class_="col",
                                                   recursive=False)
    panel_w = (CONTENT_W - Inches(0.40)) / 2
    panel_h_total = SLIDE_H - body_top - Inches(0.65)
    if with_lead:
        panel_h_total -= Inches(0.85)
    for i, col in enumerate(cols):
        left = PAD_X + i * (panel_w + Inches(0.40))
        h3 = col.find("h3")
        if h3:
            add_paragraph_text(slide, left, body_top + Inches(0.10),
                               panel_w, Inches(0.7),
                               extract_runs(h3),
                               font_size=20, bold=True, color=INK,
                               line_spacing=1.2)
            content_top = body_top + Inches(0.85)
        else:
            content_top = body_top + Inches(0.10)
        callout = col.find("div", class_="callout")
        ul = col.find("ul")
        if ul:
            items = [extract_runs(li) for li in ul.find_all("li")]
            add_bullets(slide, left, content_top, panel_w,
                        panel_h_total - (content_top - body_top),
                        items, font_size=13, color=INK,
                        line_spacing=1.3, space_after=8.0)
        if callout:
            cls = callout.get("class", [])
            if "callout--info" in cls:
                bg, edge, fg = INFO_BG, INFO_EDGE, INK
            elif "callout--scarlet" in cls:
                bg, edge, fg = PAPER_WHITE, SCARLET, INK
            elif "callout--good" in cls:
                bg, edge, fg = GOOD_BG, GOOD_EDGE, GOOD_TEXT
            elif "callout--warn" in cls:
                bg, edge, fg = WARN_BG, WARN_EDGE, WARN_TEXT
            else:
                bg, edge, fg = PAPER_WHITE, RULE_LIGHT, INK
            cl_top = content_top
            cl_h = body_top + panel_h_total - content_top - Inches(0.10)
            add_rect(slide, left, cl_top, panel_w, cl_h, bg,
                     line=RULE_LIGHT if bg == PAPER_WHITE else None,
                     line_width=0.5 if bg == PAPER_WHITE else None)
            add_rect(slide, left, cl_top, Inches(0.10), cl_h, edge)
            h4 = callout.find("h4")
            paras = callout.find_all("p")
            cur = cl_top + Inches(0.30)
            if h4:
                add_textbox(slide, left + Inches(0.30), cur,
                            panel_w - Inches(0.60), Inches(0.4),
                            text=collapse_text(h4.get_text()),
                            font_size=15, bold=True, color=fg)
                cur += Inches(0.55)
            for ptag in paras:
                runs = extract_runs(ptag)
                add_paragraph_text(slide, left + Inches(0.30), cur,
                                   panel_w - Inches(0.60),
                                   Inches(2.0),
                                   runs, font_size=12.5, color=fg,
                                   line_spacing=1.45)
                cur += Inches(0.85)
    if with_lead:
        lead = sec.find("p", class_="lead")
        if lead:
            add_paragraph_text(slide, PAD_X,
                               SLIDE_H - Inches(1.40),
                               CONTENT_W, Inches(0.85),
                               extract_runs(lead),
                               font_size=14, color=INK,
                               line_spacing=1.4)


def render_slide_28_closing(slide, spec: SlideSpec) -> None:
    sec = spec.section
    eyebrow = collapse_text(sec.find("div", class_="eyebrow").get_text())
    title = collapse_text(sec.find("h1").get_text())
    ask = sec.find("div", class_="ask")
    ask_label = collapse_text(ask.find("span", class_="label").get_text())
    ask_runs = extract_runs(ask.find("p"))
    build_closing(slide,
                  eyebrow=eyebrow,
                  title=title,
                  ask_label=ask_label,
                  ask_text_runs=ask_runs)


# Slides 14 and 21 share the same big-list-with-tail layout
def render_slide_14_four_questions(slide, spec):
    tail = spec.section.find("p", class_="lead")
    tail_runs = extract_runs(tail) if tail else None
    render_big_list_plain(slide, spec, tail_runs=tail_runs, tail_size=12.5)


def render_slide_21_three_checks(slide, spec):
    tail = spec.section.find("p", class_="lead")
    tail_runs = extract_runs(tail) if tail else None
    render_big_list_plain(slide, spec, tail_runs=tail_runs, tail_size=14)


# ---------------------------------------------------------------------------
# Slide dispatch table
# ---------------------------------------------------------------------------


SLIDE_BUILDERS: list[Callable] = [
    render_slide_1_cover,                # 1 Cover
    render_slide_2_audience,             # 2 Audience shift
    render_slide_3_modes,                # 3 Two delivery modes
    render_slide_4_agenda,               # 4 Agenda
    render_slide_5_core,                 # 5 Core message
    render_divider,                      # 6 Module 1 divider
    render_slide_7_data,                 # 7 Data wall
    render_slide_8_sitrep,               # 8 SITREP item
    render_divider,                      # 9 Module 2 divider
    render_two_cols_with_heads,          # 10 Yes / what kills
    lambda s, sp: render_big_list_in_card(s, sp),  # 11 What NOT to do
    render_slide_12_guardrails,          # 12 Guard rails
    render_divider,                      # 13 Module 3 divider
    render_slide_14_four_questions,      # 14 Four questions
    render_scenario_slide,               # 15 Exercise A scenario
    render_debrief_slide,                # 16 Exercise A debrief
    render_scenario_slide,               # 17 Exercise B scenario
    render_debrief_slide,                # 18 Exercise B debrief
    render_divider,                      # 19 Module 4 divider
    render_slide_20_apprentice,          # 20 Apprentice problem
    render_slide_21_three_checks,        # 21 Three supervision checks
    render_slide_22_exercise_c,          # 22 Exercise C
    render_divider,                      # 23 Module 5 divider
    render_slide_24_qref,                # 24 Quick reference
    render_slide_25_tools,               # 25 Approved tools
    lambda s, sp: render_two_cols_callouts(s, sp, with_lead=True),  # 26
    render_two_cols_callouts,            # 27 Week 6 preview
    render_slide_28_closing,             # 28 Closing
]


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------


def is_dark_layout(index: int, total: int) -> bool:
    """Return True for slides that use dark backgrounds (no top brand bar)."""
    # Cover (1), dividers (6, 9, 13, 19, 23), closing (28)
    return index in {1, 6, 9, 13, 19, 23, total}


def build_pptx() -> Path:
    if not SOURCE_HTML.exists():
        raise FileNotFoundError(f"Missing source HTML: {SOURCE_HTML}")
    OUTPUT_PPTX.parent.mkdir(parents=True, exist_ok=True)

    specs = parse_source()
    if len(specs) != len(SLIDE_BUILDERS):
        raise RuntimeError(
            f"Source deck has {len(specs)} slides but builder list has "
            f"{len(SLIDE_BUILDERS)}."
        )

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]  # Blank

    for spec, builder in zip(specs, SLIDE_BUILDERS):
        slide = prs.slides.add_slide(blank_layout)
        # Default light background
        on_dark = is_dark_layout(spec.index, spec.total)
        if not on_dark:
            set_slide_background(slide, PAPER)
            add_top_brand_bar(slide)
        # Render content
        builder(slide, spec)
        # Foot — for non-cover/divider/closing
        if not on_dark:
            add_foot(slide, spec.course_label, spec.index, spec.total,
                     on_dark=False)
        # Speaker notes (always)
        notes = extract_notes_text(spec.section)
        set_speaker_notes(slide, notes)

    prs.save(OUTPUT_PPTX)
    return OUTPUT_PPTX


def verify(path: Path) -> None:
    """Re-open the file and assert basic structural invariants."""
    prs = Presentation(path)
    width_in = prs.slide_width / 914400
    height_in = prs.slide_height / 914400
    assert abs(width_in - 13.333) < 0.01, f"width {width_in}"
    assert abs(height_in - 7.5) < 0.01, f"height {height_in}"
    n = len(prs.slides)
    assert n == 28, f"expected 28 slides, got {n}"
    notes_populated = 0
    text_frames = 0
    for s in prs.slides:
        for shape in s.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                text_frames += 1
        notes_text = s.notes_slide.notes_text_frame.text.strip()
        if notes_text:
            notes_populated += 1
    assert notes_populated == 28, f"expected 28 notes, got {notes_populated}"
    print(f"OK: {n} slides, {notes_populated} have notes, "
          f"{text_frames} editable text frames, "
          f"{width_in:.3f}\" x {height_in:.3f}\"")


def main() -> int:
    out = build_pptx()
    print(f"Wrote {out}")
    verify(out)
    return 0


if __name__ == "__main__":
    sys.exit(main())
