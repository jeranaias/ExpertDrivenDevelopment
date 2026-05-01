#!/usr/bin/env python3
"""Build a native, editable PowerPoint deck for Week 1 — AI Fluency Fundamentals.

Reads ``docs/decks/week-1-ai-fluency.html`` and emits
``docs/pptx/week-1-ai-fluency.pptx``. Every text frame, table, and bullet
on every slide is a real, editable PowerPoint object — there are no
images of slides. Each slide's ``<aside class="notes">`` block is written
to that slide's native PowerPoint speaker-notes pane.

Usage:
    python3 scripts/build_pptx_week_1.py

This script is intentionally self-contained — it does not import from any
sibling ``build_pptx_week_*.py`` files. Mild duplication across the six
week scripts is acceptable per the task brief.
"""
from __future__ import annotations

import re
import sys
from pathlib import Path

from bs4 import BeautifulSoup, NavigableString, Tag
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
HTML_PATH = ROOT / "docs" / "decks" / "week-1-ai-fluency.html"
OUT_PATH = ROOT / "docs" / "pptx" / "week-1-ai-fluency.pptx"

# ============================================================
# Theme — Week 1 dark / scarlet / gold palette.
# Mirrors the tokens in docs/decks/css/deck.css so the PPTX
# stays visually aligned with the HTML deck.
# ============================================================
COL_INK = RGBColor(0x1A, 0x1A, 0x1A)
COL_INK_SOFT = RGBColor(0x2A, 0x2A, 0x2A)
COL_SCARLET = RGBColor(0xCC, 0x00, 0x00)
COL_SCARLET_DARK = RGBColor(0xA3, 0x00, 0x00)
COL_GOLD = RGBColor(0xF5, 0xD1, 0x30)
COL_GOLD_DARK = RGBColor(0xD4, 0xB1, 0x1A)
COL_PAPER = RGBColor(0xFF, 0xFF, 0xFF)
COL_PAPER_WARM = RGBColor(0xFA, 0xF9, 0xF6)
COL_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
COL_TEXT_MUTE = RGBColor(0x55, 0x55, 0x55)
COL_TEXT_FAINT = RGBColor(0x88, 0x88, 0x88)
COL_RULE = RGBColor(0xE5, 0xE3, 0xDD)
COL_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COL_WHITE_78 = RGBColor(0xC8, 0xC8, 0xC8)
COL_WHITE_55 = RGBColor(0x8C, 0x8C, 0x8C)
COL_WHITE_RULE = RGBColor(0x40, 0x40, 0x40)
COL_CARD_BG = RGBColor(0xF6, 0xF4, 0xEE)
COL_CARD_OK_BG = RGBColor(0xEE, 0xF5, 0xEE)
COL_GOLD_PALE = RGBColor(0xFB, 0xEB, 0x8E)

# PowerPoint-safe fonts. The brief allows Calibri/Aptos/Inter; Calibri ships
# with every modern PowerPoint install on every platform and falls back
# gracefully on Keynote and LibreOffice Impress.
FONT_DISPLAY = "Calibri"
FONT_BODY = "Calibri"

# Slide canvas — 16:9 widescreen, the standard PowerPoint widescreen size.
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
PAD_X = Inches(0.95)
PAD_Y = Inches(0.55)
BAR_H = Inches(0.08)
BAR_H_COVER = Inches(0.14)
FOOT_FROM_BOTTOM = Inches(0.32)

COURSE_TAG = "COURSE 1 · AI FLUENCY FUNDAMENTALS"


# ============================================================
# HTML helpers
# ============================================================

def inline_runs(node):
    """Walk inline content of a BS4 node into a list of (text, bold, italic)
    runs. <br> becomes a ('\\n', ...) sentinel that splits paragraphs.
    Whitespace in each text segment is collapsed."""
    runs = []

    def walk(n, bold, italic):
        if isinstance(n, NavigableString):
            t = str(n)
            if t:
                runs.append((t, bold, italic))
            return
        if not isinstance(n, Tag):
            return
        if n.name == "br":
            runs.append(("\n", bold, italic))
            return
        nb = bold or n.name in ("strong", "b")
        ni = italic or n.name in ("em", "i")
        for c in n.children:
            walk(c, nb, ni)

    if node is None:
        return runs
    for c in node.children:
        walk(c, False, False)

    cleaned = []
    for text, b, i in runs:
        if text == "\n":
            cleaned.append((text, b, i))
            continue
        text = re.sub(r"\s+", " ", text)
        cleaned.append((text, b, i))
    # collapse leading/trailing whitespace runs
    while cleaned and cleaned[0][0].strip() == "" and cleaned[0][0] != "\n":
        cleaned.pop(0)
    while cleaned and cleaned[-1][0].strip() == "" and cleaned[-1][0] != "\n":
        cleaned.pop()
    return cleaned


def plain_text(node) -> str:
    if node is None:
        return ""
    return re.sub(r"\s+", " ", node.get_text(" ", strip=True)).strip()


def runs_to_string(runs) -> str:
    return "".join(r[0] for r in runs).replace("\n", " ").strip()


# ============================================================
# python-pptx text helpers
# ============================================================

def _split_runs_into_paragraphs(runs):
    """Split runs at '\\n' sentinels into a list of paragraphs (each a list
    of runs)."""
    paragraphs = [[]]
    for text, b, i in runs:
        if text == "\n":
            paragraphs.append([])
            continue
        if not text:
            continue
        paragraphs[-1].append((text, b, i))
    return paragraphs or [[]]


def set_text_frame(tf, runs_or_text, *, font_name=FONT_BODY, font_size=18,
                   color=COL_TEXT, bold=False, italic=False, align=None,
                   line_spacing=None, anchor=None, space_after=None):
    """Replace tf contents with the given runs (list of (text,bold,italic)
    tuples) or a plain string. Always produces native PowerPoint runs."""
    if isinstance(runs_or_text, str):
        runs = [(runs_or_text, bold, italic)]
    else:
        runs = list(runs_or_text)
    paragraphs = _split_runs_into_paragraphs(runs)

    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor

    # Reset the first paragraph
    p0 = tf.paragraphs[0]
    p0.text = ""
    if align is not None:
        p0.alignment = align
    if line_spacing is not None:
        p0.line_spacing = line_spacing
    if space_after is not None:
        p0.space_after = space_after

    for idx, par_runs in enumerate(paragraphs):
        p = p0 if idx == 0 else tf.add_paragraph()
        if idx > 0:
            if align is not None:
                p.alignment = align
            if line_spacing is not None:
                p.line_spacing = line_spacing
            if space_after is not None:
                p.space_after = space_after
        if not par_runs:
            r = p.add_run()
            r.text = ""
            r.font.name = font_name
            r.font.size = Pt(font_size)
            r.font.color.rgb = color
            continue
        for text, b, i in par_runs:
            r = p.add_run()
            r.text = text
            r.font.name = font_name
            r.font.size = Pt(font_size)
            r.font.color.rgb = color
            r.font.bold = bool(b or bold)
            r.font.italic = bool(i or italic)


# ============================================================
# Shape helpers
# ============================================================

def add_rect(slide, x, y, w, h, fill_color, *, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.shadow.inherit = False
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        if line_width is not None:
            shape.line.width = line_width
    return shape


def add_textbox(slide, x, y, w, h):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    return box


def set_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


# ============================================================
# Slide chrome — present on every slide
# ============================================================

def add_top_bar(slide, *, tall=False):
    h = BAR_H_COVER if tall else BAR_H
    scarlet_w = int(SLIDE_W * 0.70)
    add_rect(slide, 0, 0, scarlet_w, h, COL_SCARLET)
    add_rect(slide, scarlet_w, 0, SLIDE_W - scarlet_w, h, COL_GOLD)


def add_foot(slide, slide_num, total, *, dark_bg=False):
    color_course = COL_WHITE_55 if dark_bg else COL_TEXT_MUTE
    color_num = COL_WHITE_55 if dark_bg else COL_TEXT_FAINT
    y = SLIDE_H - FOOT_FROM_BOTTOM
    box_l = add_textbox(slide, PAD_X, y, Inches(8), Inches(0.25))
    set_text_frame(box_l.text_frame, COURSE_TAG, font_size=9, color=color_course)
    box_r = add_textbox(slide, SLIDE_W - PAD_X - Inches(2), y, Inches(2), Inches(0.25))
    set_text_frame(
        box_r.text_frame, f"{slide_num:02d} / {total:02d}",
        font_size=9, color=color_num, align=PP_ALIGN.RIGHT,
    )


def add_eyebrow(slide, x, y, w, text, *, dark_bg=False, gold=False):
    box = add_textbox(slide, x, y, w, Inches(0.32))
    color = COL_GOLD if (dark_bg or gold) else COL_SCARLET
    set_text_frame(box.text_frame, text.upper(), font_size=11, bold=True, color=color)
    return box


def add_title(slide, x, y, w, h, runs, *, color=COL_INK, size=36,
              line_spacing=1.05):
    box = add_textbox(slide, x, y, w, h)
    set_text_frame(
        box.text_frame, runs, font_name=FONT_DISPLAY, font_size=size,
        bold=True, color=color, line_spacing=line_spacing,
    )
    return box


# ============================================================
# Bullet / paragraph rendering
# ============================================================

def render_bullets(slide, x, y, w, h, items, *, color=COL_TEXT, size=18,
                   numbered=False, line_spacing=1.25):
    """Render a list of bullet items (each item is a list of runs) into one
    text frame with a bullet character before each."""
    box = add_textbox(slide, x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    # First paragraph is implicit
    first = True
    for idx, item_runs in enumerate(items, start=1):
        if first:
            p = tf.paragraphs[0]
            p.text = ""
            first = False
        else:
            p = tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(6)
        prefix = f"{idx}.  " if numbered else "▪  "
        # First run: prefix in scarlet (or gold) for visual rhythm
        pr = p.add_run()
        pr.text = prefix
        pr.font.name = FONT_BODY
        pr.font.size = Pt(size)
        pr.font.color.rgb = COL_SCARLET if not numbered else COL_INK
        pr.font.bold = True
        # Now render the item runs
        # Items may contain '\n'; flatten into spaces because bullets are 1-line logical
        for text, b, i in item_runs:
            if text == "\n":
                continue
            r = p.add_run()
            r.text = text
            r.font.name = FONT_BODY
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bool(b)
            r.font.italic = bool(i)
    return box


def render_paragraphs(slide, x, y, w, h, paragraphs, *, color=COL_TEXT,
                      size=18, line_spacing=1.3, space_after=10, italic=False):
    """Render multiple paragraphs (each a list of runs)."""
    box = add_textbox(slide, x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for par_runs in paragraphs:
        if first:
            p = tf.paragraphs[0]
            p.text = ""
            first = False
        else:
            p = tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        for text, b, i in par_runs:
            if text == "\n":
                continue
            r = p.add_run()
            r.text = text
            r.font.name = FONT_BODY
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bool(b)
            r.font.italic = bool(i or italic)
    return box


# ============================================================
# Layout builders — one per `slide--*` kind in the source HTML
# ============================================================

def build_cover(slide, num, total, data):
    set_background(slide, COL_INK)
    add_top_bar(slide, tall=True)
    x = PAD_X
    y = Inches(2.1)
    if data.get("eyebrow"):
        eb = add_textbox(slide, x, y, Inches(11), Inches(0.4))
        set_text_frame(eb.text_frame, data["eyebrow"].upper(),
                       font_size=12, bold=True, color=COL_GOLD)
        y += Inches(0.55)
    title_runs = data.get("title_runs") or [(data.get("title", ""), True, False)]
    add_title(slide, x, y, Inches(12), Inches(2.6), title_runs,
              color=COL_WHITE, size=80, line_spacing=1.0)
    y += Inches(2.4)
    # Scarlet rule
    add_rect(slide, x, y, Inches(1.6), Inches(0.06), COL_SCARLET)
    y += Inches(0.3)
    if data.get("subtitle"):
        sb = add_textbox(slide, x, y, Inches(11), Inches(1.0))
        set_text_frame(sb.text_frame, data["subtitle"], font_size=20,
                       color=COL_WHITE_78, line_spacing=1.3)
        y += Inches(0.95)
    if data.get("byline"):
        bb = add_textbox(slide, x, y, Inches(11), Inches(0.4))
        set_text_frame(bb.text_frame, data["byline"].upper(), font_size=10,
                       color=COL_WHITE_55)
    add_foot(slide, num, total, dark_bg=True)


def build_section(slide, num, total, data):
    set_background(slide, COL_INK)
    add_top_bar(slide)
    # Big module number on the left
    num_box = add_textbox(slide, PAD_X, Inches(1.1), Inches(5.5), Inches(5.0))
    set_text_frame(num_box.text_frame, data.get("module_num", ""),
                   font_name=FONT_DISPLAY, font_size=240, bold=True,
                   color=COL_SCARLET, line_spacing=0.85,
                   anchor=MSO_ANCHOR.MIDDLE)
    rx = Inches(7.0)
    ry = Inches(2.0)
    rw = Inches(5.6)
    add_eyebrow(slide, rx, ry, rw, data.get("eyebrow", ""), dark_bg=True)
    ry += Inches(0.5)
    add_title(slide, rx, ry, rw, Inches(2.8), data.get("title_runs", []),
              color=COL_WHITE, size=50, line_spacing=1.0)
    ry += Inches(2.7)
    # Meta pills
    px = rx
    for pill in data.get("meta_pills", []):
        text = pill["text"].upper()
        gold = pill.get("gold", False)
        # Approximate width — Calibri 10pt average char is ~0.06"
        approx_w = Inches(0.30 + 0.075 * len(text))
        if approx_w > Inches(5.5):
            approx_w = Inches(5.5)
        rect = add_rect(
            slide, px, ry, approx_w, Inches(0.36),
            COL_GOLD if gold else COL_INK,
            line_color=None if gold else COL_WHITE_RULE,
            line_width=Pt(0.75) if not gold else None,
        )
        tf = rect.text_frame
        tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
        set_text_frame(tf, text, font_size=10, bold=gold,
                       color=COL_INK if gold else COL_WHITE_78,
                       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        px += approx_w + Inches(0.18)
    add_foot(slide, num, total, dark_bg=True)


def build_content(slide, num, total, data):
    set_background(slide, COL_PAPER)
    add_top_bar(slide)
    x = PAD_X
    y = Inches(0.85)
    if data.get("eyebrow"):
        add_eyebrow(slide, x, y, Inches(11), data["eyebrow"])
        y += Inches(0.40)
    add_title(slide, x, y, Inches(11.4), Inches(1.6),
              data.get("title_runs", []), size=34, line_spacing=1.05)
    # Estimate title height based on character count
    title_text = runs_to_string(data.get("title_runs", []))
    title_lines = max(1, (len(title_text) // 55) + 1)
    y += Inches(0.60 * title_lines + 0.30)

    body_top = y
    body_h = SLIDE_H - body_top - Inches(0.7)
    cy = body_top
    for block in data.get("body_blocks", []):
        if block["type"] == "bullets":
            n = len(block["items"])
            h = Inches(min(0.55 * n + 0.2, body_h / Inches(1)))
            render_bullets(slide, x, cy, Inches(11.4), h, block["items"],
                           size=18, line_spacing=1.25)
            cy += h
        elif block["type"] == "numbered":
            n = len(block["items"])
            h = Inches(min(0.55 * n + 0.2, body_h / Inches(1)))
            render_bullets(slide, x, cy, Inches(11.4), h, block["items"],
                           numbered=True, size=18)
            cy += h
        elif block["type"] == "lede":
            lh = Inches(0.85)
            render_paragraphs(slide, x, cy + Inches(0.15), Inches(11.4), lh,
                              [block["runs"]], size=22, color=COL_INK,
                              line_spacing=1.25)
            cy += lh + Inches(0.15)
        elif block["type"] == "subtitle":
            sh = Inches(0.9)
            render_paragraphs(slide, x, cy + Inches(0.1), Inches(11), sh,
                              [block["runs"]], size=15, color=COL_TEXT_MUTE,
                              line_spacing=1.35)
            cy += sh + Inches(0.1)
        elif block["type"] == "para":
            render_paragraphs(slide, x, cy, Inches(11), Inches(0.9),
                              [block["runs"]], size=18)
            cy += Inches(0.9)
    add_foot(slide, num, total)


def build_stat(slide, num, total, data):
    set_background(slide, COL_PAPER_WARM)
    add_top_bar(slide)
    # Big number on the left
    nb = add_textbox(slide, PAD_X, Inches(1.5), Inches(6.2), Inches(4.8))
    # The stat may include <small> markup represented in runs as bold/italic
    # neither — we just render the plain runs but force the "small" portion smaller.
    # In practice the markup looks like: "80" + small "%". We approximate by
    # rendering the whole thing in one go.
    runs = data.get("stat_runs", [])
    # Split numeric vs trailing alpha portion — keeps the size relationship
    # close to the source.
    full = runs_to_string(runs)
    # Heuristic: split at first non-digit/non-minus/non-decimal char from the right.
    m = re.match(r"^([\u2212\-+]?[\d\.,]+)\s*(.*)$", full)
    if m:
        big = m.group(1)
        small = m.group(2)
    else:
        big = full
        small = ""
    tf = nb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ""
    p.line_spacing = 0.9
    r = p.add_run()
    r.text = big
    r.font.name = FONT_DISPLAY
    r.font.size = Pt(170)
    r.font.bold = True
    r.font.color.rgb = COL_SCARLET
    if small:
        r2 = p.add_run()
        r2.text = " " + small
        r2.font.name = FONT_DISPLAY
        r2.font.size = Pt(70)
        r2.font.bold = True
        r2.font.color.rgb = COL_INK

    # Right column
    rx = Inches(7.6); ry = Inches(2.0); rw = Inches(5.0)
    add_eyebrow(slide, rx, ry, rw, data.get("eyebrow", ""))
    ry += Inches(0.4)
    add_title(slide, rx, ry, rw, Inches(2.0), data.get("title_runs", []),
              size=26, line_spacing=1.15)
    ry += Inches(1.6)
    if data.get("subtitle_runs"):
        render_paragraphs(slide, rx, ry, rw, Inches(2.5),
                          [data["subtitle_runs"]], size=14,
                          color=COL_TEXT_MUTE, line_spacing=1.4)
    add_foot(slide, num, total)


def build_quote(slide, num, total, data):
    set_background(slide, COL_PAPER)
    add_top_bar(slide)
    cx = Inches(2.0); cy = Inches(1.6); cw = Inches(9.3)
    # Big gold quote mark
    mb = add_textbox(slide, cx, cy, Inches(2), Inches(1.8))
    set_text_frame(mb.text_frame, "\u201C", font_name=FONT_DISPLAY,
                   font_size=180, bold=True, color=COL_GOLD,
                   line_spacing=0.6)
    qy = cy + Inches(1.7)
    add_title(slide, cx, qy, cw, Inches(3.0), data.get("quote_runs", []),
              size=44, line_spacing=1.1)
    ay = qy + Inches(2.6)
    if data.get("attr"):
        ab = add_textbox(slide, cx, ay, cw, Inches(0.5))
        set_text_frame(ab.text_frame, data["attr"].upper(), font_size=11,
                       color=COL_TEXT_MUTE, bold=True)
    add_foot(slide, num, total)


def build_two(slide, num, total, data):
    set_background(slide, COL_PAPER)
    add_top_bar(slide)
    x = PAD_X
    y = Inches(0.85)
    if data.get("eyebrow"):
        add_eyebrow(slide, x, y, Inches(11), data["eyebrow"])
        y += Inches(0.4)
    add_title(slide, x, y, Inches(11.4), Inches(1.5),
              data.get("title_runs", []), size=30, line_spacing=1.1)
    title_text = runs_to_string(data.get("title_runs", []))
    title_lines = max(1, (len(title_text) // 70) + 1)
    y += Inches(0.55 * title_lines + 0.35)

    col_w = Inches(5.4)
    gap = Inches(0.4)
    col_h = SLIDE_H - y - Inches(0.85)

    if "compare_cells" in data:
        cells = data["compare_cells"]
        for idx, cell in enumerate(cells[:2]):
            cx = x + idx * (col_w + gap)
            tone = cell["tone"]
            if tone == "bad":
                bar_color = COL_TEXT_FAINT
                label_color = COL_TEXT_MUTE
            elif tone == "good":
                bar_color = COL_SCARLET
                label_color = COL_SCARLET
            else:
                bar_color = COL_TEXT_FAINT
                label_color = COL_TEXT_MUTE
            # Header bar
            add_rect(slide, cx, y, col_w, Inches(0.05), bar_color)
            # Label
            lb = add_textbox(slide, cx, y + Inches(0.1), col_w, Inches(0.4))
            set_text_frame(lb.text_frame, cell["label"].upper(), font_size=11,
                           bold=True, color=label_color)
            # Body — quote-block style, light card
            body_y = y + Inches(0.55)
            body_h = col_h - Inches(0.6)
            card = add_rect(slide, cx, body_y, col_w, body_h, COL_CARD_BG)
            tf = card.text_frame
            tf.margin_left = Inches(0.25); tf.margin_right = Inches(0.25)
            tf.margin_top = Inches(0.2); tf.margin_bottom = Inches(0.2)
            tf.vertical_anchor = MSO_ANCHOR.TOP
            set_text_frame(tf, cell["body_runs"], font_size=15,
                           color=COL_TEXT, line_spacing=1.4)
    elif "two_cols" in data:
        cells = data["two_cols"]
        for idx, cell in enumerate(cells[:2]):
            cx = x + idx * (col_w + gap)
            tone = cell["tone"]
            if tone == "scarlet":
                tcolor = COL_SCARLET
            elif tone == "gold":
                tcolor = COL_GOLD_DARK
            else:
                tcolor = COL_TEXT_MUTE
            # Title
            tb = add_textbox(slide, cx, y, col_w, Inches(0.45))
            set_text_frame(tb.text_frame, cell["title"].upper(), font_size=12,
                           bold=True, color=tcolor)
            # Underline rule
            add_rect(slide, cx, y + Inches(0.42), col_w, Inches(0.025), tcolor)
            # Body
            body_y = y + Inches(0.6)
            body_h = col_h - Inches(0.6)
            render_paragraphs(slide, cx, body_y, col_w, body_h,
                              cell["paragraphs"], size=15, color=COL_TEXT,
                              line_spacing=1.4, space_after=10)
    add_foot(slide, num, total)


def build_check(slide, num, total, data):
    set_background(slide, COL_INK)
    add_top_bar(slide)
    x = PAD_X
    y = Inches(0.95)
    # Badge — gold pill
    badge_text = data.get("badge", "Knowledge Check").upper()
    bw = Inches(0.30 + 0.085 * len(badge_text))
    if bw > Inches(5):
        bw = Inches(5)
    badge = add_rect(slide, x, y, bw, Inches(0.42), COL_GOLD)
    tf = badge.text_frame
    tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    set_text_frame(tf, badge_text, font_size=11, bold=True, color=COL_INK,
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.85)

    # Question
    qb = add_textbox(slide, x, y, Inches(11.4), Inches(2.2))
    set_text_frame(qb.text_frame, data.get("question", ""),
                   font_name=FONT_DISPLAY, font_size=30, bold=True,
                   color=COL_WHITE, line_spacing=1.15)
    q_text = data.get("question", "")
    q_lines = max(1, (len(q_text) // 50) + 1)
    y += Inches(0.55 * q_lines + 0.35)

    # Answer
    if data.get("answer_paragraphs"):
        # Gold rule above
        add_rect(slide, x, y, Inches(2.0), Inches(0.04), COL_GOLD)
        y += Inches(0.2)
        render_paragraphs(slide, x, y, Inches(11.4),
                          SLIDE_H - y - Inches(0.8),
                          data["answer_paragraphs"], size=16, color=COL_WHITE_78,
                          line_spacing=1.4, space_after=8)
    add_foot(slide, num, total, dark_bg=True)


def build_exercise(slide, num, total, data):
    set_background(slide, COL_GOLD)
    # Ink-colored top bar (per CSS)
    add_rect(slide, 0, 0, SLIDE_W, BAR_H, COL_INK)
    x = PAD_X
    y = Inches(0.9)
    # Badge + time row
    badge_text = data.get("badge", "Exercise").upper()
    bw = Inches(0.30 + 0.085 * len(badge_text))
    if bw > Inches(3.5):
        bw = Inches(3.5)
    badge = add_rect(slide, x, y, bw, Inches(0.4), COL_INK)
    tf = badge.text_frame
    tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    set_text_frame(tf, badge_text, font_size=11, bold=True, color=COL_GOLD,
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if data.get("time"):
        timeb = add_textbox(slide, x + bw + Inches(0.25), y + Inches(0.04),
                            Inches(8), Inches(0.4))
        set_text_frame(timeb.text_frame, data["time"].upper(), font_size=11,
                       bold=True, color=COL_INK)
    y += Inches(0.7)
    # Title
    add_title(slide, x, y, Inches(11.4), Inches(1.0),
              data.get("title_runs", []), size=36, line_spacing=1.05,
              color=COL_INK)
    y += Inches(1.0)

    cols = data.get("ex_cols", [])
    col_w = Inches(5.4)
    gap = Inches(0.4)
    col_h = SLIDE_H - y - Inches(0.85)
    for idx, col in enumerate(cols[:2]):
        cx = x + idx * (col_w + gap)
        # Heading
        hb = add_textbox(slide, cx, y, col_w, Inches(0.45))
        set_text_frame(hb.text_frame, col.get("heading", "").upper(),
                       font_size=12, bold=True, color=COL_INK)
        # Items
        items = col.get("items", [])
        list_kind = col.get("kind", "ul")
        items_y = y + Inches(0.5)
        items_h = col_h - Inches(0.5)
        render_bullets(slide, cx, items_y, col_w, items_h, items,
                       numbered=(list_kind == "ol"),
                       size=15, line_spacing=1.3, color=COL_INK_SOFT)
    add_foot(slide, num, total)


def build_debrief(slide, num, total, data):
    set_background(slide, COL_PAPER)
    add_top_bar(slide)
    x = PAD_X
    y = Inches(0.85)
    if data.get("eyebrow"):
        add_eyebrow(slide, x, y, Inches(11), data["eyebrow"])
        y += Inches(0.4)
    add_title(slide, x, y, Inches(11.4), Inches(1.0),
              data.get("title_runs", []), size=30, line_spacing=1.05)
    y += Inches(0.85)

    cards = data.get("cards", [])
    col_w = Inches(5.4)
    gap = Inches(0.4)
    col_h = SLIDE_H - y - Inches(0.85)
    for idx, card in enumerate(cards[:2]):
        cx = x + idx * (col_w + gap)
        ok = card.get("ok", False)
        bg = COL_CARD_OK_BG if ok else COL_CARD_BG
        accent = RGBColor(0x2E, 0x7D, 0x32) if ok else COL_SCARLET
        # Card background
        rect = add_rect(slide, cx, y, col_w, col_h, bg)
        # Accent bar on top
        add_rect(slide, cx, y, col_w, Inches(0.06), accent)
        # Label
        lb = add_textbox(slide, cx + Inches(0.25), y + Inches(0.18),
                         col_w - Inches(0.5), Inches(0.4))
        set_text_frame(lb.text_frame, card.get("label", "").upper(),
                       font_size=11, bold=True, color=accent)
        # Items
        items = card.get("items", [])
        items_y = y + Inches(0.7)
        items_h = col_h - Inches(0.85)
        render_bullets(slide, cx + Inches(0.25), items_y,
                       col_w - Inches(0.5), items_h, items,
                       size=14, line_spacing=1.35, color=COL_TEXT)
    add_foot(slide, num, total)


def build_recap(slide, num, total, data):
    set_background(slide, COL_PAPER)
    add_top_bar(slide)
    x = PAD_X
    y = Inches(0.85)
    if data.get("eyebrow"):
        add_eyebrow(slide, x, y, Inches(11), data["eyebrow"])
        y += Inches(0.4)
    add_title(slide, x, y, Inches(11.4), Inches(0.9),
              data.get("title_runs", []), size=30, line_spacing=1.05)
    y += Inches(1.0)

    cells = data.get("recap_cells", [])
    cols = 3
    rows = (len(cells) + cols - 1) // cols
    grid_w = Inches(11.4)
    cell_w = (grid_w - Inches(0.4)) / cols
    cell_h = Inches((SLIDE_H - y - Inches(0.85)) / Emu(1) / 914400 / max(rows, 1))
    cell_h = (SLIDE_H - y - Inches(0.85) - Inches(0.2 * (rows - 1))) / max(rows, 1)
    for idx, cell in enumerate(cells):
        r = idx // cols
        c = idx % cols
        cx = x + c * (cell_w + Inches(0.2))
        cy = y + r * (cell_h + Inches(0.2))
        tone = cell.get("tone", "neutral")
        if tone == "scarlet":
            bg = COL_SCARLET; num_color = COL_WHITE; label_color = COL_WHITE
        elif tone == "gold":
            bg = COL_GOLD; num_color = COL_INK; label_color = COL_INK
        else:
            bg = COL_CARD_BG; num_color = COL_INK; label_color = COL_TEXT
        rect = add_rect(slide, cx, cy, cell_w, cell_h, bg)
        # Number
        nb = add_textbox(slide, cx + Inches(0.2), cy + Inches(0.15),
                         cell_w - Inches(0.4), Inches(1.3))
        big = runs_to_string(cell.get("num_runs", []))
        m = re.match(r"^([\u2212\-+]?[\d]+)\s*(.*)$", big)
        if m:
            head = m.group(1); tail = m.group(2)
        else:
            head = big; tail = ""
        tf = nb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = ""
        p.line_spacing = 0.9
        r1 = p.add_run()
        r1.text = head
        r1.font.name = FONT_DISPLAY
        r1.font.size = Pt(60)
        r1.font.bold = True
        r1.font.color.rgb = num_color
        if tail:
            r2 = p.add_run()
            r2.text = tail
            r2.font.name = FONT_DISPLAY
            r2.font.size = Pt(28)
            r2.font.bold = True
            r2.font.color.rgb = num_color
        # Label
        lb = add_textbox(slide, cx + Inches(0.2), cy + Inches(1.45),
                         cell_w - Inches(0.4), cell_h - Inches(1.55))
        set_text_frame(lb.text_frame, cell.get("label_runs", []),
                       font_size=12, color=label_color, line_spacing=1.3)
    add_foot(slide, num, total)


def build_grid(slide, num, total, data):
    """Either an agenda (rows of time/title/meta) or a skill grid (3x2)."""
    set_background(slide, COL_PAPER)
    add_top_bar(slide)
    x = PAD_X
    y = Inches(0.85)
    if data.get("eyebrow"):
        add_eyebrow(slide, x, y, Inches(11), data["eyebrow"])
        y += Inches(0.4)
    add_title(slide, x, y, Inches(11.4), Inches(0.95),
              data.get("title_runs", []), size=30, line_spacing=1.05)
    y += Inches(0.95)

    if "agenda_rows" in data:
        rows = data["agenda_rows"]
        # Use a real PowerPoint table — this is the agenda table.
        n = len(rows)
        col_widths = [Inches(2.4), Inches(6.5), Inches(2.5)]
        table_w = sum(col_widths, Emu(0))
        row_h = Inches(0.55)
        table_h = row_h * n
        max_h = SLIDE_H - y - Inches(0.85)
        if table_h > max_h:
            row_h = Inches(max_h / Emu(1) / 914400 / n) if n else row_h
            table_h = row_h * n
        tbl_shape = slide.shapes.add_table(n, 3, x, y, table_w, table_h)
        tbl = tbl_shape.table
        for ci, w in enumerate(col_widths):
            tbl.columns[ci].width = w
        for ri, row in enumerate(rows):
            tbl.rows[ri].height = row_h
            cells = (row["time"], row["title"], row["meta"])
            for ci, val in enumerate(cells):
                cell = tbl.cell(ri, ci)
                cell.fill.solid()
                cell.fill.fore_color.rgb = COL_PAPER if ri % 2 == 0 else COL_PAPER_WARM
                cell.margin_left = Inches(0.15)
                cell.margin_right = Inches(0.15)
                cell.margin_top = Inches(0.06)
                cell.margin_bottom = Inches(0.06)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                if ci == 0:
                    set_text_frame(cell.text_frame, val, font_size=12,
                                   bold=True, color=COL_TEXT_MUTE)
                elif ci == 1:
                    color = COL_TEXT_FAINT if row.get("muted") else COL_INK
                    set_text_frame(cell.text_frame, val, font_size=14,
                                   bold=not row.get("muted"), color=color)
                else:
                    set_text_frame(cell.text_frame, val, font_size=11,
                                   color=COL_TEXT_MUTE)
    elif "skill_cells" in data:
        cells = data["skill_cells"]
        cols = 3
        rows = (len(cells) + cols - 1) // cols
        grid_w = Inches(11.4)
        cell_w = (grid_w - Inches(0.4)) / cols
        avail_h = SLIDE_H - y - Inches(0.85)
        cell_h = (avail_h - Inches(0.2 * (rows - 1))) / max(rows, 1)
        for idx, sk in enumerate(cells):
            r = idx // cols
            c = idx % cols
            cx = x + c * (cell_w + Inches(0.2))
            cy = y + r * (cell_h + Inches(0.2))
            rect = add_rect(slide, cx, cy, cell_w, cell_h, COL_CARD_BG)
            # Top accent
            add_rect(slide, cx, cy, cell_w, Inches(0.06), COL_SCARLET)
            inner_x = cx + Inches(0.25)
            inner_w = cell_w - Inches(0.5)
            # Number
            nb = add_textbox(slide, inner_x, cy + Inches(0.2),
                             inner_w, Inches(0.5))
            set_text_frame(nb.text_frame, sk.get("num", ""), font_size=18,
                           bold=True, color=COL_SCARLET)
            # Name
            nameb = add_textbox(slide, inner_x, cy + Inches(0.7),
                                inner_w, Inches(0.55))
            set_text_frame(nameb.text_frame, sk.get("name", ""),
                           font_name=FONT_DISPLAY, font_size=20, bold=True,
                           color=COL_INK)
            # Desc
            descb = add_textbox(slide, inner_x, cy + Inches(1.3),
                                inner_w, cell_h - Inches(1.45))
            set_text_frame(descb.text_frame, sk.get("desc", ""),
                           font_size=12, color=COL_TEXT, line_spacing=1.35)
    add_foot(slide, num, total)


def build_closing(slide, num, total, data):
    set_background(slide, COL_INK)
    add_top_bar(slide, tall=True)
    x = PAD_X
    y = Inches(2.2)
    if data.get("eyebrow"):
        eb = add_textbox(slide, x, y, Inches(11), Inches(0.4))
        set_text_frame(eb.text_frame, data["eyebrow"].upper(),
                       font_size=12, bold=True, color=COL_GOLD)
        y += Inches(0.55)
    add_title(slide, x, y, Inches(11.4), Inches(1.8),
              data.get("title_runs", []), color=COL_WHITE, size=70,
              line_spacing=1.0)
    y += Inches(1.7)
    add_rect(slide, x, y, Inches(1.6), Inches(0.06), COL_SCARLET)
    y += Inches(0.3)
    if data.get("subtitle"):
        sb = add_textbox(slide, x, y, Inches(11.4), Inches(1.5))
        set_text_frame(sb.text_frame, data["subtitle"], font_size=18,
                       color=COL_WHITE_78, line_spacing=1.35)
        y += Inches(1.4)
    if data.get("byline"):
        bb = add_textbox(slide, x, y, Inches(11), Inches(0.4))
        set_text_frame(bb.text_frame, data["byline"].upper(), font_size=10,
                       color=COL_WHITE_55)
    add_foot(slide, num, total, dark_bg=True)


# ============================================================
# Notes — write speaker notes to the native PowerPoint notes pane
# ============================================================

def attach_notes(slide, notes_paragraphs):
    notes_tf = slide.notes_slide.notes_text_frame
    if not notes_paragraphs:
        notes_tf.text = ""
        return
    notes_tf.text = ""
    first = True
    for par_runs in notes_paragraphs:
        if first:
            p = notes_tf.paragraphs[0]
            p.text = ""
            first = False
        else:
            p = notes_tf.add_paragraph()
        for text, b, i in par_runs:
            if text == "\n":
                continue
            r = p.add_run()
            r.text = text
            r.font.name = FONT_BODY
            r.font.size = Pt(12)
            r.font.bold = bool(b)
            r.font.italic = bool(i)


# ============================================================
# Parse the source HTML into per-slide dicts
# ============================================================

def parse_slides(html_path):
    soup = BeautifulSoup(html_path.read_text(encoding="utf-8"), "lxml")
    slides = []
    for sec in soup.select("section.slide"):
        classes = sec.get("class", [])
        kind = None
        for c in classes:
            if c.startswith("slide--"):
                kind = c.replace("slide--", "")
                break
        slide_id = sec.get("id", "")
        notes_aside = sec.find("aside", class_="notes")
        notes_runs = []
        if notes_aside:
            for p in notes_aside.find_all("p", recursive=False):
                notes_runs.append(inline_runs(p))

        data = {"id": slide_id, "kind": kind, "notes": notes_runs}

        if kind == "cover":
            cc = sec.find("div", class_="cover__content")
            if cc:
                data["eyebrow"] = plain_text(cc.find("p", class_="eyebrow"))
                tt = cc.find(["h1", "h2"])
                data["title_runs"] = inline_runs(tt) if tt else []
                data["title"] = plain_text(tt)
                data["subtitle"] = plain_text(cc.find("p", class_="subtitle"))
                data["byline"] = plain_text(cc.find("p", class_="byline"))

        elif kind == "section":
            mod = sec.find("div", class_="section__module")
            body = sec.find("div", class_="section__body")
            data["module_num"] = plain_text(mod)
            if body:
                data["eyebrow"] = plain_text(body.find("p", class_="eyebrow"))
                tt = body.find("h2")
                data["title_runs"] = inline_runs(tt) if tt else []
                data["title"] = plain_text(tt)
                pills = []
                meta = body.find("div", class_="section__meta")
                if meta:
                    for span in meta.find_all("span", class_="meta-pill"):
                        gold = "meta-pill--gold" in span.get("class", [])
                        pills.append({"text": plain_text(span), "gold": gold})
                data["meta_pills"] = pills

        elif kind == "content":
            head = sec.find("div", class_="content__head")
            body = sec.find("div", class_="content__body")
            if head:
                data["eyebrow"] = plain_text(head.find("p", class_="eyebrow"))
                tt = head.find("h2")
                data["title_runs"] = inline_runs(tt) if tt else []
            blocks = []
            if body:
                for child in body.find_all(recursive=False):
                    if child.name == "ul":
                        items = [inline_runs(li) for li in child.find_all("li", recursive=False)]
                        blocks.append({"type": "bullets", "items": items})
                    elif child.name == "ol":
                        items = [inline_runs(li) for li in child.find_all("li", recursive=False)]
                        blocks.append({"type": "numbered", "items": items})
                    elif child.name == "p":
                        cls = child.get("class", [])
                        if "lede" in cls:
                            blocks.append({"type": "lede", "runs": inline_runs(child)})
                        elif "subtitle" in cls:
                            blocks.append({"type": "subtitle", "runs": inline_runs(child)})
                        else:
                            blocks.append({"type": "para", "runs": inline_runs(child)})
            data["body_blocks"] = blocks

        elif kind == "stat":
            num = sec.find("div", class_="stat__num")
            body = sec.find("div", class_="stat__body")
            data["stat_runs"] = inline_runs(num) if num else []
            if body:
                data["eyebrow"] = plain_text(body.find("p", class_="eyebrow"))
                tt = body.find("h2")
                data["title_runs"] = inline_runs(tt) if tt else []
                st = body.find("p", class_="subtitle")
                data["subtitle_runs"] = inline_runs(st) if st else []

        elif kind == "quote":
            inner = sec.find("div", class_="quote__inner")
            if inner:
                qt = inner.find("p", class_="quote__text")
                attr = inner.find("p", class_="quote__attr")
                data["quote_runs"] = inline_runs(qt) if qt else []
                data["attr"] = plain_text(attr)

        elif kind == "two":
            head = sec.find("div", class_="two__head")
            if head:
                data["eyebrow"] = plain_text(head.find("p", class_="eyebrow"))
                tt = head.find("h2")
                data["title_runs"] = inline_runs(tt) if tt else []
            compare = sec.find("div", class_="compare")
            cols = sec.find("div", class_="two__cols")
            if compare:
                cells = []
                for cell in compare.find_all("div", class_="compare__cell", recursive=False):
                    cls = cell.get("class", [])
                    h4 = cell.find("h4")
                    qb = cell.find("div", class_="quote-block")
                    cells.append({
                        "label": plain_text(h4),
                        "tone": "bad" if "compare__cell--bad" in cls
                                 else "good" if "compare__cell--good" in cls
                                 else "neutral",
                        "body_runs": inline_runs(qb) if qb else [],
                    })
                data["compare_cells"] = cells
            elif cols:
                cells = []
                for div in cols.find_all("div", recursive=False):
                    title = div.find("div", class_="col__title")
                    body = div.find("div", class_="col__body")
                    tone = "neutral"
                    if title:
                        tcls = title.get("class", [])
                        if "col__title--scarlet" in tcls:
                            tone = "scarlet"
                        elif "col__title--gold" in tcls:
                            tone = "gold"
                    body_paras = []
                    if body:
                        for p in body.find_all("p", recursive=False):
                            body_paras.append(inline_runs(p))
                    cells.append({
                        "title": plain_text(title),
                        "tone": tone,
                        "paragraphs": body_paras,
                    })
                data["two_cols"] = cells

        elif kind == "check":
            badge = sec.find(class_="check__badge")
            q = sec.find(class_="check__q")
            ans = sec.find(class_="check__answer")
            data["badge"] = plain_text(badge)
            data["question"] = plain_text(q)
            ans_paras = []
            if ans:
                for p in ans.find_all("p", recursive=False):
                    ans_paras.append(inline_runs(p))
            data["answer_paragraphs"] = ans_paras

        elif kind == "exercise":
            badge = sec.find(class_="ex__badge")
            time_el = sec.find(class_="ex__time")
            title = sec.find("h2")
            body = sec.find("div", class_="ex__body")
            data["badge"] = plain_text(badge)
            data["time"] = plain_text(time_el)
            data["title_runs"] = inline_runs(title) if title else []
            cols = []
            if body:
                for col in body.find_all("div", class_="ex__col", recursive=False):
                    h3 = col.find("h3")
                    items_el = col.find(["ol", "ul"])
                    items = []
                    list_kind = items_el.name if items_el else "ul"
                    if items_el:
                        items = [inline_runs(li) for li in items_el.find_all("li", recursive=False)]
                    cols.append({
                        "heading": plain_text(h3),
                        "kind": list_kind,
                        "items": items,
                    })
            data["ex_cols"] = cols

        elif kind == "debrief":
            head = sec.find("div", class_="debrief__head")
            grid = sec.find("div", class_="debrief__grid")
            if head:
                data["eyebrow"] = plain_text(head.find("p", class_="eyebrow"))
                tt = head.find("h2")
                data["title_runs"] = inline_runs(tt) if tt else []
            cards = []
            if grid:
                for card in grid.find_all("div", class_="debrief__card", recursive=False):
                    ok = "debrief__card--ok" in card.get("class", [])
                    h4 = card.find("h4")
                    ul = card.find("ul")
                    items = []
                    if ul:
                        items = [inline_runs(li) for li in ul.find_all("li", recursive=False)]
                    cards.append({"label": plain_text(h4), "ok": ok, "items": items})
            data["cards"] = cards

        elif kind == "recap":
            head = sec.find("div", class_="recap__head")
            if head:
                data["eyebrow"] = plain_text(head.find("p", class_="eyebrow"))
                tt = head.find("h2")
                data["title_runs"] = inline_runs(tt) if tt else []
            cells = []
            for cell in sec.select(".recap__cell"):
                cls = cell.get("class", [])
                tone = "neutral"
                if "recap__cell--scarlet" in cls:
                    tone = "scarlet"
                elif "recap__cell--gold" in cls:
                    tone = "gold"
                num = cell.find(class_="recap__num")
                lab = cell.find(class_="recap__lab")
                cells.append({
                    "num_runs": inline_runs(num) if num else [],
                    "label_runs": inline_runs(lab) if lab else [],
                    "tone": tone,
                })
            data["recap_cells"] = cells

        elif kind == "grid":
            head = sec.find("div", class_="grid__head")
            if head:
                data["eyebrow"] = plain_text(head.find("p", class_="eyebrow"))
                tt = head.find("h2")
                data["title_runs"] = inline_runs(tt) if tt else []
            rows_el = sec.find("div", class_="grid__rows")
            skill_grid = sec.find("div", class_="skill-grid")
            if rows_el:
                rows = []
                children = list(rows_el.find_all("span", recursive=False))
                for i in range(0, len(children), 3):
                    chunk = children[i:i + 3]
                    if len(chunk) == 3:
                        rows.append({
                            "time": plain_text(chunk[0]),
                            "title": plain_text(chunk[1]),
                            "muted": "u-text-mute" in chunk[1].get("class", []),
                            "meta": plain_text(chunk[2]),
                        })
                data["agenda_rows"] = rows
            elif skill_grid:
                cells = []
                for c in skill_grid.find_all("div", class_="skill-cell", recursive=False):
                    cells.append({
                        "num": plain_text(c.find(class_="sk-num")),
                        "name": plain_text(c.find(class_="sk-name")),
                        "desc": plain_text(c.find(class_="sk-desc")),
                    })
                data["skill_cells"] = cells

        elif kind == "closing":
            inner = sec.find("div", class_="closing__inner")
            if inner:
                data["eyebrow"] = plain_text(inner.find("p", class_="eyebrow"))
                tt = inner.find("h2")
                data["title_runs"] = inline_runs(tt) if tt else []
                data["subtitle"] = plain_text(inner.find("p", class_="subtitle"))
                data["byline"] = plain_text(inner.find("p", class_="byline"))

        slides.append(data)
    return slides


# ============================================================
# Build dispatcher
# ============================================================

BUILDERS = {
    "cover": build_cover,
    "section": build_section,
    "content": build_content,
    "stat": build_stat,
    "quote": build_quote,
    "two": build_two,
    "check": build_check,
    "exercise": build_exercise,
    "debrief": build_debrief,
    "recap": build_recap,
    "grid": build_grid,
    "closing": build_closing,
}


def build_pptx(slides, out_path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]  # Blank layout
    total = len(slides)
    for idx, data in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank_layout)
        kind = data.get("kind") or "content"
        builder = BUILDERS.get(kind, build_content)
        builder(slide, idx, total, data)
        attach_notes(slide, data.get("notes", []))
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return prs


# ============================================================
# Verification — re-open and confirm structural invariants
# ============================================================

def verify(out_path, expected_count):
    prs = Presentation(str(out_path))
    assert prs.slide_width == SLIDE_W, (
        f"slide_width = {prs.slide_width}, expected {SLIDE_W}")
    assert prs.slide_height == SLIDE_H, (
        f"slide_height = {prs.slide_height}, expected {SLIDE_H}")
    assert len(prs.slides) == expected_count, (
        f"slide count = {len(prs.slides)}, expected {expected_count}")
    for i, slide in enumerate(prs.slides, start=1):
        text_frames = [s for s in slide.shapes if s.has_text_frame]
        assert text_frames, f"slide {i} has no text frames"
        notes_text = slide.notes_slide.notes_text_frame.text
        assert notes_text.strip(), f"slide {i} has empty speaker notes"
    return prs


def main():
    if not HTML_PATH.exists():
        print(f"ERROR: source HTML not found at {HTML_PATH}", file=sys.stderr)
        return 1
    print(f"Parsing {HTML_PATH}")
    slides = parse_slides(HTML_PATH)
    print(f"  -> {len(slides)} slides found")
    print(f"Building {OUT_PATH}")
    build_pptx(slides, OUT_PATH)
    print("Verifying output")
    verify(OUT_PATH, len(slides))
    size_kb = OUT_PATH.stat().st_size / 1024
    print(f"OK — wrote {OUT_PATH} ({size_kb:.1f} KB, {len(slides)} slides)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
