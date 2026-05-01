"""
Build native PPTX — Week 2 (Builder Orientation)

Reads docs/decks/week-2-builder-orientation.html and produces a real,
editable .pptx file at docs/pptx/week-2-builder-orientation.pptx.

All on-slide text is native PowerPoint text. Speaker notes from each
<div class="speaker-notes"> block become native PowerPoint speaker notes.
Self-contained: no imports from sibling build_pptx_week_*.py files.
"""
from __future__ import annotations

import html
import re
import sys
from pathlib import Path
from typing import Iterable

from bs4 import BeautifulSoup, NavigableString, Tag
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt


REPO_ROOT = Path(__file__).resolve().parent.parent
SRC = REPO_ROOT / "docs" / "decks" / "week-2-builder-orientation.html"
OUT_DIR = REPO_ROOT / "docs" / "pptx"
OUT = OUT_DIR / "week-2-builder-orientation.pptx"


# ---------- Theme ---------------------------------------------------------

# Week 2 palette (mirrors body.w2-deck variables in deck.css)
SCARLET = RGBColor(0xCC, 0x00, 0x00)
SCARLET_DK = RGBColor(0xA3, 0x00, 0x00)
GOLD = RGBColor(0xF5, 0xD1, 0x30)
GOLD_DK = RGBColor(0xD4, 0xB1, 0x1A)
AMBER = RGBColor(0xF5, 0xA6, 0x23)
INK = RGBColor(0x1A, 0x1A, 0x1A)
INK_SOFT = RGBColor(0x2A, 0x2A, 0x2A)
PAPER = RGBColor(0xF8, 0xF7, 0xF5)
PAPER_DK = RGBColor(0xEC, 0xEB, 0xE6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x6E, 0x6E, 0x6E)
RULE = RGBColor(0xD9, 0xD8, 0xD4)
WHITE_78 = RGBColor(0xC8, 0xC8, 0xC8)
WHITE_55 = RGBColor(0x8C, 0x8C, 0x8C)

# PowerPoint-safe fonts. Inter is not in PowerPoint by default;
# Calibri/Cascadia ship with PowerPoint and approximate the deck well.
FONT_DISPLAY = "Calibri"
FONT_BODY = "Calibri"
FONT_MONO = "Cascadia Mono"

# Slide size: 16:9 widescreen, 13.333 x 7.5 inches.
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------- HTML parsing helpers ------------------------------------------


def _txt(node: Tag | NavigableString | None) -> str:
    """Collapse element text, decode entities, normalise whitespace."""
    if node is None:
        return ""
    if isinstance(node, NavigableString):
        s = str(node)
    else:
        s = node.get_text(" ", strip=False)
    s = html.unescape(s)
    s = s.replace("\xa0", " ")
    s = re.sub(r"\s+", " ", s).strip()
    return s


def _multiline(node: Tag) -> str:
    """Like _txt but turns <br> into newlines."""
    parts: list[str] = []
    for child in node.descendants:
        if isinstance(child, NavigableString):
            parts.append(str(child))
        elif child.name == "br":
            parts.append("\n")
    s = html.unescape("".join(parts))
    s = s.replace("\xa0", " ")
    lines = [re.sub(r"[ \t]+", " ", ln).strip() for ln in s.split("\n")]
    return "\n".join(lines).strip("\n")


def _li_lines(ul: Tag) -> list[str]:
    out = []
    for li in ul.find_all("li", recursive=False):
        out.append(_txt(li))
    return out


def _slide_kind(section: Tag) -> str:
    classes = section.get("class", [])
    for c in classes:
        if c.startswith("slide--"):
            return c.replace("slide--", "")
    return "content"


def _notes_text(section: Tag) -> str:
    """Convert the <div class="speaker-notes"> block into plain text suitable
    for the PowerPoint notes pane. Paragraphs separated by blank lines, and
    <ul>/<li> rendered as dashed bullets."""
    notes_div = section.find("div", class_="speaker-notes")
    if notes_div is None:
        return ""
    chunks: list[str] = []
    for child in notes_div.children:
        if isinstance(child, NavigableString):
            t = _txt(child)
            if t:
                chunks.append(t)
            continue
        if child.name == "p":
            chunks.append(_txt(child))
        elif child.name in ("ul", "ol"):
            for li in child.find_all("li", recursive=False):
                chunks.append(f"  • {_txt(li)}")
        elif child.name == "div":
            chunks.append(_txt(child))
    out = "\n\n".join(c for c in chunks if c)
    return out


# ---------- Drawing primitives --------------------------------------------


def _bg(slide, color: RGBColor) -> None:
    """Force a full-bleed background by drawing a rectangle the full slide
    size. (slide.background.fill is unreliable across PowerPoint vs Keynote
    when no master is involved.)"""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False


def _band(slide, x, y, w, h, color: RGBColor):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    rect.shadow.inherit = False
    return rect


def _txtbox(
    slide,
    x,
    y,
    w,
    h,
    text: str,
    *,
    font=FONT_BODY,
    size: int = 18,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor = INK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    leading: float = 1.2,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if text else [""]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = leading
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def _bullets(
    slide,
    x,
    y,
    w,
    h,
    items: Iterable[str],
    *,
    bullet_color: RGBColor = SCARLET,
    text_color: RGBColor = INK,
    size: int = 16,
    bold_lead: bool = False,
    leading: float = 1.25,
):
    """Render bullet items as a single textbox using a simple en-dash bullet
    glyph. Keeps everything editable as native PowerPoint text."""
    items = list(items)
    if not items:
        return None
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = leading
        p.space_after = Pt(6)
        # Bullet marker
        marker = p.add_run()
        marker.text = "▎ "  # solid left bar like the deck's red bullet
        marker.font.name = FONT_BODY
        marker.font.size = Pt(size)
        marker.font.bold = True
        marker.font.color.rgb = bullet_color
        # Body
        body = p.add_run()
        body.text = item
        body.font.name = FONT_BODY
        body.font.size = Pt(size)
        body.font.bold = bold_lead
        body.font.color.rgb = text_color
    return tb


def _set_notes(slide, text: str) -> None:
    """Set native PowerPoint speaker notes for the slide."""
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.clear()
    if not text:
        notes_tf.text = ""
        return
    paragraphs = text.split("\n\n")
    for i, para in enumerate(paragraphs):
        p = notes_tf.paragraphs[0] if i == 0 else notes_tf.add_paragraph()
        run = p.add_run()
        run.text = para
        run.font.name = FONT_BODY
        run.font.size = Pt(11)
        run.font.color.rgb = INK


# ---------- Slide builders ------------------------------------------------


# Common pad: 0.6" left/right, 0.5" top/bottom for content pages.
PAD_X = Inches(0.6)
PAD_Y = Inches(0.5)
INNER_W = SLIDE_W - 2 * PAD_X


def _foot(slide, course: str, slide_num: int, total: int, *, on_dark: bool = False):
    """Course tag + slide number on every slide."""
    color = WHITE_55 if on_dark else MUTED
    _txtbox(
        slide,
        PAD_X,
        SLIDE_H - Inches(0.32),
        Inches(8),
        Inches(0.25),
        course,
        size=8,
        color=color,
        bold=True,
    )
    _txtbox(
        slide,
        SLIDE_W - PAD_X - Inches(2),
        SLIDE_H - Inches(0.32),
        Inches(2),
        Inches(0.25),
        f"{slide_num} / {total}",
        size=8,
        color=color,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )


def _topbar(slide):
    """Scarlet/gold brand bar across top edge (0.07" tall)."""
    bar_h = Inches(0.07)
    split = SLIDE_W - Inches(2.0)
    _band(slide, 0, 0, split, bar_h, SCARLET)
    _band(slide, split, 0, SLIDE_W - split, bar_h, GOLD)


def _scarlet_edge(slide):
    """Vertical scarlet edge ribbon on the left."""
    _band(slide, 0, 0, Inches(0.16), SLIDE_H, SCARLET)


# --- Layout: COVER --------------------------------------------------------


def build_cover(prs, parsed: dict, slide_num: int, total: int):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, INK)
    # Top bar a touch thicker on cover
    _band(s, 0, 0, SLIDE_W - Inches(2), Inches(0.13), SCARLET)
    _band(s, SLIDE_W - Inches(2), 0, Inches(2), Inches(0.13), GOLD)

    # Body block (top portion above the scarlet meta band)
    body_top = Inches(0.9)
    body_h = SLIDE_H - body_top - Inches(2.05)

    _txtbox(
        s,
        Inches(0.85),
        body_top,
        Inches(11.5),
        Inches(0.4),
        parsed["eyebrow"],
        size=11,
        bold=True,
        color=GOLD,
    )
    _txtbox(
        s,
        Inches(0.85),
        body_top + Inches(0.55),
        Inches(11.5),
        Inches(2.6),
        parsed["title"],
        font=FONT_DISPLAY,
        size=64,
        bold=True,
        color=WHITE,
        leading=0.95,
    )
    _txtbox(
        s,
        Inches(0.85),
        body_top + Inches(3.3),
        Inches(11.5),
        Inches(1.4),
        parsed["sub"],
        size=18,
        color=WHITE_78,
        leading=1.3,
    )

    # Bottom scarlet meta band with 4 cells
    band_top = SLIDE_H - Inches(2.0)
    _band(s, 0, band_top, SLIDE_W, Inches(2.0), SCARLET)
    cells = parsed["meta_cells"]
    cell_w = (SLIDE_W - Inches(1.7)) / 4
    cell_x0 = Inches(0.85)
    cell_y = band_top + Inches(0.45)
    for i, (label, value) in enumerate(cells):
        x = cell_x0 + cell_w * i
        _txtbox(
            s,
            x,
            cell_y,
            cell_w - Inches(0.2),
            Inches(0.3),
            label.upper(),
            size=9,
            bold=True,
            color=WHITE_78,
        )
        _txtbox(
            s,
            x,
            cell_y + Inches(0.4),
            cell_w - Inches(0.2),
            Inches(0.95),
            value,
            font=FONT_DISPLAY,
            size=16,
            bold=True,
            color=WHITE,
            leading=1.15,
        )

    _foot(s, "Course 2 · Builder Orientation", slide_num, total, on_dark=True)
    return s


def parse_cover(section: Tag) -> dict:
    eyebrow = _txt(section.select_one(".cover-eyebrow"))
    title = _multiline(section.select_one(".cover-title"))
    sub = _txt(section.select_one(".cover-sub"))
    meta_cells: list[tuple[str, str]] = []
    for cell in section.select(".cover-meta__cell"):
        meta_cells.append(
            (_txt(cell.select_one(".label")), _txt(cell.select_one(".value")))
        )
    return {
        "eyebrow": eyebrow,
        "title": title,
        "sub": sub,
        "meta_cells": meta_cells,
    }


# --- Layout: SECTION DIVIDER ---------------------------------------------


def parse_section(section: Tag) -> dict:
    return {
        "tag": _txt(section.select_one(".module-tag")),
        "title": _multiline(section.select_one(".module-title")),
        "meta": [_txt(d) for d in section.select(".module-meta > div")],
    }


def build_section(prs, parsed: dict, slide_num: int, total: int):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, INK)
    # Gold accent block
    _band(s, Inches(0.85), Inches(0.95), Inches(0.6), Inches(0.05), GOLD)
    _txtbox(
        s,
        Inches(1.6),
        Inches(0.78),
        Inches(10),
        Inches(0.4),
        parsed["tag"],
        size=12,
        bold=True,
        color=GOLD,
    )
    _txtbox(
        s,
        Inches(0.85),
        Inches(2.4),
        Inches(11.5),
        Inches(3.5),
        parsed["title"],
        font=FONT_DISPLAY,
        size=64,
        bold=True,
        color=WHITE,
        leading=0.95,
    )
    if parsed["meta"]:
        _txtbox(
            s,
            Inches(0.85),
            Inches(6.0),
            Inches(7),
            Inches(0.4),
            parsed["meta"][0],
            size=14,
            color=WHITE_78,
        )
        if len(parsed["meta"]) > 1:
            # Gold "duration"-style chip
            chip_text = parsed["meta"][1]
            chip_w = Inches(min(6.0, max(3.0, 0.16 * len(chip_text) + 0.6)))
            chip = slide_chip(s, Inches(0.85), Inches(6.5), chip_w, Inches(0.45),
                              chip_text, fg=GOLD, border=GOLD)

    _foot(s, "Course 2 · Builder Orientation", slide_num, total, on_dark=True)
    return s


def slide_chip(slide, x, y, w, h, text: str, *, fg: RGBColor, border: RGBColor):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    rect.fill.background()
    rect.line.color.rgb = border
    rect.line.width = Pt(1.5)
    rect.shadow.inherit = False
    tf = rect.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT_BODY
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = fg
    return rect


# --- Layout: CONTENT ------------------------------------------------------


def parse_content(section: Tag) -> dict:
    eyebrow = _txt(section.select_one(".slide__eyebrow"))
    title = _multiline(section.select_one(".slide__title"))
    # Detect content pattern: bullets, two-column rules+text, columns of cards,
    # five-card row, callouts.
    # Two-column "h3 + ul" pattern (slides 2, 5/bridge etc.)
    two_h3 = section.select(":scope > div > div > h3")
    cols = []
    if two_h3:
        for col_div in section.select(":scope > div > div"):
            h3 = col_div.find("h3")
            ul = col_div.find("ul", class_="bullets")
            p = col_div.find("p")
            if not h3 and not ul and not p:
                continue
            cols.append({
                "head": _txt(h3) if h3 else "",
                "head_color": _h3_color(h3) if h3 else SCARLET,
                "items": _li_lines(ul) if ul else [],
                "para": _txt(p) if p else "",
            })
    # Card grid (.col-card)
    cards = []
    for c in section.select(".col-card"):
        cards.append({
            "label": _txt(c.select_one(".col-card__label")),
            "label_color": _label_color(c.select_one(".col-card__label")),
            "title": _txt(c.select_one(".col-card__title")),
            "body": _txt(c.select_one(".col-card__body")),
        })
    # Bullets directly under section
    top_bullets = []
    direct_ul = section.find("ul", class_="bullets")
    if direct_ul:
        top_bullets = _li_lines(direct_ul)
    # Callouts
    callouts = []
    for cal in section.select(".callout"):
        kind = "gold" if "callout--gold" in cal.get("class", []) else (
            "amber" if "callout--amber" in cal.get("class", []) else "scarlet"
        )
        callouts.append({"text": _txt(cal), "kind": kind})
    return {
        "eyebrow": eyebrow,
        "title": title,
        "two_cols": cols,
        "cards": cards,
        "bullets": top_bullets if not cols else [],
        "callouts": callouts,
    }


def _h3_color(h3: Tag) -> RGBColor:
    style = (h3.get("style") or "").lower()
    if "c-gold" in style or "gold" in style:
        return GOLD_DK
    return SCARLET


def _label_color(label_tag: Tag | None) -> RGBColor:
    if label_tag is None:
        return SCARLET
    cls = label_tag.get("class", [])
    if "col-card__label--gold" in cls:
        return GOLD_DK
    return SCARLET


def build_content(prs, parsed: dict, slide_num: int, total: int):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, PAPER)
    _topbar(s)
    _scarlet_edge(s)

    y = Inches(0.7)
    if parsed["eyebrow"]:
        _txtbox(
            s,
            PAD_X,
            y,
            INNER_W,
            Inches(0.32),
            parsed["eyebrow"].upper(),
            size=11,
            bold=True,
            color=SCARLET,
        )
        y += Inches(0.42)
    if parsed["title"]:
        title_h = Inches(1.4)
        _txtbox(
            s,
            PAD_X,
            y,
            INNER_W,
            title_h,
            parsed["title"],
            font=FONT_DISPLAY,
            size=34,
            bold=True,
            color=INK,
            leading=1.05,
        )
        y += title_h + Inches(0.15)

    body_top = y
    body_h = SLIDE_H - y - Inches(0.6)

    if parsed["two_cols"]:
        col_w = (INNER_W - Inches(0.5)) / len(parsed["two_cols"])
        for i, col in enumerate(parsed["two_cols"]):
            cx = PAD_X + (col_w + Inches(0.5)) * i
            cy = body_top
            _txtbox(
                s,
                cx,
                cy,
                col_w,
                Inches(0.32),
                col["head"].upper(),
                size=11,
                bold=True,
                color=col["head_color"],
            )
            cy += Inches(0.42)
            if col["items"]:
                _bullets(
                    s,
                    cx,
                    cy,
                    col_w,
                    body_h - Inches(0.6),
                    col["items"],
                    bullet_color=col["head_color"],
                    size=14,
                )
            elif col["para"]:
                _txtbox(
                    s,
                    cx,
                    cy,
                    col_w,
                    body_h - Inches(0.6),
                    col["para"],
                    size=14,
                    color=INK,
                    leading=1.35,
                )
    elif parsed["cards"]:
        cards = parsed["cards"]
        n = len(cards)
        # Choose grid columns: 5 across for slide 11, 3 across for slide 7,
        # 2 across for 4-card layouts, 4 across for the peer-review check grid.
        if n >= 5:
            cols = n
            rows = 1
        elif n == 6:
            cols, rows = 3, 2
        elif n == 4:
            cols, rows = 2, 2
        elif n == 3:
            cols, rows = 3, 1
        else:
            cols, rows = n, 1
        gap_x = Inches(0.18)
        gap_y = Inches(0.18)
        cw = (INNER_W - gap_x * (cols - 1)) / cols
        ch = (body_h - gap_y * (rows - 1) - Inches(0.6)) / rows
        for i, c in enumerate(cards):
            r = i // cols
            cc = i % cols
            cx = PAD_X + (cw + gap_x) * cc
            cy = body_top + (ch + gap_y) * r
            _build_card(s, cx, cy, cw, ch, c)
    elif parsed["bullets"]:
        _bullets(
            s,
            PAD_X,
            body_top,
            INNER_W,
            body_h - Inches(0.6),
            parsed["bullets"],
            bullet_color=SCARLET,
            size=14,
            leading=1.3,
        )

    if parsed["callouts"]:
        cy = SLIDE_H - Inches(1.2)
        for cal in parsed["callouts"][:1]:  # one callout fits cleanly
            _build_callout(s, PAD_X, cy, INNER_W, Inches(0.6), cal)

    _foot(s, "Course 2 · Builder Orientation", slide_num, total)
    return s


def _build_card(slide, x, y, w, h, card: dict):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = WHITE
    rect.line.color.rgb = RULE
    rect.line.width = Pt(1)
    rect.shadow.inherit = False
    inner_pad = Inches(0.18)
    cx = x + inner_pad
    cy = y + inner_pad
    cw = w - 2 * inner_pad
    if card.get("label"):
        _txtbox(
            slide,
            cx,
            cy,
            cw,
            Inches(0.28),
            card["label"].upper(),
            size=9,
            bold=True,
            color=card.get("label_color") or SCARLET,
        )
        cy += Inches(0.32)
    if card.get("title"):
        title_h = Inches(0.55)
        _txtbox(
            slide,
            cx,
            cy,
            cw,
            title_h,
            card["title"],
            font=FONT_DISPLAY,
            size=14,
            bold=True,
            color=INK,
            leading=1.1,
        )
        cy += title_h
    if card.get("body"):
        body_h = (y + h) - cy - inner_pad
        _txtbox(
            slide,
            cx,
            cy,
            cw,
            body_h,
            card["body"],
            size=11,
            color=INK_SOFT,
            leading=1.3,
        )


def _build_callout(slide, x, y, w, h, cal: dict):
    color = {"gold": GOLD, "amber": AMBER, "scarlet": SCARLET}.get(
        cal["kind"], SCARLET
    )
    _band(slide, x, y, Inches(0.08), h, color)
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x + Inches(0.08), y, w - Inches(0.08), h
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    bg.shadow.inherit = False
    _txtbox(
        slide,
        x + Inches(0.25),
        y + Inches(0.08),
        w - Inches(0.4),
        h - Inches(0.16),
        cal["text"],
        size=12,
        color=INK,
        leading=1.3,
    )


# --- Layout: QUOTE --------------------------------------------------------


def parse_quote(section: Tag) -> dict:
    bq = section.find("blockquote")
    attrib = section.find(class_="attrib")
    return {"quote": _txt(bq), "attrib": _txt(attrib)}


def build_quote(prs, parsed: dict, slide_num: int, total: int):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, PAPER)
    _topbar(s)

    _txtbox(
        s,
        Inches(0.8),
        Inches(0.6),
        Inches(2.5),
        Inches(2.0),
        "“",
        font=FONT_DISPLAY,
        size=120,
        bold=True,
        color=GOLD,
        leading=0.6,
    )
    _txtbox(
        s,
        Inches(1.5),
        Inches(2.2),
        Inches(11),
        Inches(3.6),
        parsed["quote"],
        font=FONT_DISPLAY,
        size=36,
        bold=True,
        color=INK,
        leading=1.15,
    )
    _txtbox(
        s,
        Inches(1.5),
        Inches(6.1),
        Inches(11),
        Inches(0.4),
        parsed["attrib"].upper(),
        size=10,
        bold=True,
        color=MUTED,
    )
    _foot(s, "Course 2 · Builder Orientation", slide_num, total)
    return s


# --- Layout: LIVE-BUILD framing ------------------------------------------


def parse_live_build(section: Tag) -> dict:
    tag = _txt(section.select_one(".lb-tag"))
    title = _multiline(section.select_one(".slide__title"))
    lede = _txt(section.select_one(".slide__lede"))
    cols = []
    for col_div in section.select(":scope > div > div"):
        h3 = col_div.find("h3")
        ul = col_div.find("ul", class_="bullets")
        if not h3 and not ul:
            continue
        cols.append({
            "head": _txt(h3) if h3 else "",
            "items": _li_lines(ul) if ul else [],
        })
    return {"tag": tag, "title": title, "lede": lede, "cols": cols}


def build_live_build(prs, parsed: dict, slide_num: int, total: int):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, PAPER)
    _topbar(s)
    _scarlet_edge(s)

    # Red live-build tag
    tag_w = Inches(min(5.0, 0.13 * len(parsed["tag"]) + 0.6))
    tag = slide_chip(s, PAD_X, Inches(0.55), tag_w, Inches(0.4),
                     parsed["tag"].upper(), fg=WHITE, border=SCARLET)
    tag.fill.solid()
    tag.fill.fore_color.rgb = SCARLET
    _txtbox(
        s,
        PAD_X,
        Inches(1.1),
        INNER_W,
        Inches(1.4),
        parsed["title"],
        font=FONT_DISPLAY,
        size=30,
        bold=True,
        color=INK,
        leading=1.1,
    )
    if parsed["lede"]:
        _txtbox(
            s,
            PAD_X,
            Inches(2.5),
            INNER_W,
            Inches(0.8),
            parsed["lede"],
            size=14,
            color=INK_SOFT,
            leading=1.35,
        )
    if parsed["cols"]:
        col_w = (INNER_W - Inches(0.5)) / len(parsed["cols"])
        for i, col in enumerate(parsed["cols"][:2]):
            cx = PAD_X + (col_w + Inches(0.5)) * i
            cy = Inches(3.6)
            _txtbox(
                s,
                cx,
                cy,
                col_w,
                Inches(0.32),
                col["head"].upper(),
                size=10,
                bold=True,
                color=SCARLET,
            )
            _bullets(
                s,
                cx,
                cy + Inches(0.42),
                col_w,
                Inches(2.8),
                col["items"],
                bullet_color=SCARLET if i == 0 else GOLD_DK,
                size=12,
            )
    _foot(s, "Course 2 · Builder Orientation", slide_num, total)
    return s


# --- Layout: EXERCISE -----------------------------------------------------


def parse_exercise(section: Tag) -> dict:
    label = _txt(section.select_one(".ex-banner .label"))
    timer = _txt(section.select_one(".ex-banner .timer"))
    prompt = _txt(section.select_one(".ex-prompt"))
    task = _txt(section.select_one(".ex-task"))
    bullets = []
    direct_ul = section.find("ul", class_="bullets")
    if direct_ul:
        bullets = _li_lines(direct_ul)
    cards = []
    for c in section.select(".col-card"):
        cards.append({
            "label": _txt(c.select_one(".col-card__label")),
            "label_color": _label_color(c.select_one(".col-card__label")),
            "title": _txt(c.select_one(".col-card__title")),
            "body": _txt(c.select_one(".col-card__body")),
        })
    callouts = []
    for cal in section.select(".callout"):
        kind = "gold" if "callout--gold" in cal.get("class", []) else (
            "amber" if "callout--amber" in cal.get("class", []) else "scarlet"
        )
        callouts.append({"text": _txt(cal), "kind": kind})
    # The whiteboard exercise (slide 10) has nested h3+ul + h3+p
    side_cols = []
    for col_div in section.select(".ex-body > div > div"):
        h3 = col_div.find("h3")
        ul = col_div.find("ul", class_="bullets")
        p = col_div.find("p")
        if not (h3 or ul or p):
            continue
        side_cols.append({
            "head": _txt(h3) if h3 else "",
            "head_color": _h3_color(h3) if h3 else SCARLET,
            "items": _li_lines(ul) if ul else [],
            "para": _txt(p) if p else "",
        })
    return {
        "label": label,
        "timer": timer,
        "prompt": prompt,
        "task": task,
        "bullets": bullets,
        "cards": cards,
        "callouts": callouts,
        "side_cols": side_cols,
    }


def build_exercise(prs, parsed: dict, slide_num: int, total: int):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, PAPER)
    # Gold banner across top
    banner_h = Inches(0.7)
    _band(s, 0, 0, SLIDE_W, banner_h, GOLD)
    _txtbox(
        s,
        PAD_X,
        Inches(0.18),
        Inches(8),
        Inches(0.4),
        parsed["label"].upper(),
        size=12,
        bold=True,
        color=INK,
    )
    if parsed["timer"]:
        timer_w = Inches(2.2)
        timer_chip = slide_chip(
            s,
            SLIDE_W - PAD_X - timer_w,
            Inches(0.13),
            timer_w,
            Inches(0.42),
            parsed["timer"].upper(),
            fg=INK,
            border=INK,
        )

    body_top = banner_h + Inches(0.4)
    _txtbox(
        s,
        PAD_X,
        body_top,
        INNER_W,
        Inches(1.3),
        parsed["prompt"],
        font=FONT_DISPLAY,
        size=30,
        bold=True,
        color=INK,
        leading=1.1,
    )
    body_top += Inches(1.4)
    if parsed["task"]:
        _txtbox(
            s,
            PAD_X,
            body_top,
            INNER_W,
            Inches(0.8),
            parsed["task"],
            size=14,
            color=INK_SOFT,
            leading=1.3,
        )
        body_top += Inches(0.7)

    remaining_h = SLIDE_H - body_top - Inches(0.6)
    if parsed["side_cols"]:
        cw = (INNER_W - Inches(0.5)) / len(parsed["side_cols"])
        for i, col in enumerate(parsed["side_cols"]):
            cx = PAD_X + (cw + Inches(0.5)) * i
            _txtbox(
                s,
                cx,
                body_top,
                cw,
                Inches(0.32),
                col["head"].upper(),
                size=10,
                bold=True,
                color=col["head_color"],
            )
            inner_y = body_top + Inches(0.42)
            if col["items"]:
                _bullets(
                    s,
                    cx,
                    inner_y,
                    cw,
                    remaining_h - Inches(0.6),
                    col["items"],
                    bullet_color=col["head_color"],
                    size=12,
                )
            elif col["para"]:
                _txtbox(
                    s,
                    cx,
                    inner_y,
                    cw,
                    remaining_h - Inches(0.6),
                    col["para"],
                    size=13,
                    color=INK,
                    leading=1.35,
                )
    elif parsed["bullets"]:
        _bullets(
            s,
            PAD_X,
            body_top,
            INNER_W,
            remaining_h - Inches(0.4),
            parsed["bullets"],
            bullet_color=SCARLET,
            size=13,
            leading=1.3,
        )
    elif parsed["cards"]:
        cards = parsed["cards"]
        n = len(cards)
        if n == 4:
            cols, rows = 2, 2
        elif n == 3:
            cols, rows = 3, 1
        elif n == 2:
            cols, rows = 2, 1
        else:
            cols, rows = min(4, n), (n + 3) // 4
        gap = Inches(0.18)
        cw = (INNER_W - gap * (cols - 1)) / cols
        ch = (remaining_h - gap * (rows - 1) - Inches(0.4)) / rows
        for i, c in enumerate(cards):
            r = i // cols
            cc = i % cols
            cx = PAD_X + (cw + gap) * cc
            cy = body_top + (ch + gap) * r
            _build_card(s, cx, cy, cw, ch, c)

    if parsed["callouts"]:
        cy = SLIDE_H - Inches(1.0)
        _build_callout(s, PAD_X, cy, INNER_W, Inches(0.6), parsed["callouts"][0])

    _foot(s, "Course 2 · Builder Orientation", slide_num, total)
    return s


# --- Layout: LIVE-PARK ----------------------------------------------------


def parse_live_park(section: Tag) -> dict:
    bar = section.select(".park-bar > span")
    eyebrow = _txt(section.select_one(".park-eyebrow"))
    headline = _txt(section.select_one(".park-headline"))
    sub = _txt(section.select_one(".park-sub"))
    items = []
    for it in section.select(".park-checklist__item"):
        num = _txt(it.select_one(".num"))
        # body is the rest of the text after .num
        full = _txt(it)
        body = full[len(num):].strip() if full.startswith(num) else full
        items.append({"label": num, "body": body})
    foot = [_txt(s) for s in section.select(".park-foot > span")]
    return {
        "bar_left": _txt(bar[0]) if len(bar) > 0 else "",
        "bar_right": _txt(bar[1]) if len(bar) > 1 else "",
        "eyebrow": eyebrow,
        "headline": headline,
        "sub": sub,
        "items": items,
        "foot": foot,
    }


def build_live_park(prs, parsed: dict, slide_num: int, total: int):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, INK)
    # Top scarlet bar
    bar_h = Inches(0.55)
    _band(s, 0, 0, SLIDE_W, bar_h, SCARLET)
    _txtbox(
        s,
        PAD_X,
        Inches(0.13),
        Inches(8),
        Inches(0.35),
        parsed["bar_left"].upper(),
        size=11,
        bold=True,
        color=WHITE,
    )
    _txtbox(
        s,
        SLIDE_W - PAD_X - Inches(6),
        Inches(0.13),
        Inches(6),
        Inches(0.35),
        parsed["bar_right"].upper(),
        size=11,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.RIGHT,
    )

    body_top = bar_h + Inches(0.5)
    _txtbox(
        s,
        PAD_X,
        body_top,
        INNER_W,
        Inches(0.4),
        parsed["eyebrow"].upper(),
        size=11,
        bold=True,
        color=GOLD,
    )
    body_top += Inches(0.4)
    _txtbox(
        s,
        PAD_X,
        body_top,
        INNER_W,
        Inches(1.6),
        parsed["headline"],
        font=FONT_DISPLAY,
        size=46,
        bold=True,
        color=WHITE,
        leading=1.0,
    )
    body_top += Inches(1.5)
    _txtbox(
        s,
        PAD_X,
        body_top,
        INNER_W,
        Inches(1.1),
        parsed["sub"],
        size=14,
        color=WHITE_78,
        leading=1.35,
    )
    body_top += Inches(1.0)

    if parsed["items"]:
        cols = len(parsed["items"])
        gap = Inches(0.25)
        cw = (INNER_W - gap * (cols - 1)) / cols
        ch = Inches(1.5)
        for i, it in enumerate(parsed["items"]):
            cx = PAD_X + (cw + gap) * i
            cy = body_top
            rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, cy, cw, ch)
            rect.fill.background()
            rect.line.color.rgb = WHITE_55
            rect.line.width = Pt(1)
            rect.shadow.inherit = False
            inner = Inches(0.18)
            _txtbox(
                s,
                cx + inner,
                cy + inner,
                cw - 2 * inner,
                Inches(0.28),
                it["label"].upper(),
                size=9,
                bold=True,
                color=GOLD,
            )
            _txtbox(
                s,
                cx + inner,
                cy + inner + Inches(0.32),
                cw - 2 * inner,
                ch - Inches(0.5),
                it["body"],
                size=11,
                color=WHITE,
                leading=1.3,
            )

    # Bottom faint band
    foot_top = SLIDE_H - Inches(0.7)
    _band(s, 0, foot_top, SLIDE_W, Inches(0.7), RGBColor(0x26, 0x26, 0x26))
    if parsed["foot"]:
        _txtbox(
            s,
            PAD_X,
            foot_top + Inches(0.22),
            Inches(8),
            Inches(0.3),
            parsed["foot"][0].upper(),
            size=9,
            bold=True,
            color=WHITE_55,
        )
        if len(parsed["foot"]) > 1:
            _txtbox(
                s,
                SLIDE_W - PAD_X - Inches(6),
                foot_top + Inches(0.22),
                Inches(6),
                Inches(0.3),
                parsed["foot"][1].upper(),
                size=9,
                bold=True,
                color=WHITE_55,
                align=PP_ALIGN.RIGHT,
            )
    return s


# --- Layout: CHECKPOINT ---------------------------------------------------


def parse_checkpoint(section: Tag) -> dict:
    bar = section.select(".cp-banner > span")
    title = _txt(section.select_one(".cp-title"))
    lede = _txt(section.select_one(".slide__lede"))
    items = []
    for it in section.select(".cp-grid__item"):
        items.append({
            "h": _txt(it.select_one(".h")),
            "t": _txt(it.select_one(".t")),
        })
    return {
        "bar_left": _txt(bar[0]) if len(bar) > 0 else "Back to the deck",
        "bar_right": _txt(bar[1]) if len(bar) > 1 else "",
        "title": title,
        "lede": lede,
        "items": items,
    }


def build_checkpoint(prs, parsed: dict, slide_num: int, total: int):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, PAPER)
    bar_h = Inches(0.55)
    _band(s, 0, 0, SLIDE_W, bar_h, INK)
    _txtbox(
        s,
        PAD_X,
        Inches(0.14),
        Inches(8),
        Inches(0.32),
        parsed["bar_left"].upper(),
        size=11,
        bold=True,
        color=GOLD,
    )
    _txtbox(
        s,
        SLIDE_W - PAD_X - Inches(6),
        Inches(0.14),
        Inches(6),
        Inches(0.32),
        parsed["bar_right"].upper(),
        size=11,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.RIGHT,
    )

    y = bar_h + Inches(0.5)
    _txtbox(
        s,
        PAD_X,
        y,
        INNER_W,
        Inches(1.0),
        parsed["title"],
        font=FONT_DISPLAY,
        size=34,
        bold=True,
        color=INK,
        leading=1.1,
    )
    y += Inches(1.0)
    if parsed["lede"]:
        _txtbox(
            s,
            PAD_X,
            y,
            INNER_W,
            Inches(0.7),
            parsed["lede"],
            size=14,
            color=INK_SOFT,
            leading=1.35,
        )
        y += Inches(0.8)

    if parsed["items"]:
        cols = len(parsed["items"])
        gap = Inches(0.25)
        cw = (INNER_W - gap * (cols - 1)) / cols
        ch = SLIDE_H - y - Inches(0.6)
        for i, it in enumerate(parsed["items"]):
            cx = PAD_X + (cw + gap) * i
            cy = y
            rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, cy, cw, ch)
            rect.fill.solid()
            rect.fill.fore_color.rgb = WHITE
            rect.line.color.rgb = RULE
            rect.line.width = Pt(0.75)
            rect.shadow.inherit = False
            # red top border emphasised by a band
            _band(s, cx, cy, cw, Inches(0.06), SCARLET)
            inner = Inches(0.22)
            _txtbox(
                s,
                cx + inner,
                cy + inner + Inches(0.06),
                cw - 2 * inner,
                Inches(0.32),
                it["h"].upper(),
                size=11,
                bold=True,
                color=SCARLET,
            )
            _txtbox(
                s,
                cx + inner,
                cy + inner + Inches(0.5),
                cw - 2 * inner,
                ch - Inches(0.7),
                it["t"],
                size=12,
                color=INK,
                leading=1.35,
            )
    _foot(s, "Course 2 · Builder Orientation", slide_num, total)
    return s


# --- Layout: BREAK --------------------------------------------------------


def parse_break(section: Tag) -> dict:
    return {
        "eyebrow": _txt(section.select_one(".break-eyebrow")),
        "title": _txt(section.select_one(".break-title")),
        "sub": _txt(section.select_one(".break-sub")),
    }


def build_break(prs, parsed: dict, slide_num: int, total: int):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, GOLD)
    _txtbox(
        s,
        Inches(0),
        Inches(2.4),
        SLIDE_W,
        Inches(0.5),
        parsed["eyebrow"].upper(),
        size=14,
        bold=True,
        color=INK_SOFT,
        align=PP_ALIGN.CENTER,
    )
    _txtbox(
        s,
        Inches(0),
        Inches(3.0),
        SLIDE_W,
        Inches(2.5),
        parsed["title"],
        font=FONT_DISPLAY,
        size=140,
        bold=True,
        color=INK,
        align=PP_ALIGN.CENTER,
        leading=0.95,
    )
    _txtbox(
        s,
        Inches(0),
        Inches(5.6),
        SLIDE_W,
        Inches(0.7),
        parsed["sub"],
        size=18,
        color=INK,
        align=PP_ALIGN.CENTER,
        leading=1.3,
    )
    return s


# --- Layout: BREAKS (When Something Breaks - amber) ----------------------


def parse_breaks(section: Tag) -> dict:
    banner = section.select(".br-banner > span")
    title = _multiline(section.select_one(".br-title"))
    lede = _txt(section.select_one(".br-lede"))
    steps = []
    for st in section.select(".br-step"):
        steps.append({
            "num": _txt(st.select_one(".num")),
            "h": _txt(st.select_one(".h")),
            "b": _txt(st.select_one(".b")),
        })
    code = section.select_one(".code-card")
    code_text = code.get_text() if code else ""
    if code_text:
        code_text = html.unescape(code_text).strip("\n")
    callouts = []
    for cal in section.select(".callout"):
        kind = "gold" if "callout--gold" in cal.get("class", []) else (
            "amber" if "callout--amber" in cal.get("class", []) else "scarlet"
        )
        callouts.append({"text": _txt(cal), "kind": kind})
    side_bullets = []
    side_ul = section.select_one(".br-body ul.bullets")
    if side_ul:
        for li in side_ul.find_all("li", recursive=False):
            side_bullets.append(_txt(li))
    side_para = section.select_one(".br-body > div > div p:not(.br-lede)")
    return {
        "banner_left": _txt(banner[0]) if len(banner) > 0 else "",
        "banner_right": _txt(banner[1]) if len(banner) > 1 else "",
        "title": title,
        "lede": lede,
        "steps": steps,
        "code": code_text,
        "side_bullets": side_bullets,
        "side_para": _txt(side_para) if side_para else "",
        "callouts": callouts,
    }


def build_breaks(prs, parsed: dict, slide_num: int, total: int):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, INK)
    # amber banner
    banner_h = Inches(0.55)
    _band(s, 0, 0, SLIDE_W, banner_h, AMBER)
    _txtbox(
        s,
        PAD_X,
        Inches(0.14),
        Inches(8),
        Inches(0.32),
        parsed["banner_left"].upper(),
        size=11,
        bold=True,
        color=INK,
    )
    _txtbox(
        s,
        SLIDE_W - PAD_X - Inches(6),
        Inches(0.14),
        Inches(6),
        Inches(0.32),
        parsed["banner_right"].upper(),
        size=10,
        bold=True,
        color=INK,
        align=PP_ALIGN.RIGHT,
    )

    y = banner_h + Inches(0.45)
    _txtbox(
        s,
        PAD_X,
        y,
        INNER_W,
        Inches(1.4),
        parsed["title"],
        font=FONT_DISPLAY,
        size=36,
        bold=True,
        color=AMBER,
        leading=1.05,
    )
    y += Inches(1.4)
    if parsed["lede"]:
        _txtbox(
            s,
            PAD_X,
            y,
            INNER_W,
            Inches(0.9),
            parsed["lede"],
            size=14,
            color=WHITE_78,
            leading=1.35,
        )
        y += Inches(0.85)

    remaining_h = SLIDE_H - y - Inches(0.6)
    if parsed["steps"] and not parsed["code"]:
        cols = len(parsed["steps"])
        gap = Inches(0.22)
        cw = (INNER_W - gap * (cols - 1)) / cols
        ch = min(remaining_h, Inches(2.6))
        for i, st in enumerate(parsed["steps"]):
            cx = PAD_X + (cw + gap) * i
            cy = y
            rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, cy, cw, ch)
            rect.fill.solid()
            rect.fill.fore_color.rgb = RGBColor(0x33, 0x28, 0x14)
            rect.line.color.rgb = AMBER
            rect.line.width = Pt(1)
            rect.shadow.inherit = False
            inner = Inches(0.22)
            _txtbox(
                s,
                cx + inner,
                cy + inner,
                cw - 2 * inner,
                Inches(0.3),
                st["num"].upper(),
                size=9,
                bold=True,
                color=AMBER,
            )
            _txtbox(
                s,
                cx + inner,
                cy + inner + Inches(0.35),
                cw - 2 * inner,
                Inches(0.6),
                st["h"],
                font=FONT_DISPLAY,
                size=14,
                bold=True,
                color=WHITE,
                leading=1.15,
            )
            _txtbox(
                s,
                cx + inner,
                cy + inner + Inches(0.95),
                cw - 2 * inner,
                ch - Inches(1.2),
                st["b"],
                size=11,
                color=WHITE_78,
                leading=1.35,
            )
    elif parsed["code"]:
        # 2-column: code on left, bullets/para on right
        col_w = (INNER_W - Inches(0.4)) * 0.58
        right_x = PAD_X + col_w + Inches(0.4)
        right_w = INNER_W - col_w - Inches(0.4)
        ch = min(remaining_h, Inches(3.4))
        # code panel
        rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, PAD_X, y, col_w, ch)
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(0x10, 0x10, 0x10)
        rect.line.color.rgb = AMBER
        rect.line.width = Pt(0.75)
        rect.shadow.inherit = False
        _txtbox(
            s,
            PAD_X + Inches(0.2),
            y + Inches(0.18),
            col_w - Inches(0.4),
            ch - Inches(0.36),
            parsed["code"],
            font=FONT_MONO,
            size=11,
            color=WHITE,
            leading=1.35,
        )
        if parsed["side_bullets"]:
            _bullets(
                s,
                right_x,
                y,
                right_w,
                ch - Inches(0.6),
                parsed["side_bullets"],
                bullet_color=AMBER,
                text_color=WHITE,
                size=12,
                leading=1.35,
            )
        if parsed["side_para"]:
            _txtbox(
                s,
                right_x,
                y + ch - Inches(0.7),
                right_w,
                Inches(0.7),
                parsed["side_para"],
                size=11,
                italic=True,
                color=WHITE_78,
                leading=1.3,
            )

    if parsed["callouts"]:
        cy = SLIDE_H - Inches(1.0)
        cal = parsed["callouts"][0]
        # On dark backgrounds we draw a darker callout
        color = {"gold": GOLD, "amber": AMBER, "scarlet": SCARLET}.get(
            cal["kind"], AMBER
        )
        _band(s, PAD_X, cy, Inches(0.08), Inches(0.6), color)
        bg = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            PAD_X + Inches(0.08),
            cy,
            INNER_W - Inches(0.08),
            Inches(0.6),
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0x26, 0x26, 0x26)
        bg.line.fill.background()
        bg.shadow.inherit = False
        _txtbox(
            s,
            PAD_X + Inches(0.25),
            cy + Inches(0.08),
            INNER_W - Inches(0.4),
            Inches(0.45),
            cal["text"],
            size=11,
            color=WHITE,
            leading=1.3,
        )

    _foot(s, "Course 2 · Builder Orientation", slide_num, total, on_dark=True)
    return s


# --- Layout: KNOWLEDGE CHECK ---------------------------------------------


def parse_knowledge(section: Tag) -> dict:
    eyebrow = _txt(section.select_one(".slide__eyebrow"))
    title = _multiline(section.select_one(".slide__title"))
    cards = []
    for kc in section.select(".kc-card"):
        cards.append({
            "n": _txt(kc.select_one(".qnum")),
            "q": _txt(kc.select_one(".qtext")),
            "a": _txt(kc.select_one(".qansw")),
        })
    return {"eyebrow": eyebrow, "title": title, "cards": cards}


def build_knowledge(prs, parsed: dict, slide_num: int, total: int):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, PAPER)
    _topbar(s)
    _scarlet_edge(s)

    y = Inches(0.7)
    if parsed["eyebrow"]:
        _txtbox(
            s,
            PAD_X,
            y,
            INNER_W,
            Inches(0.32),
            parsed["eyebrow"].upper(),
            size=11,
            bold=True,
            color=SCARLET,
        )
        y += Inches(0.4)
    _txtbox(
        s,
        PAD_X,
        y,
        INNER_W,
        Inches(0.9),
        parsed["title"],
        font=FONT_DISPLAY,
        size=30,
        bold=True,
        color=INK,
        leading=1.1,
    )
    y += Inches(1.0)

    cards = parsed["cards"]
    cols, rows = 2, 2
    gap = Inches(0.22)
    cw = (INNER_W - gap) / cols
    ch = (SLIDE_H - y - Inches(0.6) - gap) / rows
    for i, c in enumerate(cards):
        r = i // cols
        cc = i % cols
        cx = PAD_X + (cw + gap) * cc
        cy = y + (ch + gap) * r
        rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, cy, cw, ch)
        rect.fill.solid()
        rect.fill.fore_color.rgb = WHITE
        rect.line.color.rgb = RULE
        rect.line.width = Pt(0.75)
        rect.shadow.inherit = False
        inner = Inches(0.22)
        # Red circle with number
        circ_d = Inches(0.6)
        circle = s.shapes.add_shape(
            MSO_SHAPE.OVAL, cx + inner, cy + inner, circ_d, circ_d
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = SCARLET
        circle.line.fill.background()
        circle.shadow.inherit = False
        ctf = circle.text_frame
        ctf.margin_left = Emu(0)
        ctf.margin_right = Emu(0)
        ctf.margin_top = Emu(0)
        ctf.margin_bottom = Emu(0)
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = c["n"]
        cr.font.name = FONT_DISPLAY
        cr.font.size = Pt(20)
        cr.font.bold = True
        cr.font.color.rgb = WHITE
        # Question + Answer
        text_x = cx + inner + circ_d + Inches(0.18)
        text_w = cw - (text_x - cx) - inner
        _txtbox(
            s,
            text_x,
            cy + inner,
            text_w,
            Inches(0.7),
            c["q"],
            size=14,
            bold=True,
            color=INK,
            leading=1.25,
        )
        _txtbox(
            s,
            text_x,
            cy + inner + Inches(0.7),
            text_w,
            ch - Inches(0.95),
            c["a"],
            size=11,
            color=MUTED,
            leading=1.35,
        )

    _foot(s, "Course 2 · Builder Orientation", slide_num, total)
    return s


# --- Layout: CLOSING ------------------------------------------------------


def parse_closing(section: Tag) -> dict:
    bar = section.select(".cl-bar > span")
    eyebrow = _txt(section.select_one(".cl-eyebrow"))
    title = _multiline(section.select_one(".cl-title"))
    sub = _txt(section.select_one(".cl-sub"))
    modules = []
    for m in section.select(".cl-modules__item"):
        modules.append({
            "lbl": _txt(m.select_one(".lbl")),
            "dur": _txt(m.select_one(".dur")),
        })
    foot = [_txt(s) for s in section.select(".cl-foot > span")]
    return {
        "bar_left": _txt(bar[0]) if len(bar) > 0 else "",
        "bar_right": _txt(bar[1]) if len(bar) > 1 else "",
        "eyebrow": eyebrow,
        "title": title,
        "sub": sub,
        "modules": modules,
        "foot": foot,
    }


def build_closing(prs, parsed: dict, slide_num: int, total: int):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, INK)
    bar_h = Inches(0.55)
    _band(s, 0, 0, SLIDE_W, bar_h, SCARLET)
    _txtbox(
        s,
        PAD_X,
        Inches(0.14),
        Inches(8),
        Inches(0.32),
        parsed["bar_left"].upper(),
        size=11,
        bold=True,
        color=WHITE,
    )
    _txtbox(
        s,
        SLIDE_W - PAD_X - Inches(6),
        Inches(0.14),
        Inches(6),
        Inches(0.32),
        parsed["bar_right"].upper(),
        size=11,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.RIGHT,
    )

    body_top = bar_h + Inches(0.55)
    body_h = SLIDE_H - body_top - Inches(0.9)
    col_w = (INNER_W - Inches(0.6)) / 2

    # Left column
    _txtbox(
        s,
        PAD_X,
        body_top,
        col_w,
        Inches(0.4),
        parsed["eyebrow"].upper(),
        size=11,
        bold=True,
        color=GOLD,
    )
    _txtbox(
        s,
        PAD_X,
        body_top + Inches(0.5),
        col_w,
        Inches(2.6),
        parsed["title"],
        font=FONT_DISPLAY,
        size=46,
        bold=True,
        color=WHITE,
        leading=0.95,
    )
    _txtbox(
        s,
        PAD_X,
        body_top + Inches(3.3),
        col_w,
        Inches(1.4),
        parsed["sub"],
        size=14,
        color=WHITE_78,
        leading=1.35,
    )

    # Right column: modules list
    rx = PAD_X + col_w + Inches(0.6)
    item_h = body_h / max(len(parsed["modules"]), 1)
    item_h = min(item_h, Inches(0.55))
    for i, m in enumerate(parsed["modules"]):
        iy = body_top + item_h * i
        _band(s, rx, iy, Inches(0.05), Inches(0.4), GOLD)
        _txtbox(
            s,
            rx + Inches(0.18),
            iy + Inches(0.04),
            col_w - Inches(2),
            Inches(0.32),
            m["lbl"],
            size=12,
            bold=True,
            color=WHITE,
        )
        _txtbox(
            s,
            rx + col_w - Inches(2),
            iy + Inches(0.04),
            Inches(2),
            Inches(0.32),
            m["dur"].upper(),
            size=9,
            bold=True,
            color=WHITE_55,
            align=PP_ALIGN.RIGHT,
        )

    # Bottom faint band
    foot_top = SLIDE_H - Inches(0.6)
    _band(s, 0, foot_top, SLIDE_W, Inches(0.6), RGBColor(0x26, 0x26, 0x26))
    if parsed["foot"]:
        _txtbox(
            s,
            PAD_X,
            foot_top + Inches(0.18),
            Inches(8),
            Inches(0.3),
            parsed["foot"][0].upper(),
            size=9,
            bold=True,
            color=WHITE_55,
        )
        if len(parsed["foot"]) > 1:
            _txtbox(
                s,
                SLIDE_W - PAD_X - Inches(6),
                foot_top + Inches(0.18),
                Inches(6),
                Inches(0.3),
                parsed["foot"][1].upper(),
                size=9,
                bold=True,
                color=WHITE_55,
                align=PP_ALIGN.RIGHT,
            )
    return s


# ---------- Top-level orchestration ---------------------------------------


BUILDERS: dict[str, tuple[callable, callable]] = {
    "cover": (parse_cover, build_cover),
    "section": (parse_section, build_section),
    "content": (parse_content, build_content),
    "quote": (parse_quote, build_quote),
    "live-build": (parse_live_build, build_live_build),
    "exercise": (parse_exercise, build_exercise),
    "live-park": (parse_live_park, build_live_park),
    "checkpoint": (parse_checkpoint, build_checkpoint),
    "break": (parse_break, build_break),
    "breaks": (parse_breaks, build_breaks),
    "knowledge": (parse_knowledge, build_knowledge),
    "closing": (parse_closing, build_closing),
}


def build() -> Path:
    if not SRC.exists():
        raise SystemExit(f"Source HTML not found: {SRC}")
    OUT_DIR.mkdir(parents=True, exist_ok=True)

    soup = BeautifulSoup(SRC.read_text(encoding="utf-8"), "html.parser")
    sections = soup.select("section.slide")
    total = len(sections)
    if total != 28:
        print(f"WARNING: expected 28 slides, found {total}", file=sys.stderr)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for idx, sec in enumerate(sections, start=1):
        kind = _slide_kind(sec)
        builder = BUILDERS.get(kind)
        if builder is None:
            print(f"WARNING: no builder for kind={kind!r} (slide {idx})", file=sys.stderr)
            # Generic fallback: dump title + text
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _bg(slide, PAPER)
            _topbar(slide)
            _txtbox(
                slide,
                PAD_X,
                Inches(1.0),
                INNER_W,
                Inches(2),
                _txt(sec.select_one("[data-title]")) or kind,
                font=FONT_DISPLAY,
                size=28,
                bold=True,
                color=INK,
            )
        else:
            parser, builder_fn = builder
            parsed = parser(sec)
            slide = builder_fn(prs, parsed, idx, total)

        notes = _notes_text(sec)
        _set_notes(slide, notes)

    prs.save(OUT)
    return OUT


def verify(path: Path) -> None:
    """Open the produced .pptx and assert basic structural promises."""
    prs = Presentation(path)
    assert prs.slide_width == SLIDE_W, (
        f"slide width mismatch: {prs.slide_width} vs {SLIDE_W}"
    )
    assert prs.slide_height == SLIDE_H, (
        f"slide height mismatch: {prs.slide_height} vs {SLIDE_H}"
    )
    n = len(prs.slides)
    assert n == 28, f"expected 28 slides, got {n}"
    for i, slide in enumerate(prs.slides, start=1):
        # Notes populated
        notes_text = slide.notes_slide.notes_text_frame.text
        assert notes_text and notes_text.strip(), f"slide {i} has empty notes"
        # At least one editable text frame present
        text_frames = [sh for sh in slide.shapes if sh.has_text_frame and sh.text_frame.text.strip()]
        assert text_frames, f"slide {i} has no editable text frames"
    print(f"OK · {n} slides · notes populated · 16:9 (13.333 x 7.5 in)")


if __name__ == "__main__":
    out = build()
    verify(out)
    print(f"Wrote {out.relative_to(REPO_ROOT)}")
