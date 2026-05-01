"""
Build a native, editable PowerPoint deck for Week 6 — Full-Stack AI-Assisted
Development — directly from the source HTML at docs/decks/week-6-fullstack.html.

The output lives at docs/pptx/week-6-fullstack.pptx.

Design goals:
  * Real PowerPoint text frames, bullets, and tables (NOT images of slides).
  * Speaker notes from <aside class="notes"> land in each slide's notes pane
    as native PowerPoint speaker notes.
  * 16:9 widescreen (13.333 x 7.5 inches).
  * Brand identity: white background with scarlet (#CC0000) and gold (#F5D130)
    accents — a careful approximation of the Week 6 capstone palette. The
    cover and closing slides get a soft gold/red wash that reads as the
    capstone's white->pink hue. Capstone badge treatment preserved on the
    cover, assessment divider, program-close divider, and final slide.
  * PowerPoint-safe fonts only (Calibri / Consolas).
  * Self-contained: no imports from sibling build_pptx_week_*.py files.

Run:
    python3 scripts/build_pptx_week_6.py
"""

from __future__ import annotations

import re
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterable, List, Optional, Sequence, Tuple

from bs4 import BeautifulSoup, NavigableString, Tag
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree


# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
ROOT = Path(__file__).resolve().parent.parent
HTML_SRC = ROOT / "docs" / "decks" / "week-6-fullstack.html"
PPTX_OUT = ROOT / "docs" / "pptx" / "week-6-fullstack.pptx"


# ---------------------------------------------------------------------------
# Brand palette
# ---------------------------------------------------------------------------
SCARLET = RGBColor(0xCC, 0x00, 0x00)
SCARLET_DARK = RGBColor(0xA3, 0x00, 0x00)
GOLD = RGBColor(0xF5, 0xD1, 0x30)
GOLD_DARK = RGBColor(0xD4, 0xB1, 0x1A)
INK = RGBColor(0x1A, 0x1A, 0x1A)
INK_2 = RGBColor(0x4A, 0x4A, 0x4A)
INK_3 = RGBColor(0x6E, 0x6E, 0x6E)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
PAPER_2 = RGBColor(0xF8, 0xF7, 0xF5)
PAPER_3 = RGBColor(0xEF, 0xEE, 0xE9)
LINE = RGBColor(0xD9, 0xD8, 0xD4)
LINE_2 = RGBColor(0xE8, 0xE7, 0xE3)
INFO = RGBColor(0x1A, 0x3A, 0x6B)
INFO_BG = RGBColor(0xE8, 0xF0, 0xFE)
SUCCESS = RGBColor(0x0D, 0x65, 0x2D)
SUCCESS_BG = RGBColor(0xE6, 0xF4, 0xEA)

# Soft pink-ish wash used for cover/close (gold + scarlet at low opacity over white).
COVER_WASH = RGBColor(0xFD, 0xF1, 0xEC)

FONT_SANS = "Calibri"
FONT_MONO = "Consolas"

# Slide geometry (EMU helpers)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# HTML parsing — slide model
# ---------------------------------------------------------------------------

@dataclass
class TextRun:
    text: str
    bold: bool = False
    italic: bool = False
    mono: bool = False
    color: Optional[RGBColor] = None


@dataclass
class Para:
    runs: List[TextRun] = field(default_factory=list)
    bullet: bool = False
    level: int = 0
    align: Optional[str] = None  # "left" | "center" | "right"


@dataclass
class TableSpec:
    headers: List[str]
    rows: List[List[Tuple[str, Optional[str]]]]  # cell text + optional class


@dataclass
class ArchLayer:
    no: str
    name: str
    desc_runs: List[TextRun]
    is_built: bool = False
    is_current: bool = False


@dataclass
class Slide:
    kind: str  # cover|section|break|close|editor|debrief|build|content|null
    eyebrow: Optional[str] = None
    minutes: Optional[str] = None
    title: Optional[str] = None
    subtitle: Optional[str] = None
    foot_left: Optional[str] = None
    foot_right: Optional[str] = None
    speaker_notes: str = ""

    # Cover-specific
    cover_brand: Optional[str] = None
    cover_stamp: Optional[str] = None
    cover_h1_lines: List[Tuple[str, bool]] = field(default_factory=list)  # (text, accent?)
    cover_lede: Optional[str] = None
    cover_meta: List[Tuple[str, str]] = field(default_factory=list)  # (strong, after)
    cover_pip_label: Optional[str] = None
    cover_pip_num: Optional[str] = None

    # Section-specific
    section_eyebrow: Optional[str] = None
    section_no: Optional[str] = None
    section_of: Optional[str] = None
    section_name: Optional[str] = None
    section_pills: List[Tuple[str, str]] = field(default_factory=list)  # (text, kind)
    section_blurb: Optional[str] = None

    # Break-specific
    break_lab: Optional[str] = None
    break_num: Optional[str] = None
    break_num_unit: Optional[str] = None
    break_sub: Optional[str] = None

    # Close-specific
    close_stamp: Optional[str] = None
    close_h1_runs: List[TextRun] = field(default_factory=list)
    close_lede: Optional[str] = None
    close_sig: Optional[str] = None

    # Editor-specific
    editor_badge: Optional[str] = None
    editor_h2_runs: List[TextRun] = field(default_factory=list)
    editor_what_runs: List[TextRun] = field(default_factory=list)
    editor_timer: Optional[str] = None

    # Body content (for content/build/debrief slides)
    body_blocks: List[dict] = field(default_factory=list)


# --- Inline run extraction ----------------------------------------------------

def _runs_from_node(
    node, *, bold: bool = False, italic: bool = False, mono: bool = False,
    color: Optional[RGBColor] = None,
) -> List[TextRun]:
    """Walk a BeautifulSoup node and return a flat list of TextRuns,
    preserving inline formatting (strong, em, code)."""
    runs: List[TextRun] = []
    if isinstance(node, NavigableString):
        text = str(node)
        if text:
            # Collapse internal whitespace but keep meaningful spaces.
            text = re.sub(r"\s+", " ", text)
            runs.append(TextRun(text=text, bold=bold, italic=italic, mono=mono, color=color))
        return runs
    if not isinstance(node, Tag):
        return runs
    name = node.name.lower()
    if name in ("strong", "b"):
        for c in node.children:
            runs.extend(_runs_from_node(c, bold=True, italic=italic, mono=mono, color=color))
    elif name in ("em", "i"):
        for c in node.children:
            runs.extend(_runs_from_node(c, bold=bold, italic=True, mono=mono, color=color))
    elif name == "code":
        for c in node.children:
            runs.extend(_runs_from_node(c, bold=bold, italic=italic, mono=True, color=color))
    elif name in ("br",):
        runs.append(TextRun(text="\n", bold=bold, italic=italic, mono=mono, color=color))
    else:
        for c in node.children:
            runs.extend(_runs_from_node(c, bold=bold, italic=italic, mono=mono, color=color))
    return runs


def _runs_from_tag(tag: Tag) -> List[TextRun]:
    runs = _runs_from_node(tag)
    # Tidy: strip leading/trailing whitespace on first/last run
    if runs:
        runs[0].text = runs[0].text.lstrip()
        runs[-1].text = runs[-1].text.rstrip()
    return [r for r in runs if r.text]


def _plain_text(tag: Optional[Tag]) -> str:
    if tag is None:
        return ""
    return re.sub(r"\s+", " ", tag.get_text(" ", strip=True)).strip()


# --- Speaker notes extraction -------------------------------------------------

def _notes_text(slide_tag: Tag) -> str:
    """Convert the <aside class='notes'> block into plain readable speaker notes."""
    notes = slide_tag.find("aside", class_="notes")
    if notes is None:
        return ""
    chunks: List[str] = []
    for child in notes.children:
        if isinstance(child, NavigableString):
            t = str(child).strip()
            if t:
                chunks.append(t)
            continue
        if not isinstance(child, Tag):
            continue
        cname = child.name.lower()
        cls = " ".join(child.get("class") or [])
        if cname in ("h1", "h2", "h3", "h4", "h5", "h6"):
            chunks.append("")  # blank line before heading
            chunks.append(_plain_text(child).upper())
        elif cname == "p":
            text = _plain_text(child)
            if "cue" in cls:
                # Already prefixed "TRANSITION → ..." inside a strong tag
                chunks.append("")
                chunks.append(text)
            else:
                chunks.append(text)
        elif cname in ("ul", "ol"):
            for li in child.find_all("li", recursive=False):
                bullet = "  - " if cname == "ul" else "  • "
                chunks.append(bullet + _plain_text(li))
        elif "cue" in cls:
            chunks.append("")
            chunks.append(_plain_text(child))
        else:
            t = _plain_text(child)
            if t:
                chunks.append(t)
    # Compact triple+ blank lines
    out: List[str] = []
    blank_run = 0
    for line in chunks:
        if line == "":
            blank_run += 1
            if blank_run > 1:
                continue
            out.append("")
        else:
            blank_run = 0
            out.append(line)
    return "\n".join(out).strip()


# --- Body block extraction ----------------------------------------------------

def _extract_arch(arch_tag: Tag) -> List[ArchLayer]:
    layers: List[ArchLayer] = []
    for layer in arch_tag.find_all("div", class_="arch__layer", recursive=False):
        cls = " ".join(layer.get("class") or [])
        no = _plain_text(layer.find("div", class_="arch__no"))
        name_div = layer.find("div", class_="arch__name")
        # Drop the trailing checkmark inside arch__check from the rendered name.
        if name_div is not None:
            name_clone = BeautifulSoup(str(name_div), "lxml").find("div")
            chk = name_clone.find("span", class_="arch__check")
            if chk is not None:
                chk.decompose()
            name = _plain_text(name_clone)
        else:
            name = ""
        desc_div = layer.find("div", class_="arch__desc")
        desc_runs = _runs_from_tag(desc_div) if desc_div is not None else []
        layers.append(ArchLayer(
            no=no, name=name, desc_runs=desc_runs,
            is_built="is-built" in cls, is_current="is-current" in cls,
        ))
    return layers


def _extract_table(table_tag: Tag) -> TableSpec:
    headers: List[str] = []
    thead = table_tag.find("thead")
    if thead is not None:
        for th in thead.find_all("th"):
            headers.append(_plain_text(th))
    rows: List[List[Tuple[str, Optional[str]]]] = []
    tbody = table_tag.find("tbody") or table_tag
    for tr in tbody.find_all("tr"):
        if tr.find_parent("thead") is not None:
            continue
        cells: List[Tuple[str, Optional[str]]] = []
        for td in tr.find_all(["td", "th"]):
            cls = " ".join(td.get("class") or []) or None
            cells.append((_plain_text(td), cls))
        if cells:
            rows.append(cells)
    return TableSpec(headers=headers, rows=rows)


def _extract_card(card_tag: Tag) -> dict:
    cls = " ".join(card_tag.get("class") or [])
    variant = "default"
    if "card--ink" in cls:
        variant = "ink"
    elif "card--gold" in cls:
        variant = "gold"
    elif "card--accent" in cls:
        variant = "accent"
    h3 = card_tag.find("h3")
    title = _plain_text(h3) if h3 is not None else None
    content_blocks: List[dict] = []
    for child in card_tag.children:
        if not isinstance(child, Tag):
            continue
        cn = child.name.lower()
        if cn == "h3":
            continue
        if cn == "p":
            content_blocks.append({"type": "p", "runs": _runs_from_tag(child)})
        elif cn in ("ul", "ol"):
            items = []
            for li in child.find_all("li", recursive=False):
                items.append(_runs_from_tag(li))
            content_blocks.append({
                "type": "list", "ordered": cn == "ol", "items": items,
            })
    return {"type": "card", "variant": variant, "title": title, "content": content_blocks}


def _extract_pillrow(pillrow_tag: Tag) -> List[Tuple[str, str]]:
    pills: List[Tuple[str, str]] = []
    for pill in pillrow_tag.find_all("span", class_="pill"):
        cls = " ".join(pill.get("class") or [])
        kind = "default"
        if "pill--scarlet" in cls:
            kind = "scarlet"
        elif "pill--gold" in cls:
            kind = "gold"
        elif "pill--ink" in cls:
            kind = "ink"
        pills.append((_plain_text(pill), kind))
    return pills


def _extract_stats(stats_tag: Tag) -> List[dict]:
    out = []
    for stat in stats_tag.find_all("div", class_="stat", recursive=False):
        out.append({
            "num": _plain_text(stat.find("div", class_="num")),
            "lab": _plain_text(stat.find("div", class_="lab")),
            "desc": _plain_text(stat.find("div", class_="desc")),
        })
    return out


def _extract_agenda(agenda_tag: Tag) -> List[Tuple[str, List[TextRun]]]:
    """Return list of (time, name-runs) — name may include a small description span."""
    items: List[Tuple[str, List[TextRun]]] = []
    for item in agenda_tag.find_all("div", class_="agenda__item", recursive=False):
        time = _plain_text(item.find("span", class_="agenda__time"))
        name_span = item.find("span", class_="agenda__name")
        runs: List[TextRun] = []
        if name_span is not None:
            # Inner span holds the description in muted style; everything else is the headline.
            inner = name_span.find("span")
            inner_text = _plain_text(inner) if inner is not None else ""
            if inner is not None:
                inner.decompose()
            head = _plain_text(name_span)
            runs.append(TextRun(text=head, bold=True))
            if inner_text:
                runs.append(TextRun(text="  " + inner_text, color=INK_3))
        items.append((time, runs))
    return items


def _extract_pullquote(pq_tag: Tag) -> Tuple[str, str]:
    attr = pq_tag.find("span", class_="attr")
    attr_text = _plain_text(attr) if attr is not None else ""
    # Drop the attr span out of the body for a clean quote string.
    pq_clone = BeautifulSoup(str(pq_tag), "lxml")
    inner_attr = pq_clone.find("span", class_="attr")
    if inner_attr is not None:
        inner_attr.decompose()
    quote = _plain_text(pq_clone)
    return quote, attr_text


def _extract_body_blocks(body_tag: Tag) -> List[dict]:
    """Walk the .slide__body subtree and return an ordered list of block specs."""
    blocks: List[dict] = []

    def classify(child: Tag) -> Optional[dict]:
        """Return a block dict for `child` if it matches a known component,
        otherwise None (caller should recurse or treat as paragraph)."""
        cls_list = child.get("class") or []
        cls = " ".join(cls_list)
        cn = child.name.lower()
        if cn == "table" and "compare" in cls:
            return {"type": "table", "spec": _extract_table(child)}
        if "arch" in cls_list:
            return {"type": "arch", "layers": _extract_arch(child)}
        if "stats" in cls_list:
            return {"type": "stats", "items": _extract_stats(child)}
        if "agenda" in cls_list:
            return {"type": "agenda", "items": _extract_agenda(child)}
        if "card" in cls_list:
            return _extract_card(child)
        if "prompt" in cls_list:
            return {
                "type": "prompt",
                "label": _plain_text(child.find("strong")),
                "runs": _runs_from_tag_excluding(child, "strong"),
            }
        if "pillrow" in cls_list:
            return {"type": "pillrow", "pills": _extract_pillrow(child)}
        if "pullquote" in cls_list:
            quote, attr = _extract_pullquote(child)
            return {"type": "pullquote", "quote": quote, "attr": attr}
        if "arch-callout" in cls_list:
            return {"type": "arch_callout", "runs": _runs_from_tag(child)}
        if cn == "ul" and "bullets-lg" in cls:
            items = [_runs_from_tag(li) for li in child.find_all("li", recursive=False)]
            return {"type": "bullets", "items": items}
        if cn == "ul" and "checklist" in cls:
            items = [_runs_from_tag(li) for li in child.find_all("li", recursive=False)]
            return {"type": "checklist", "items": items}
        if cn in ("ul", "ol"):
            items = [_runs_from_tag(li) for li in child.find_all("li", recursive=False)]
            return {"type": "list", "ordered": cn == "ol", "items": items}
        if cn == "p":
            runs = _runs_from_tag(child)
            if runs:
                return {"type": "para", "runs": runs, "cls": cls}
            return None
        if cn == "span" and "build-banner" in cls:
            return {"type": "build_banner", "text": _plain_text(child)}
        return None

    def walk(node: Tag):
        for child in node.children:
            if isinstance(child, NavigableString):
                t = str(child).strip()
                if t:
                    blocks.append({"type": "para", "runs": [TextRun(text=t)]})
                continue
            if not isinstance(child, Tag):
                continue
            cls_list = child.get("class") or []
            cls = " ".join(cls_list)
            cn = child.name.lower()

            if "cols-2" in cls_list:
                # Two-column layout — process each child column.
                col_kind = "5050"
                if "cols-2--6040" in cls:
                    col_kind = "6040"
                elif "cols-2--4060" in cls:
                    col_kind = "4060"
                cols: List[List[dict]] = []
                for col in child.find_all("div", recursive=False):
                    saved = blocks[:]
                    blocks.clear()
                    # Check the column itself first — many decks place a
                    # single component (arch, card, prompt, pillrow, …)
                    # directly as the column root rather than wrapping it
                    # in a generic <div>.
                    col_block = classify(col)
                    if col_block is not None:
                        blocks.append(col_block)
                    else:
                        walk(col)
                    cols.append(blocks[:])
                    blocks.clear()
                    blocks.extend(saved)
                blocks.append({"type": "cols2", "kind": col_kind, "cols": cols})
                continue

            block = classify(child)
            if block is not None:
                blocks.append(block)
                continue

            if cn == "div":
                # Generic wrapper — recurse.
                walk(child)
            else:
                runs = _runs_from_tag(child)
                if runs:
                    blocks.append({"type": "para", "runs": runs})

    walk(body_tag)
    return blocks


def _runs_from_tag_excluding(tag: Tag, excluded_name: str) -> List[TextRun]:
    """Like _runs_from_tag but skip any direct child whose tag name matches."""
    runs: List[TextRun] = []
    for child in tag.children:
        if isinstance(child, Tag) and child.name.lower() == excluded_name.lower():
            continue
        runs.extend(_runs_from_node(child))
    if runs:
        runs[0].text = runs[0].text.lstrip()
        runs[-1].text = runs[-1].text.rstrip()
    return [r for r in runs if r.text]


# --- Slide classification ----------------------------------------------------

def _split_eyebrow(text: str) -> str:
    """Eyebrow text in HTML uses a · separator span 'dot'. Plain-text already
    drops the span; we just collapse whitespace."""
    return re.sub(r"\s+", " ", text).strip()


def _foot_parts(slide_tag: Tag) -> Tuple[Optional[str], Optional[str]]:
    foot = slide_tag.find("div", class_="slide__foot")
    if foot is None:
        return None, None
    spans = foot.find_all("span", recursive=False)
    if len(spans) >= 2:
        # First span has a brand-mark span inside; strip it.
        left_clone = BeautifulSoup(str(spans[0]), "lxml")
        bm = left_clone.find("span", class_="brand-mark")
        if bm is not None:
            bm.decompose()
        left = _plain_text(left_clone)
        right = _plain_text(spans[1])
        return left, right
    return None, None


def parse_slides(html: str) -> List[Slide]:
    soup = BeautifulSoup(html, "lxml")
    slide_tags = soup.find_all("section", class_="slide")
    slides: List[Slide] = []
    for tag in slide_tags:
        cls = set(tag.get("class") or [])
        notes_text = _notes_text(tag)
        if "slide--cover" in cls:
            slides.append(_parse_cover(tag, notes_text))
        elif "slide--section" in cls:
            slides.append(_parse_section(tag, notes_text))
        elif "slide--break" in cls:
            slides.append(_parse_break(tag, notes_text))
        elif "slide--close" in cls:
            slides.append(_parse_close(tag, notes_text))
        elif "slide--editor" in cls:
            slides.append(_parse_editor(tag, notes_text))
        else:
            kind = "content"
            if "slide--debrief" in cls:
                kind = "debrief"
            elif "slide--build" in cls:
                kind = "build"
            slides.append(_parse_content(tag, kind, notes_text))
    return slides


def _parse_cover(tag: Tag, notes: str) -> Slide:
    s = Slide(kind="cover", speaker_notes=notes)
    inner = tag.find("div", class_="cover-inner")
    top = inner.find("div", class_="cover-top")
    if top is not None:
        brand = top.find("div", class_="brand")
        s.cover_brand = _plain_text(brand) if brand else None
        stamp = top.find("div", class_="stamp")
        s.cover_stamp = _plain_text(stamp) if stamp else None
    mid = inner.find("div", class_="cover-mid")
    if mid is not None:
        h1 = mid.find("h1")
        if h1 is not None:
            # Walk h1 in order, splitting accent vs non-accent.
            for child in h1.children:
                if isinstance(child, NavigableString):
                    txt = str(child).strip()
                    if txt:
                        s.cover_h1_lines.append((txt, False))
                elif isinstance(child, Tag):
                    if "h1-accent" in (child.get("class") or []):
                        s.cover_h1_lines.append((_plain_text(child), True))
                    else:
                        t = _plain_text(child)
                        if t:
                            s.cover_h1_lines.append((t, False))
        lede = mid.find("p", class_="lede")
        s.cover_lede = _plain_text(lede) if lede else None
    bot = inner.find("div", class_="cover-bottom")
    if bot is not None:
        meta = bot.find("div", class_="meta")
        if meta is not None:
            for d in meta.find_all("div", recursive=False):
                strong = d.find("strong")
                strong_text = _plain_text(strong) if strong else ""
                clone = BeautifulSoup(str(d), "lxml").find("div")
                s_clone = clone.find("strong")
                if s_clone is not None:
                    s_clone.decompose()
                rest = _plain_text(clone)
                s.cover_meta.append((strong_text, rest))
        pip = bot.find("div", class_="week-pip")
        if pip is not None:
            num = pip.find("span", class_="six")
            s.cover_pip_num = _plain_text(num) if num else None
            clone = BeautifulSoup(str(pip), "lxml").find("div")
            num_clone = clone.find("span", class_="six")
            if num_clone is not None:
                num_clone.decompose()
            s.cover_pip_label = _plain_text(clone)
    return s


def _parse_section(tag: Tag, notes: str) -> Slide:
    s = Slide(kind="section", speaker_notes=notes)
    left = tag.find("div", class_="section-left")
    right = tag.find("div", class_="section-right")
    if left is not None:
        eb = left.find("div", class_="module-eyebrow")
        s.section_eyebrow = _plain_text(eb) if eb else None
        no = left.find("div", class_="module-no")
        if no is not None:
            of = no.find("span", class_="of")
            of_text = _plain_text(of) if of else ""
            clone = BeautifulSoup(str(no), "lxml").find("div")
            of_clone = clone.find("span", class_="of")
            if of_clone is not None:
                of_clone.decompose()
            s.section_no = _plain_text(clone)
            s.section_of = of_text
    if right is not None:
        name = right.find("div", class_="module-name")
        s.section_name = _plain_text(name) if name else None
        meta = right.find("div", class_="module-meta")
        if meta is not None:
            for pill in meta.find_all("span", class_="pill"):
                pcls = " ".join(pill.get("class") or [])
                kind = "outline"
                if "pill--build" in pcls:
                    kind = "build"
                elif "pill--gold" in pcls:
                    kind = "gold"
                s.section_pills.append((_plain_text(pill), kind))
        blurb = right.find("p", class_="module-blurb")
        s.section_blurb = _plain_text(blurb) if blurb else None
    return s


def _parse_break(tag: Tag, notes: str) -> Slide:
    s = Slide(kind="break", speaker_notes=notes)
    inner = tag.find("div", class_="break-inner")
    if inner is not None:
        lab = inner.find("div", class_="lab")
        s.break_lab = _plain_text(lab) if lab else None
        num = inner.find("div", class_="num")
        if num is not None:
            small = num.find("small")
            small_text = _plain_text(small) if small else ""
            clone = BeautifulSoup(str(num), "lxml").find("div")
            sc = clone.find("small")
            if sc is not None:
                sc.decompose()
            s.break_num = _plain_text(clone)
            s.break_num_unit = small_text
        sub = inner.find("div", class_="sub")
        s.break_sub = _plain_text(sub) if sub else None
    return s


def _parse_close(tag: Tag, notes: str) -> Slide:
    s = Slide(kind="close", speaker_notes=notes)
    inner = tag.find("div", class_="close-inner")
    if inner is not None:
        stamp = inner.find("span", class_="stamp")
        s.close_stamp = _plain_text(stamp) if stamp else None
        h1 = inner.find("h1")
        if h1 is not None:
            # Convert <br> + <em> into runs with line breaks.
            s.close_h1_runs = []
            for child in h1.children:
                if isinstance(child, NavigableString):
                    t = re.sub(r"\s+", " ", str(child))
                    if t.strip():
                        s.close_h1_runs.append(TextRun(text=t, bold=True))
                elif isinstance(child, Tag):
                    if child.name.lower() == "br":
                        s.close_h1_runs.append(TextRun(text="\n", bold=True))
                    elif child.name.lower() == "em":
                        s.close_h1_runs.append(
                            TextRun(text=_plain_text(child), bold=True, color=GOLD))
                    else:
                        s.close_h1_runs.append(TextRun(text=_plain_text(child), bold=True))
            # Trim leading/trailing whitespace adjacent to line breaks.
            for i, run in enumerate(s.close_h1_runs):
                if run.text == "\n":
                    if i > 0:
                        s.close_h1_runs[i - 1].text = s.close_h1_runs[i - 1].text.rstrip()
                    if i + 1 < len(s.close_h1_runs):
                        s.close_h1_runs[i + 1].text = s.close_h1_runs[i + 1].text.lstrip()
        lede = inner.find("p", class_="lede")
        s.close_lede = _plain_text(lede) if lede else None
        sig = inner.find("div", class_="sig")
        s.close_sig = _plain_text(sig) if sig else None
    return s


def _parse_editor(tag: Tag, notes: str) -> Slide:
    s = Slide(kind="editor", speaker_notes=notes)
    inner = tag.find("div", class_="editor-inner")
    if inner is not None:
        badge = inner.find("span", class_="badge")
        s.editor_badge = _plain_text(badge) if badge else None
        h2 = inner.find("h2")
        if h2 is not None:
            for child in h2.children:
                if isinstance(child, NavigableString):
                    t = re.sub(r"\s+", " ", str(child))
                    if t.strip():
                        s.editor_h2_runs.append(TextRun(text=t, bold=True))
                elif isinstance(child, Tag):
                    if child.name.lower() == "em":
                        s.editor_h2_runs.append(
                            TextRun(text=_plain_text(child), bold=True, color=GOLD))
                    else:
                        s.editor_h2_runs.append(TextRun(text=_plain_text(child), bold=True))
            # Trim leading space on first run / trailing space on last run.
            if s.editor_h2_runs:
                s.editor_h2_runs[0].text = s.editor_h2_runs[0].text.lstrip()
                s.editor_h2_runs[-1].text = s.editor_h2_runs[-1].text.rstrip()
        what = inner.find("p", class_="what")
        if what is not None:
            s.editor_what_runs = _runs_from_tag(what)
        timer = inner.find("div", class_="timer")
        s.editor_timer = _plain_text(timer) if timer else None
    return s


def _parse_content(tag: Tag, kind: str, notes: str) -> Slide:
    s = Slide(kind=kind, speaker_notes=notes)
    head = tag.find("div", class_="slide__head")
    if head is not None:
        eb = head.find("div", class_="slide__eyebrow")
        s.eyebrow = _split_eyebrow(_plain_text(eb)) if eb else None
        mn = head.find("div", class_="slide__minutes")
        s.minutes = _plain_text(mn) if mn else None
    title = tag.find("h2", class_="slide__title")
    if title is not None:
        s.title = _plain_text(title)
    sub = tag.find("p", class_="slide__subtitle")
    if sub is not None:
        s.subtitle = _plain_text(sub)
    body = tag.find("div", class_="slide__body")
    if body is not None:
        s.body_blocks = _extract_body_blocks(body)
        # Some slides (e.g. 0.6 Target) put a trailing <p class="slide__subtitle"> after
        # stats; pick that up too.
        trailing_subs = body.find_all("p", class_="slide__subtitle", recursive=False)
        for t_sub in trailing_subs:
            runs = _runs_from_tag(t_sub)
            if runs:
                # Avoid duplicating block we already captured via _extract_body_blocks
                # (body walker treats <p> as a generic para — this is a guard).
                pass
    s.foot_left, s.foot_right = _foot_parts(tag)
    return s


# ---------------------------------------------------------------------------
# Rendering — low-level helpers
# ---------------------------------------------------------------------------

def _set_solid_bg(slide, color: RGBColor):
    """Paint the entire slide background."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    # Send to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def _add_rect(slide, x, y, w, h, fill: Optional[RGBColor] = None,
              line: Optional[RGBColor] = None, line_w: float = 0.0):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w) if line_w else Pt(0.75)
    shape.shadow.inherit = False
    return shape


def _add_rounded(slide, x, y, w, h, fill: Optional[RGBColor] = None,
                 line: Optional[RGBColor] = None, line_w: float = 0.0,
                 corner: float = 0.08):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    # Tweak corner radius via adjustment.
    try:
        shape.adjustments[0] = corner
    except Exception:
        pass
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w) if line_w else Pt(0.75)
    shape.shadow.inherit = False
    return shape


# --- Pill width auto-fit ----------------------------------------------------
# Pill labels were historically sized by character count (e.g.
# Inches(0.30 + 0.10 * len(text))), which left long labels with trailing
# whitespace and made short labels look cramped. We now (a) measure the
# label's actual rendered advance width using a real TrueType font and (b)
# tag the pill's text frame with PowerPoint's spAutoFit so PowerPoint snaps
# the shape exactly to its text content on render. The pre-computed width is
# the best-guess starting size for renderers that ignore spAutoFit
# (LibreOffice, Google Slides, Keynote thumbnails, etc.).

# Calibri isn't installed on most build hosts (incl. this Linux container),
# so we measure with DejaVu Sans Bold and apply a Calibri-vs-DejaVu width
# correction factor. DejaVu Bold is consistently ~16% wider than Calibri
# Bold across the ASCII range we use for pill labels.
_PILL_FONT_PATH = "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"
_PILL_FONT_CORRECTION = 0.86

_PILL_LEFT_MARGIN_IN = 0.12
_PILL_RIGHT_MARGIN_IN = 0.12
# Small safety buffer so the user's installed Calibri (which can vary
# slightly from our measured metric) never clips the label by a hair.
_PILL_TEXT_BUFFER_IN = 0.06

try:  # Pillow is declared in pyproject.toml; the guard covers stripped envs.
    from PIL import ImageFont as _ImageFont
    _PILL_FONT_CACHE: dict = {}
    _HAVE_PIL_FONT = True
except Exception:
    _HAVE_PIL_FONT = False
    _PILL_FONT_CACHE = {}


def _pill_text_width_inches(text: str, *, size: int) -> float:
    """Return the rendered advance width of ``text`` in Calibri Bold at ``size`` pt.

    Falls back to a per-pt-size character estimate when Pillow or the font
    file is unavailable.
    """
    if _HAVE_PIL_FONT:
        try:
            font = _PILL_FONT_CACHE.get(size)
            if font is None:
                font = _ImageFont.truetype(_PILL_FONT_PATH, size)
                _PILL_FONT_CACHE[size] = font
            # PIL's truetype font returns advance width in pixels at 72 DPI,
            # which means px == pt for our purposes (1 pt = 1/72 inch).
            px = font.getlength(text)
            return (px / 72.0) * _PILL_FONT_CORRECTION
        except Exception:
            pass
    # Fallback: ~0.0078 inches per char per pt (rough Calibri Bold avg).
    return size * 0.0078 * max(1, len(text))


def _pill_width(text: str, *, size: int) -> Inches:
    """Snug width for a pill containing ``text`` at the given font ``size``."""
    text_in = _pill_text_width_inches(text, size=size)
    return Inches(text_in
                  + _PILL_LEFT_MARGIN_IN
                  + _PILL_RIGHT_MARGIN_IN
                  + _PILL_TEXT_BUFFER_IN)


def _add_pill(slide, x, y, w, h, text: str, *, fill: Optional[RGBColor],
              text_color: RGBColor, border: Optional[RGBColor] = None,
              size: int = 9):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        shape.adjustments[0] = 0.5
    except Exception:
        pass
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if border is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = border
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.margin_left = Inches(_PILL_LEFT_MARGIN_IN)
    tf.margin_right = Inches(_PILL_RIGHT_MARGIN_IN)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = False
    # Tell PowerPoint to snap the shape's width/height to the rendered text.
    # Combined with word_wrap=False this resizes width tightly around the
    # label, eliminating trailing whitespace and avoiding clipping when our
    # pre-computed width differs from PowerPoint's actual Calibri metric.
    try:
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT_SANS
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = text_color
    return shape


def _add_textbox(slide, x, y, w, h):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    return box, tf


def _set_para(p, runs: Sequence[TextRun], *, size: int, color: RGBColor,
              align: Optional[str] = None, line_spacing: float = 1.15,
              space_after: int = 0, font: str = FONT_SANS,
              bullet: bool = False, indent_level: int = 0):
    p.line_spacing = line_spacing
    if space_after:
        p.space_after = Pt(space_after)
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    elif align == "left":
        p.alignment = PP_ALIGN.LEFT
    if bullet:
        _enable_bullet(p, indent_level)
    for r in runs:
        run = p.add_run()
        run.text = r.text
        run.font.name = FONT_MONO if r.mono else font
        run.font.size = Pt(size)
        run.font.bold = r.bold
        run.font.italic = r.italic
        run.font.color.rgb = r.color or color


def _enable_bullet(p, level: int = 0):
    """Apply a real PowerPoint round-bullet to the paragraph."""
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    pPr.set("indent", "-228600")  # ~0.25"
    pPr.set("marL", str(228600 * (level + 1)))
    pPr.set("lvl", str(level))
    # Remove any existing bullet child
    for tag in ("buChar", "buAutoNum", "buNone"):
        existing = pPr.find(qn(f"a:{tag}"))
        if existing is not None:
            pPr.remove(existing)
    buFont = pPr.find(qn("a:buFont"))
    if buFont is None:
        buFont = etree.SubElement(pPr, qn("a:buFont"))
    buFont.set("typeface", "Arial")
    buChar = etree.SubElement(pPr, qn("a:buChar"))
    buChar.set("char", "•")


def _enable_numbered(p, level: int = 0):
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    pPr.set("indent", "-228600")
    pPr.set("marL", str(228600 * (level + 1)))
    pPr.set("lvl", str(level))
    for tag in ("buChar", "buAutoNum", "buNone"):
        existing = pPr.find(qn(f"a:{tag}"))
        if existing is not None:
            pPr.remove(existing)
    buAutoNum = etree.SubElement(pPr, qn("a:buAutoNum"))
    buAutoNum.set("type", "arabicPeriod")


def _add_hairline(slide, x, y, w, color: RGBColor, weight: float = 1.0):
    """A thin horizontal rule."""
    shape = slide.shapes.add_connector(1, x, y, x + w, y)
    shape.line.color.rgb = color
    shape.line.width = Pt(weight)
    return shape


# ---------------------------------------------------------------------------
# Per-slide layout builders
# ---------------------------------------------------------------------------

def _draw_top_brand_bar(slide):
    """Scarlet+gold brand bar at the very top of every standard slide."""
    bar_h = Inches(0.07)
    scarlet_w = SLIDE_W * 0.7
    gold_w = SLIDE_W - scarlet_w
    _add_rect(slide, 0, 0, scarlet_w, bar_h, fill=SCARLET)
    _add_rect(slide, scarlet_w, 0, gold_w, bar_h, fill=GOLD)


def _draw_head(slide, eyebrow: Optional[str], minutes: Optional[str]):
    """Top eyebrow (red, uppercase) + minutes (gray, uppercase) + bottom rule."""
    margin_x = Inches(0.6)
    y = Inches(0.30)
    if eyebrow:
        box, tf = _add_textbox(slide, margin_x, y, SLIDE_W - 2 * margin_x - Inches(2.5), Inches(0.35))
        tf.margin_left = 0
        tf.word_wrap = False
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = eyebrow.upper()
        r.font.name = FONT_SANS
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = SCARLET
        # Letter spacing isn't directly supported in python-pptx; we rely on caps + bold.
    if minutes:
        box, tf = _add_textbox(slide, SLIDE_W - margin_x - Inches(2.5), y,
                               Inches(2.5), Inches(0.35))
        tf.margin_right = 0
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = minutes.upper()
        r.font.name = FONT_SANS
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = INK_3
    _add_hairline(slide, margin_x, Inches(0.7), SLIDE_W - 2 * margin_x, LINE_2, weight=1.0)


def _draw_title(slide, title: Optional[str], subtitle: Optional[str]):
    margin_x = Inches(0.6)
    y = Inches(0.85)
    if title:
        box, tf = _add_textbox(slide, margin_x, y, SLIDE_W - 2 * margin_x, Inches(0.85))
        p = tf.paragraphs[0]
        p.line_spacing = 1.05
        r = p.add_run()
        r.text = title
        r.font.name = FONT_SANS
        r.font.size = Pt(34)
        r.font.bold = True
        r.font.color.rgb = INK
    if subtitle:
        box, tf = _add_textbox(slide, margin_x, Inches(1.74),
                               SLIDE_W - 2 * margin_x, Inches(0.6))
        p = tf.paragraphs[0]
        p.line_spacing = 1.25
        r = p.add_run()
        r.text = subtitle
        r.font.name = FONT_SANS
        r.font.size = Pt(15)
        r.font.bold = False
        r.font.color.rgb = INK_2


def _draw_foot(slide, left: Optional[str], right: Optional[str]):
    if not (left or right):
        return
    margin_x = Inches(0.6)
    y = Inches(7.05)
    _add_hairline(slide, margin_x, y, SLIDE_W - 2 * margin_x, LINE_2, weight=0.75)
    fy = Inches(7.13)
    if left:
        box, tf = _add_textbox(slide, margin_x, fy, Inches(7), Inches(0.3))
        tf.word_wrap = False
        p = tf.paragraphs[0]
        # Tiny scarlet brand-mark dot
        r0 = p.add_run()
        r0.text = "▍ "
        r0.font.name = FONT_SANS
        r0.font.size = Pt(9)
        r0.font.color.rgb = SCARLET
        r0.font.bold = True
        r = p.add_run()
        r.text = left.upper()
        r.font.name = FONT_SANS
        r.font.size = Pt(9)
        r.font.color.rgb = INK_3
        r.font.bold = False
    if right:
        box, tf = _add_textbox(slide, SLIDE_W - margin_x - Inches(5), fy,
                               Inches(5), Inches(0.3))
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = right.upper()
        r.font.name = FONT_SANS
        r.font.size = Pt(9)
        r.font.color.rgb = INK_3


# --- Body-region renderers ----------------------------------------------------

BODY_X = Inches(0.6)
BODY_Y = Inches(2.45)
BODY_W = SLIDE_W - Inches(1.2)
BODY_H = Inches(4.45)


def _render_bullets(slide, x, y, w, h, items: List[List[TextRun]],
                    *, size: int = 14, color: RGBColor = INK,
                    line_spacing: float = 1.30, space_after: int = 6):
    box, tf = _add_textbox(slide, x, y, w, h)
    tf.word_wrap = True
    first = True
    for runs in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        _set_para(p, runs, size=size, color=color,
                  line_spacing=line_spacing, space_after=space_after,
                  bullet=True)
        # Make bullet glyph scarlet via the bullet color attribute.
        pPr = p._pPr
        if pPr is not None:
            buClr = pPr.find(qn("a:buClr"))
            if buClr is not None:
                pPr.remove(buClr)
            buClr = etree.SubElement(pPr, qn("a:buClr"))
            srgb = etree.SubElement(buClr, qn("a:srgbClr"))
            srgb.set("val", "CC0000")
            # reorder buClr before buFont
            pPr.remove(buClr)
            pPr.insert(0, buClr)
    return box


def _render_checklist(slide, x, y, w, h, items: List[List[TextRun]],
                      *, size: int = 13):
    """Native bulleted list with the success-green check marker."""
    box, tf = _add_textbox(slide, x, y, w, h)
    tf.word_wrap = True
    first = True
    for runs in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        # Prepend a leading checkmark glyph as an inline run.
        check_run = TextRun(text="✓ ", bold=True, color=SUCCESS)
        _set_para(p, [check_run] + list(runs), size=size, color=INK,
                  line_spacing=1.30, space_after=5)
    return box


def _render_paragraph(slide, x, y, w, h, runs: List[TextRun],
                      *, size: int = 14, color: RGBColor = INK_2,
                      align: Optional[str] = None, italic: bool = False):
    box, tf = _add_textbox(slide, x, y, w, h)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    _set_para(p, runs, size=size, color=color, align=align, line_spacing=1.30)


def _render_card(slide, x, y, w, h, card: dict):
    variant = card.get("variant", "default")
    if variant == "ink":
        bg = INK
        border = INK
        title_color = GOLD
        body_color = RGBColor(0xD8, 0xD8, 0xD8)
    elif variant == "accent":
        bg = PAPER
        border = LINE
        title_color = INK
        body_color = INK_2
    elif variant == "gold":
        bg = PAPER
        border = LINE
        title_color = INK
        body_color = INK_2
    else:
        bg = PAPER_2
        border = LINE
        title_color = INK
        body_color = INK_2
    _add_rounded(slide, x, y, w, h, fill=bg, line=border, line_w=0.75, corner=0.06)
    # Left accent stripe for accent / gold variants.
    if variant in ("accent", "gold"):
        stripe_color = SCARLET if variant == "accent" else GOLD
        _add_rect(slide, x, y, Inches(0.07), h, fill=stripe_color)
    pad = Inches(0.22)
    inner_x = x + pad + (Inches(0.08) if variant in ("accent", "gold") else Emu(0))
    inner_y = y + pad
    inner_w = w - pad * 2 - (Inches(0.08) if variant in ("accent", "gold") else Emu(0))
    cur_y = inner_y
    if card.get("title"):
        title_box, tf = _add_textbox(slide, inner_x, cur_y, inner_w, Inches(0.4))
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = card["title"]
        r.font.name = FONT_SANS
        r.font.size = Pt(15)
        r.font.bold = True
        r.font.color.rgb = title_color
        cur_y += Inches(0.42)
    # Body content
    body_box, body_tf = _add_textbox(slide, inner_x, cur_y, inner_w, h - (cur_y - y) - pad)
    body_tf.word_wrap = True
    first = True
    for blk in card.get("content", []):
        if blk["type"] == "p":
            if first:
                p = body_tf.paragraphs[0]
                first = False
            else:
                p = body_tf.add_paragraph()
            _set_para(p, blk["runs"], size=12, color=body_color,
                      line_spacing=1.30, space_after=6)
        elif blk["type"] == "list":
            for items_runs in blk.get("items", []):
                if first:
                    p = body_tf.paragraphs[0]
                    first = False
                else:
                    p = body_tf.add_paragraph()
                if blk.get("ordered"):
                    _set_para(p, items_runs, size=12, color=body_color,
                              line_spacing=1.25, space_after=4)
                    _enable_numbered(p)
                else:
                    _set_para(p, items_runs, size=12, color=body_color,
                              line_spacing=1.25, space_after=4, bullet=True)
                    pPr = p._pPr
                    if pPr is not None:
                        buClr = pPr.find(qn("a:buClr"))
                        if buClr is not None:
                            pPr.remove(buClr)
                        buClr = etree.SubElement(pPr, qn("a:buClr"))
                        srgb = etree.SubElement(buClr, qn("a:srgbClr"))
                        srgb.set("val", "CC0000" if variant != "ink" else "F5D130")
                        pPr.remove(buClr)
                        pPr.insert(0, buClr)


def _render_prompt(slide, x, y, w, h, label: str, runs: List[TextRun]):
    """Info-style call-out: light blue background, dark blue accent stripe."""
    _add_rounded(slide, x, y, w, h, fill=INFO_BG, line=INFO_BG, corner=0.04)
    _add_rect(slide, x, y, Inches(0.08), h, fill=INFO)
    pad = Inches(0.22)
    inner_x = x + pad + Inches(0.08)
    inner_y = y + pad
    inner_w = w - pad * 2 - Inches(0.08)
    if label:
        box, tf = _add_textbox(slide, inner_x, inner_y, inner_w, Inches(0.3))
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = label.upper()
        r.font.name = FONT_SANS
        r.font.size = Pt(9)
        r.font.bold = True
        r.font.color.rgb = INFO
        inner_y += Inches(0.32)
    body_box, body_tf = _add_textbox(slide, inner_x, inner_y, inner_w,
                                     h - (inner_y - y) - pad)
    body_tf.word_wrap = True
    p = body_tf.paragraphs[0]
    _set_para(p, runs, size=12, color=INFO, line_spacing=1.35)


def _render_pillrow(slide, x, y, w, h, pills: List[Tuple[str, str]]):
    cur_x = x
    cur_y = y
    pill_h = Inches(0.28)
    gap_x = Inches(0.08)
    gap_y = Inches(0.10)
    for text, kind in pills:
        # Snug width that fits the actual rendered label (PowerPoint also
        # snaps via spAutoFit on _add_pill's text frame).
        pill_w = _pill_width(text, size=9)
        if cur_x + pill_w > x + w:
            cur_x = x
            cur_y += pill_h + gap_y
        if cur_y + pill_h > y + h:
            break
        if kind == "scarlet":
            _add_pill(slide, cur_x, cur_y, pill_w, pill_h, text,
                      fill=SCARLET, text_color=PAPER, border=SCARLET, size=9)
        elif kind == "gold":
            _add_pill(slide, cur_x, cur_y, pill_w, pill_h, text,
                      fill=GOLD, text_color=INK, border=GOLD, size=9)
        elif kind == "ink":
            _add_pill(slide, cur_x, cur_y, pill_w, pill_h, text,
                      fill=INK, text_color=PAPER, border=INK, size=9)
        else:
            _add_pill(slide, cur_x, cur_y, pill_w, pill_h, text,
                      fill=PAPER_2, text_color=INK_2, border=LINE, size=9)
        cur_x += pill_w + gap_x


def _render_table_compare(slide, x, y, w, h, spec: TableSpec):
    headers = spec.headers
    rows = spec.rows
    n_cols = max(len(headers), max((len(r) for r in rows), default=1))
    n_rows = (1 if headers else 0) + len(rows)
    if n_rows == 0 or n_cols == 0:
        return
    table_shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h).table
    # Column widths: first column ~36%, remaining split evenly.
    if n_cols >= 3:
        first_w = int(w * 0.30)
        rest_w = (w - first_w) // (n_cols - 1)
        table_shape.columns[0].width = first_w
        for i in range(1, n_cols):
            table_shape.columns[i].width = rest_w
    else:
        col_w = w // n_cols
        for i in range(n_cols):
            table_shape.columns[i].width = col_w
    # Header row
    if headers:
        for ci, h_text in enumerate(headers):
            cell = table_shape.cell(0, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = INK
            tf = cell.text_frame
            tf.margin_left = Inches(0.1)
            tf.margin_right = Inches(0.1)
            tf.margin_top = Inches(0.05)
            tf.margin_bottom = Inches(0.05)
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = h_text.upper()
            r.font.name = FONT_SANS
            r.font.size = Pt(10)
            r.font.bold = True
            r.font.color.rgb = PAPER
    # Body rows
    for ri, row in enumerate(rows):
        target_ri = ri + (1 if headers else 0)
        for ci, (txt, cls) in enumerate(row):
            cell = table_shape.cell(target_ri, ci)
            # Alternate background
            cell.fill.solid()
            cell.fill.fore_color.rgb = PAPER if ri % 2 == 0 else PAPER_2
            tf = cell.text_frame
            tf.margin_left = Inches(0.1)
            tf.margin_right = Inches(0.1)
            tf.margin_top = Inches(0.05)
            tf.margin_bottom = Inches(0.05)
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = txt
            r.font.name = FONT_SANS
            r.font.size = Pt(11)
            if cls and "win" in cls:
                r.font.color.rgb = SUCCESS
                r.font.bold = True
            elif cls and "dim" in cls:
                r.font.color.rgb = INK_3
            else:
                r.font.color.rgb = INK


def _render_arch(slide, x, y, w, h, layers: List[ArchLayer]):
    """Four-layer architecture diagram. Built layers fill scarlet on the
    label column and bold the name; current layer adds a gold left stripe."""
    n = len(layers)
    if n == 0:
        return
    row_h = (h - Inches(0.1) * (n - 1)) / n
    label_w = Inches(1.4)
    name_w = (w - label_w) * 0.42
    desc_w = w - label_w - name_w
    for i, layer in enumerate(layers):
        ry = y + (row_h + Inches(0.10)) * i
        # Gold stripe on the very left for current
        stripe_w = Inches(0.07) if layer.is_current else Emu(0)
        if layer.is_current:
            _add_rect(slide, x, ry, Inches(0.07), row_h, fill=GOLD)
        # Label cell
        label_fill = SCARLET if layer.is_built else PAPER_2
        label_text_color = PAPER if layer.is_built else INK_2
        _add_rect(slide, x + stripe_w, ry, label_w - stripe_w, row_h,
                  fill=label_fill, line=LINE, line_w=0.5)
        box, tf = _add_textbox(slide, x + stripe_w + Inches(0.1), ry,
                               label_w - stripe_w - Inches(0.1), row_h)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = layer.no
        r.font.name = FONT_SANS
        r.font.size = Pt(9)
        r.font.bold = True
        r.font.color.rgb = label_text_color
        # Name cell
        name_fill = PAPER if layer.is_built else PAPER_2
        _add_rect(slide, x + label_w, ry, name_w, row_h,
                  fill=name_fill, line=LINE, line_w=0.5)
        box, tf = _add_textbox(slide, x + label_w + Inches(0.15), ry,
                               name_w - Inches(0.2), row_h)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = layer.name
        r.font.name = FONT_SANS
        r.font.size = Pt(13)
        r.font.bold = layer.is_built
        r.font.color.rgb = INK if layer.is_built else INK_2
        if layer.is_built:
            r2 = p.add_run()
            r2.text = "  ✓"
            r2.font.name = FONT_SANS
            r2.font.size = Pt(11)
            r2.font.color.rgb = SCARLET
            r2.font.bold = True
        # Desc cell
        _add_rect(slide, x + label_w + name_w, ry, desc_w, row_h,
                  fill=PAPER_2, line=LINE, line_w=0.5)
        box, tf = _add_textbox(slide, x + label_w + name_w + Inches(0.15), ry,
                               desc_w - Inches(0.2), row_h)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p = tf.paragraphs[0]
        _set_para(p, layer.desc_runs or [TextRun(text="")], size=10, color=INK_2,
                  line_spacing=1.20)


def _render_arch_callout(slide, x, y, w, runs: List[TextRun]):
    h = Inches(0.6)
    _add_rounded(slide, x, y, w, h, fill=PAPER_3, line=LINE, line_w=0.5, corner=0.06)
    box, tf = _add_textbox(slide, x + Inches(0.18), y + Inches(0.05),
                           w - Inches(0.36), h - Inches(0.10))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    p = tf.paragraphs[0]
    _set_para(p, runs, size=11, color=INK_2, line_spacing=1.25)


def _render_stats(slide, x, y, w, h, items: List[dict]):
    n = len(items)
    if n == 0:
        return
    gap = Inches(0.15)
    cell_w = (w - gap * (n - 1)) / n
    for i, item in enumerate(items):
        cx = x + (cell_w + gap) * i
        _add_rounded(slide, cx, y, cell_w, h, fill=PAPER_2, line=LINE, line_w=0.5, corner=0.04)
        # Number
        box, tf = _add_textbox(slide, cx + Inches(0.18), y + Inches(0.18),
                               cell_w - Inches(0.36), Inches(1.1))
        tf.word_wrap = False
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = item.get("num", "")
        r.font.name = FONT_SANS
        r.font.size = Pt(40)
        r.font.bold = True
        r.font.color.rgb = SCARLET
        # Label
        box, tf = _add_textbox(slide, cx + Inches(0.18), y + Inches(1.40),
                               cell_w - Inches(0.36), Inches(0.30))
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = item.get("lab", "").upper()
        r.font.name = FONT_SANS
        r.font.size = Pt(9)
        r.font.bold = True
        r.font.color.rgb = INK_3
        # Description
        box, tf = _add_textbox(slide, cx + Inches(0.18), y + Inches(1.70),
                               cell_w - Inches(0.36), h - Inches(1.85))
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = item.get("desc", "")
        r.font.name = FONT_SANS
        r.font.size = Pt(10)
        r.font.color.rgb = INK_2


def _render_agenda(slide, x, y, w, h, items: List[Tuple[str, List[TextRun]]]):
    if not items:
        return
    n = len(items)
    row_h = h / n
    time_w = Inches(0.9)
    for i, (time, name_runs) in enumerate(items):
        ry = y + row_h * i
        # Time cell
        _add_rect(slide, x, ry, time_w, row_h - Inches(0.04),
                  fill=PAPER_2, line=LINE, line_w=0.4)
        box, tf = _add_textbox(slide, x + Inches(0.1), ry,
                               time_w - Inches(0.2), row_h - Inches(0.04))
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = time
        r.font.name = FONT_SANS
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = SCARLET
        # Name cell
        box, tf = _add_textbox(slide, x + time_w + Inches(0.18), ry,
                               w - time_w - Inches(0.2), row_h - Inches(0.04))
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p = tf.paragraphs[0]
        _set_para(p, name_runs, size=12, color=INK, line_spacing=1.20)


def _render_pullquote(slide, x, y, w, h, quote: str, attr: str):
    _add_rounded(slide, x, y, w, h, fill=PAPER_2, line=GOLD, line_w=1.5, corner=0.04)
    box, tf = _add_textbox(slide, x + Inches(0.4), y + Inches(0.25),
                           w - Inches(0.8), h - Inches(0.5))
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "“"
    r.font.name = FONT_SANS
    r.font.size = Pt(36)
    r.font.bold = True
    r.font.color.rgb = GOLD
    p2 = tf.add_paragraph()
    p2.line_spacing = 1.30
    r = p2.add_run()
    r.text = quote
    r.font.name = FONT_SANS
    r.font.size = Pt(16)
    r.font.italic = True
    r.font.color.rgb = INK
    if attr:
        p3 = tf.add_paragraph()
        p3.space_before = Pt(8)
        r = p3.add_run()
        r.text = attr
        r.font.name = FONT_SANS
        r.font.size = Pt(10)
        r.font.color.rgb = INK_3


def _render_build_banner(slide, x, y, text: str):
    pill_w = _pill_width(text, size=9)
    _add_pill(slide, x, y, pill_w, Inches(0.28), text,
              fill=SCARLET, text_color=PAPER, border=SCARLET, size=9)


# --- Block dispatcher --------------------------------------------------------

def _render_block(slide, x, y, w, h, block: dict) -> Inches:
    """Render a single block within a region; return consumed height."""
    btype = block.get("type")
    if btype == "bullets":
        _render_bullets(slide, x, y, w, h, block["items"], size=14)
        return h
    if btype == "checklist":
        _render_checklist(slide, x, y, w, h, block["items"], size=13)
        return h
    if btype == "list":
        # Render plain list as bullets too.
        _render_bullets(slide, x, y, w, h, block["items"], size=13)
        return h
    if btype == "para":
        _render_paragraph(slide, x, y, w, h, block["runs"], size=14, color=INK_2)
        return h
    if btype == "card":
        _render_card(slide, x, y, w, h, block)
        return h
    if btype == "prompt":
        _render_prompt(slide, x, y, w, h, block.get("label", ""), block.get("runs", []))
        return h
    if btype == "pillrow":
        _render_pillrow(slide, x, y, w, h, block["pills"])
        return h
    if btype == "build_banner":
        _render_build_banner(slide, x, y, block["text"])
        return Inches(0.4)
    if btype == "arch_callout":
        _render_arch_callout(slide, x, y, w, block["runs"])
        return Inches(0.6)
    if btype == "agenda":
        _render_agenda(slide, x, y, w, h, block["items"])
        return h
    if btype == "stats":
        _render_stats(slide, x, y, w, h, block["items"])
        return h
    if btype == "pullquote":
        _render_pullquote(slide, x, y, w, h, block["quote"], block["attr"])
        return h
    if btype == "table":
        _render_table_compare(slide, x, y, w, h, block["spec"])
        return h
    if btype == "arch":
        _render_arch(slide, x, y, w, h, block["layers"])
        return h
    if btype == "cols2":
        _render_cols2(slide, x, y, w, h, block)
        return h
    return h


def _render_cols2(slide, x, y, w, h, block: dict):
    kind = block.get("kind", "5050")
    cols = block.get("cols", [])
    gap = Inches(0.30)
    if kind == "6040":
        c1_w = (w - gap) * 0.60
    elif kind == "4060":
        c1_w = (w - gap) * 0.40
    else:
        c1_w = (w - gap) * 0.50
    c2_w = w - gap - c1_w
    _render_col_blocks(slide, x, y, c1_w, h, cols[0] if len(cols) > 0 else [])
    _render_col_blocks(slide, x + c1_w + gap, y, c2_w, h, cols[1] if len(cols) > 1 else [])


def _est_text_lines(text: str, w_in: float, *, size_pt: int) -> int:
    """Word-aware estimate of how many visual lines `text` wraps to at `w_in`.

    Uses a fixed average character width tuned for Calibri Regular. Bold text
    is ~10% wider; we err on the conservative side so we don't under-estimate.
    """
    if not text or w_in <= 0:
        return 0
    avg_char_in = size_pt * 0.0070  # ~0.084" per char at 12pt Calibri
    chars_per_line = max(8, int(w_in / avg_char_in))
    lines = 0
    for paragraph in text.split("\n"):
        words = paragraph.split()
        if not words:
            lines += 1
            continue
        cur = 0
        line_count = 1
        for word in words:
            add = len(word) + (1 if cur > 0 else 0)
            if cur + add > chars_per_line and cur > 0:
                line_count += 1
                cur = len(word)
            else:
                cur += add
        lines += line_count
    return lines


def _natural_prompt_height_in(block: dict, w_in: float) -> float:
    """Estimated natural rendered height (inches) of a prompt callout."""
    pad = 0.22
    label_h = 0.32 if block.get("label") else 0.0
    inner_w = w_in - pad * 2 - 0.08
    text = "".join(r.text for r in block.get("runs", []))
    line_h = (12 / 72.0) * 1.35
    lines = _est_text_lines(text, inner_w, size_pt=12)
    text_h = lines * line_h
    # Add a small safety buffer so descenders / wider glyphs don't clip.
    return pad * 2 + label_h + text_h + 0.10


def _natural_card_height_in(block: dict, w_in: float) -> float:
    """Estimated natural rendered height (inches) of a card."""
    pad = 0.22
    has_stripe = block.get("variant") in ("accent", "gold")
    inner_w = w_in - pad * 2 - (0.08 if has_stripe else 0.0)
    h = pad * 2
    if block.get("title"):
        h += 0.42
    p_line_h = (12 / 72.0) * 1.30
    li_line_h = (12 / 72.0) * 1.25
    for blk in block.get("content", []):
        bt = blk.get("type")
        if bt == "p":
            text = "".join(r.text for r in blk.get("runs", []))
            lines = max(1, _est_text_lines(text, inner_w, size_pt=12))
            h += lines * p_line_h + (6 / 72.0)
        elif bt == "list":
            for items_runs in blk.get("items", []):
                text = "".join(r.text for r in items_runs)
                # Account for bullet indent of ~0.25"
                lines = max(1, _est_text_lines(text, inner_w - 0.25, size_pt=12))
                h += lines * li_line_h + (4 / 72.0)
    return h + 0.08


def _natural_block_height_in(block: dict, w_in: float) -> Optional[float]:
    """Return natural content height in inches for blocks we know how to size."""
    bt = block.get("type")
    if bt == "prompt":
        return _natural_prompt_height_in(block, w_in)
    if bt == "card":
        return _natural_card_height_in(block, w_in)
    return None


def _render_col_blocks(slide, x, y, w, h, blocks: List[dict]):
    """Render a vertical stack of blocks inside a column region."""
    if not blocks:
        return
    # Cap single-block columns to their natural content height so a short
    # prompt or card doesn't stretch to the whole column. This keeps the
    # build-framing slides (e.g. "Five tools, then we go") from leaving
    # several inches of empty space below a small callout.
    if len(blocks) == 1:
        natural = _natural_block_height_in(blocks[0], w / 914400)
        if natural is not None:
            natural_emu = Inches(natural)
            if natural_emu < h:
                h = natural_emu
    # Heuristic per-block height allocation, then fill the column.
    weights: List[float] = []
    for b in blocks:
        bt = b.get("type")
        if bt == "build_banner":
            weights.append(0.4)
        elif bt == "arch_callout":
            weights.append(0.6)
        elif bt == "para":
            n_chars = sum(len(r.text) for r in b.get("runs", []))
            weights.append(max(0.6, min(1.5, n_chars / 120)))
        elif bt == "card":
            n_items = sum(
                len(c.get("items", [])) for c in b.get("content", []) if c.get("type") == "list"
            ) + sum(1 for c in b.get("content", []) if c.get("type") == "p")
            weights.append(max(2.0, 1.0 + n_items * 0.6))
        elif bt == "prompt":
            n_chars = sum(len(r.text) for r in b.get("runs", []))
            weights.append(max(2.0, 1.0 + n_chars / 80))
        elif bt in ("bullets", "checklist", "list"):
            weights.append(max(2.0, 0.7 + len(b.get("items", [])) * 0.55))
        elif bt == "agenda":
            weights.append(max(3.0, len(b.get("items", [])) * 0.45))
        elif bt == "stats":
            weights.append(2.5)
        elif bt == "pullquote":
            weights.append(2.0)
        elif bt == "table":
            weights.append(3.0)
        elif bt == "arch":
            weights.append(3.0)
        elif bt == "pillrow":
            # Pillrows are visually short — give them a small weight so they
            # don't gobble half a stacked column.
            weights.append(0.5)
        else:
            weights.append(1.0)
    # Reserve fixed heights for compact blocks first, then distribute the
    # remaining height across the rest by weight.
    gap = Inches(0.18)
    fixed_h: List[Optional[int]] = []
    for b in blocks:
        bt = b.get("type")
        if bt == "build_banner":
            fixed_h.append(Inches(0.4))
        elif bt == "arch_callout":
            fixed_h.append(Inches(0.65))
        elif bt == "pillrow":
            n_pills = len(b.get("pills", []))
            rows = max(1, (n_pills + 3) // 4)
            fixed_h.append(Inches(0.30 + (rows - 1) * 0.38 + 0.1))
        else:
            fixed_h.append(None)
    flexible_weight = sum(
        wt for wt, fh in zip(weights, fixed_h) if fh is None
    ) or 1.0
    flexible_h = h - gap * (len(blocks) - 1) - sum(fh for fh in fixed_h if fh is not None)
    if flexible_h < 0:
        flexible_h = 0
    cur_y = y
    for b, weight, fh in zip(blocks, weights, fixed_h):
        if fh is not None:
            bh = fh
        else:
            bh = int(flexible_h * (weight / flexible_weight))
        _render_block(slide, x, cur_y, w, bh, b)
        cur_y += bh + gap


# --- Top-level slide builders ------------------------------------------------

def build_cover_slide(prs, slide: Slide):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_solid_bg(s, COVER_WASH)
    # Subtle scarlet wash bottom-left
    wash = _add_rect(s, 0, Inches(4.0), Inches(7), Inches(3.5),
                     fill=RGBColor(0xFC, 0xE8, 0xE3))
    # Top bar (capstone gold strip — full width gold under scarlet)
    _draw_top_brand_bar(s)

    # Top: brand + capstone stamp
    margin_x = Inches(0.7)
    box, tf = _add_textbox(s, margin_x, Inches(0.55), Inches(8), Inches(0.4))
    p = tf.paragraphs[0]
    # Scarlet vertical mark
    r0 = p.add_run()
    r0.text = "▍ "
    r0.font.name = FONT_SANS
    r0.font.size = Pt(13)
    r0.font.color.rgb = SCARLET
    r0.font.bold = True
    r = p.add_run()
    r.text = (slide.cover_brand or "").upper()
    r.font.name = FONT_SANS
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = INK
    if slide.cover_stamp:
        stamp_text = slide.cover_stamp.upper()
        stamp_w = _pill_width(stamp_text, size=10)
        # Right-align the stamp against the same right margin as before.
        stamp_x = SLIDE_W - Inches(0.7) - stamp_w
        _add_pill(s, stamp_x, Inches(0.50), stamp_w, Inches(0.36),
                  stamp_text, fill=PAPER, text_color=SCARLET,
                  border=SCARLET, size=10)
    # Mid: big H1. Auto-shrink and stretch the box to fit longer titles
    # so the gold rule never cuts through a third line. Each h1 entry may
    # wrap to multiple visual lines, so estimate visible-line count from the
    # text width relative to the box width.
    title_y = Inches(1.5)
    box_w_in = (SLIDE_W - 2 * margin_x) / 914400  # inches available for text
    h1_entries = slide.cover_h1_lines or []
    # Choose font size: shrink if total text would clearly wrap to 3+ lines
    # at 64pt. Approx avg char width = 0.45 * pt/72 inches for this sans face.
    def _visual_lines(size_pt: int) -> int:
        # Bold sans (Inter/Calibri-like) at this size renders ≈ 0.62 × pt/72"
        # per average character. Be conservative so we never under-estimate.
        char_w = 0.62 * size_pt / 72.0
        chars_per_line = max(6, int(box_w_in / char_w))
        total = 0
        for text, _ in h1_entries:
            n = max(1, len(text))
            total += (n + chars_per_line - 1) // chars_per_line
        return max(1, total)

    h1_size = 64
    while h1_size > 44 and _visual_lines(h1_size) > 2:
        h1_size -= 4
    visible_lines = _visual_lines(h1_size)
    line_h = (h1_size / 72.0) * 1.25  # 1.25× leading approximation
    title_h = Inches(max(1.6, line_h * visible_lines + 0.20))
    rule_y = title_y + title_h + Inches(0.10)
    if slide.cover_h1_lines:
        box, tf = _add_textbox(s, margin_x, title_y, SLIDE_W - 2 * margin_x, title_h)
        tf.word_wrap = True
        first = True
        for text, accent in slide.cover_h1_lines:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.line_spacing = 1.0
            r = p.add_run()
            r.text = text
            r.font.name = FONT_SANS
            r.font.size = Pt(h1_size)
            r.font.bold = True
            r.font.color.rgb = SCARLET if accent else INK
    # Gradient rule (approximated as a 2-segment scarlet→gold rectangle).
    rule_w = Inches(2.6)
    _add_rect(s, margin_x, rule_y, rule_w * 0.6, Inches(0.10), fill=SCARLET)
    _add_rect(s, margin_x + rule_w * 0.6, rule_y, rule_w * 0.4, Inches(0.10), fill=GOLD)

    lede_y = rule_y + Inches(0.30)
    if slide.cover_lede:
        box, tf = _add_textbox(s, margin_x, lede_y,
                               SLIDE_W - 2 * margin_x - Inches(0.5), Inches(1.2))
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.line_spacing = 1.35
        r = p.add_run()
        r.text = slide.cover_lede
        r.font.name = FONT_SANS
        r.font.size = Pt(15)
        r.font.color.rgb = INK_2
    # Bottom: meta + week pip
    bottom_y = Inches(6.30)
    if slide.cover_meta:
        box, tf = _add_textbox(s, margin_x, bottom_y, Inches(8.5), Inches(0.7))
        tf.word_wrap = False
        for i, (strong_text, after) in enumerate(slide.cover_meta):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            r = p.add_run()
            r.text = f"{strong_text}  "
            r.font.name = FONT_SANS
            r.font.size = Pt(12)
            r.font.bold = True
            r.font.color.rgb = INK
            if after:
                r2 = p.add_run()
                r2.text = after
                r2.font.name = FONT_SANS
                r2.font.size = Pt(10)
                r2.font.color.rgb = INK_3
    if slide.cover_pip_label:
        # Black pill with gold "6" badge
        pip_w = Inches(2.6)
        pip_h = Inches(0.5)
        pip_x = SLIDE_W - margin_x - pip_w
        pip_y = bottom_y + Inches(0.1)
        _add_rounded(s, pip_x, pip_y, pip_w, pip_h, fill=INK, line=INK,
                     corner=0.5)
        # Gold circle
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  pip_x + Inches(0.12), pip_y + Inches(0.07),
                                  Inches(0.36), Inches(0.36))
        circ.fill.solid()
        circ.fill.fore_color.rgb = GOLD
        circ.line.fill.background()
        circ.shadow.inherit = False
        ctf = circ.text_frame
        ctf.margin_left = 0
        ctf.margin_right = 0
        ctf.margin_top = 0
        ctf.margin_bottom = 0
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = slide.cover_pip_num or "6"
        cr.font.name = FONT_SANS
        cr.font.size = Pt(13)
        cr.font.bold = True
        cr.font.color.rgb = INK
        # Label
        box, tf = _add_textbox(s, pip_x + Inches(0.6), pip_y, pip_w - Inches(0.7), pip_h)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = False
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = (slide.cover_pip_label or "").upper()
        r.font.name = FONT_SANS
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = PAPER

    _attach_notes(s, slide.speaker_notes)


def build_section_slide(prs, slide: Slide):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_solid_bg(s, INK)
    # Vertical scarlet rule between halves
    half_x = SLIDE_W / 2
    _add_rect(s, half_x - Inches(0.05), Inches(1.2), Inches(0.05),
              SLIDE_H - Inches(2.4), fill=SCARLET)
    # Left: module eyebrow + big number
    margin = Inches(0.7)
    if slide.section_eyebrow:
        box, tf = _add_textbox(s, margin, Inches(1.5), half_x - margin - Inches(0.3), Inches(0.4))
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = slide.section_eyebrow.upper()
        r.font.name = FONT_SANS
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = GOLD
    if slide.section_no:
        box, tf = _add_textbox(s, margin, Inches(2.0),
                               half_x - margin - Inches(0.3), Inches(4.0))
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = slide.section_no
        r.font.name = FONT_SANS
        r.font.size = Pt(180)
        r.font.bold = True
        r.font.color.rgb = PAPER
        if slide.section_of:
            r2 = p.add_run()
            r2.text = "  " + slide.section_of
            r2.font.name = FONT_SANS
            r2.font.size = Pt(28)
            r2.font.bold = False
            r2.font.color.rgb = INK_3
    # Right: name + meta + blurb
    right_x = half_x + Inches(0.5)
    right_w = SLIDE_W - right_x - margin
    if slide.section_name:
        box, tf = _add_textbox(s, right_x, Inches(1.7), right_w, Inches(2.0))
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.line_spacing = 1.0
        r = p.add_run()
        r.text = slide.section_name
        r.font.name = FONT_SANS
        r.font.size = Pt(40)
        r.font.bold = True
        r.font.color.rgb = PAPER
    # Pills
    pill_y = Inches(4.0)
    cur_x = right_x
    pill_h = Inches(0.36)
    for text, kind in slide.section_pills:
        pw = _pill_width(text, size=10)
        if kind == "build":
            _add_pill(s, cur_x, pill_y, pw, pill_h, text,
                      fill=SCARLET, text_color=PAPER, border=SCARLET, size=10)
        elif kind == "gold":
            _add_pill(s, cur_x, pill_y, pw, pill_h, text,
                      fill=GOLD, text_color=INK, border=GOLD, size=10)
        else:
            # Outline on dark — use white text & light border for legibility.
            _add_pill(s, cur_x, pill_y, pw, pill_h, text,
                      fill=None, text_color=PAPER, border=RGBColor(0x99, 0x99, 0x99), size=10)
        cur_x += pw + Inches(0.15)
    # Blurb
    if slide.section_blurb:
        box, tf = _add_textbox(s, right_x, Inches(4.7), right_w, Inches(2.0))
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.line_spacing = 1.40
        r = p.add_run()
        r.text = slide.section_blurb
        r.font.name = FONT_SANS
        r.font.size = Pt(14)
        r.font.color.rgb = RGBColor(0xC8, 0xC8, 0xC8)

    _attach_notes(s, slide.speaker_notes)


def build_break_slide(prs, slide: Slide):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_solid_bg(s, PAPER_2)
    _draw_top_brand_bar(s)
    if slide.break_lab:
        box, tf = _add_textbox(s, 0, Inches(2.0), SLIDE_W, Inches(0.5))
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = slide.break_lab.upper()
        r.font.name = FONT_SANS
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = SCARLET
    if slide.break_num:
        box, tf = _add_textbox(s, 0, Inches(2.6), SLIDE_W, Inches(3.0))
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 0.95
        r = p.add_run()
        r.text = slide.break_num
        r.font.name = FONT_SANS
        r.font.size = Pt(220)
        r.font.bold = True
        r.font.color.rgb = INK
        if slide.break_num_unit:
            r2 = p.add_run()
            r2.text = " " + slide.break_num_unit
            r2.font.name = FONT_SANS
            r2.font.size = Pt(40)
            r2.font.bold = True
            r2.font.color.rgb = INK_3
    if slide.break_sub:
        box, tf = _add_textbox(s, Inches(2), Inches(5.9), SLIDE_W - Inches(4), Inches(0.8))
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = slide.break_sub
        r.font.name = FONT_SANS
        r.font.size = Pt(15)
        r.font.color.rgb = INK_2

    _attach_notes(s, slide.speaker_notes)


def build_close_slide(prs, slide: Slide):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_solid_bg(s, INK)
    # Soft pinkish wash overlay
    overlay = _add_rect(s, 0, Inches(4.0), SLIDE_W, Inches(3.5),
                        fill=RGBColor(0x32, 0x0A, 0x0A))
    # Top thin gold band
    _add_rect(s, 0, 0, SLIDE_W, Inches(0.07), fill=GOLD)
    # Stamp
    if slide.close_stamp:
        stamp_text = slide.close_stamp.upper()
        sw = _pill_width(stamp_text, size=11)
        _add_pill(s, (SLIDE_W - sw) / 2, Inches(1.0), sw, Inches(0.45),
                  stamp_text, fill=None, text_color=GOLD,
                  border=GOLD, size=11)
    # H1 — multi-line with em→gold
    if slide.close_h1_runs:
        box, tf = _add_textbox(s, Inches(0.8), Inches(1.7),
                               SLIDE_W - Inches(1.6), Inches(3.5))
        tf.word_wrap = True
        # Split runs at \n into separate paragraphs
        cur_runs: List[TextRun] = []
        first = True
        all_paragraphs: List[List[TextRun]] = []
        for run in slide.close_h1_runs:
            if run.text == "\n":
                all_paragraphs.append(cur_runs)
                cur_runs = []
            else:
                cur_runs.append(run)
        if cur_runs:
            all_paragraphs.append(cur_runs)
        for runs in all_paragraphs:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            p.line_spacing = 1.05
            for run in runs:
                r = p.add_run()
                r.text = run.text
                r.font.name = FONT_SANS
                r.font.size = Pt(48)
                r.font.bold = True
                r.font.color.rgb = run.color or PAPER
    # Rule
    rule_w = Inches(3.6)
    _add_rect(s, (SLIDE_W - rule_w) / 2, Inches(5.4), rule_w * 0.6, Inches(0.07), fill=SCARLET)
    _add_rect(s, (SLIDE_W - rule_w) / 2 + rule_w * 0.6, Inches(5.4),
              rule_w * 0.4, Inches(0.07), fill=GOLD)
    if slide.close_lede:
        box, tf = _add_textbox(s, Inches(1.5), Inches(5.7),
                               SLIDE_W - Inches(3.0), Inches(1.1))
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.35
        r = p.add_run()
        r.text = slide.close_lede
        r.font.name = FONT_SANS
        r.font.size = Pt(13)
        r.font.color.rgb = RGBColor(0xD8, 0xD8, 0xD8)
    if slide.close_sig:
        box, tf = _add_textbox(s, 0, Inches(7.0), SLIDE_W, Inches(0.4))
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = slide.close_sig.upper()
        r.font.name = FONT_SANS
        r.font.size = Pt(9)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

    _attach_notes(s, slide.speaker_notes)


def build_editor_slide(prs, slide: Slide):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_solid_bg(s, INK)
    if slide.editor_badge:
        badge_text = slide.editor_badge.upper()
        bw = _pill_width(badge_text, size=11)
        _add_pill(s, (SLIDE_W - bw) / 2, Inches(1.5), bw, Inches(0.42),
                  badge_text, fill=SCARLET, text_color=PAPER,
                  border=SCARLET, size=11)
    # Arrow
    box, tf = _add_textbox(s, 0, Inches(2.05), SLIDE_W, Inches(1.0))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "→"
    r.font.name = FONT_SANS
    r.font.size = Pt(60)
    r.font.bold = True
    r.font.color.rgb = GOLD
    # H2
    if slide.editor_h2_runs:
        box, tf = _add_textbox(s, Inches(0.8), Inches(3.1),
                               SLIDE_W - Inches(1.6), Inches(1.7))
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.0
        for run in slide.editor_h2_runs:
            r = p.add_run()
            r.text = run.text
            r.font.name = FONT_SANS
            r.font.size = Pt(42)
            r.font.bold = True
            r.font.color.rgb = run.color or PAPER
    # What
    if slide.editor_what_runs:
        box, tf = _add_textbox(s, Inches(1.5), Inches(5.0),
                               SLIDE_W - Inches(3.0), Inches(1.5))
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.40
        for run in slide.editor_what_runs:
            r = p.add_run()
            r.text = run.text
            r.font.name = FONT_MONO if run.mono else FONT_SANS
            r.font.size = Pt(13)
            r.font.bold = run.bold
            r.font.italic = run.italic
            r.font.color.rgb = run.color or RGBColor(0xC0, 0xC0, 0xC0)
    if slide.editor_timer:
        box, tf = _add_textbox(s, 0, Inches(6.7), SLIDE_W, Inches(0.4))
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = slide.editor_timer
        r.font.name = FONT_MONO
        r.font.size = Pt(12)
        r.font.color.rgb = GOLD

    _attach_notes(s, slide.speaker_notes)


def build_content_slide(prs, slide: Slide):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_solid_bg(s, PAPER)
    _draw_top_brand_bar(s)
    _draw_head(s, slide.eyebrow, slide.minutes)
    _draw_title(s, slide.title, slide.subtitle)
    # Body
    blocks = slide.body_blocks
    # If subtitle is missing, lift body up a little
    body_y = BODY_Y
    body_h = BODY_H
    if not slide.subtitle:
        body_y = Inches(2.05)
        body_h = SLIDE_H - body_y - Inches(0.7)
    if blocks:
        if len(blocks) == 1:
            _render_block(s, BODY_X, body_y, BODY_W, body_h, blocks[0])
        else:
            _render_col_blocks(s, BODY_X, body_y, BODY_W, body_h, blocks)
    _draw_foot(s, slide.foot_left, slide.foot_right)

    _attach_notes(s, slide.speaker_notes)


# --- Speaker notes -----------------------------------------------------------

def _attach_notes(pp_slide, text: str):
    if not text:
        return
    notes = pp_slide.notes_slide
    tf = notes.notes_text_frame
    tf.text = text
    # Apply consistent font to notes
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.name = FONT_SANS
            r.font.size = Pt(11)


# ---------------------------------------------------------------------------
# Top-level entry point
# ---------------------------------------------------------------------------

def build_presentation(slides: List[Slide]) -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for slide in slides:
        if slide.kind == "cover":
            build_cover_slide(prs, slide)
        elif slide.kind == "section":
            build_section_slide(prs, slide)
        elif slide.kind == "break":
            build_break_slide(prs, slide)
        elif slide.kind == "close":
            build_close_slide(prs, slide)
        elif slide.kind == "editor":
            build_editor_slide(prs, slide)
        else:
            build_content_slide(prs, slide)
    return prs


def main() -> int:
    if not HTML_SRC.exists():
        print(f"ERROR: source HTML not found at {HTML_SRC}", file=sys.stderr)
        return 1
    html = HTML_SRC.read_text(encoding="utf-8")
    slides = parse_slides(html)
    print(f"Parsed {len(slides)} slides from {HTML_SRC.name}")
    if len(slides) != 58:
        print(f"WARNING: expected 58 slides, got {len(slides)}", file=sys.stderr)
    prs = build_presentation(slides)
    PPTX_OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(PPTX_OUT)
    print(f"Wrote {PPTX_OUT}  ({PPTX_OUT.stat().st_size:,} bytes)")
    # Verification pass: re-open the file and sanity-check.
    verify = Presentation(PPTX_OUT)
    n = len(verify.slides)
    notes_populated = sum(
        1 for sl in verify.slides
        if sl.has_notes_slide
        and (sl.notes_slide.notes_text_frame.text or "").strip()
    )
    text_frames = sum(
        1 for sl in verify.slides for sh in sl.shapes if sh.has_text_frame
    )
    print(
        f"Verify: {n} slides, {notes_populated} with speaker notes, "
        f"{text_frames} text frames; size = "
        f"{verify.slide_width / 914400:.3f}\" x {verify.slide_height / 914400:.3f}\""
    )
    if n != 58:
        print("ERROR: slide count mismatch in output file", file=sys.stderr)
        return 2
    if notes_populated < 50:
        print(
            f"WARNING: only {notes_populated}/58 slides have populated speaker notes",
            file=sys.stderr,
        )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
